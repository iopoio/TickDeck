from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

def create_default_template():
    prs = Presentation()
    
    # 16:9 Wide Layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide Layout
    # Layout 0 is usually Title Slide
    # But for a custom template, we might want to just build it on a blank slide or a layout we define.
    # Here, we'll use a blank layout and add placeholders manually if possible, 
    # or just use the first slide of a layout.
    
    blank_slide_layout = prs.slide_layouts[6] # 6 is usually blank
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background (Dark Navy)
    # Using a rectangle shape for the background for better compatibility
    bg_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(0x1C, 0x3D, 0x5A) # Navy
    bg_rect.line.fill.background() # No border
    
    # Title (HEADLINE)
    # x:10%, y:26% -> 1.33, 1.95
    txBox = slide.shapes.add_textbox(Inches(1.33), Inches(1.95), Inches(9.6), Inches(3.15))
    tf = txBox.text_frame
    tf.text = "HEADLINE PLACEHOLDER"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(62)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White
    
    # Subtitle Placeholder (SUBHEADLINE)
    # x:10% y:74% w:72% h:12%
    subBox = slide.shapes.add_textbox(Inches(1.33), Inches(5.55), Inches(9.6), Inches(0.9))
    stf = subBox.text_frame
    stf.text = "SUBHEADLINE PLACEHOLDER"
    sp = stf.paragraphs[0]
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB) # Gray
    
    # Company Pill (Idx 0 potentially? but we're doing manual now)
    # x:10%, y:16%
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.33), Inches(1.2), Inches(2.0), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(0x5A, 0x9E, 0x86) # Sage (Accent)
    pill.line.fill.background()
    
    ptf = pill.text_frame
    ptf.text = "COMPANY NAME"
    pp = ptf.paragraphs[0]
    pp.font.size = Pt(11)
    pp.font.bold = True
    pp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # Accent Line
    # x:10%, y:70% w:10%, h:3pt (3pt = 0.0416")
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.33), Inches(5.25), Inches(1.33), Inches(0.04))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = RGBColor(0x5A, 0x9E, 0x86) # Sage
    accent_line.line.fill.background()
    
    # Logo Position Marker
    # x: 우하단 (approx 10.5, 6.5)
    logo_marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(6.5), Inches(1.8), Inches(0.5))
    logo_marker.fill.solid()
    logo_marker.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    logo_marker.fill.transparency = 0.5
    logo_marker.text = "LOGO AREA"
    
    # Bottom Bar
    # h:6% -> 7.5 * 0.06 = 0.45"
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(0x5A, 0x9E, 0x86) # Sage
    bottom_bar.line.fill.background()
    
    # Save template
    os.makedirs('static/templates', exist_ok=True)
    prs.save('static/templates/cover_default.pptx')
    print("Template created at static/templates/cover_default.pptx")

if __name__ == "__main__":
    create_default_template()
