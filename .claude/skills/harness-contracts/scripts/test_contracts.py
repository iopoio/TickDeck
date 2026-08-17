import hashlib
import importlib.util
import json
import os
import pathlib
import re
import shutil
import subprocess
import sys
import tempfile
import unittest
from unittest import mock

import contract_checks as contract_checks_module
from contract_checks import (
    ContractViolation,
    validate_c1_proposition_dag,
    validate_c2_no_validation_metadata,
    validate_c3_trend_state_transition,
    validate_c4_citation_tracker,
    validate_c5_stage_order,
    validate_c6_content_authority,
    check_c8_genre_artifacts,
    check_c9_final_review,
    check_c14_viz_intent_preserved,
    check_c15_page_count_ceiling,
    validate_all_contracts,
)


def _missing_contract_check(*_args, **_kwargs):
    raise AssertionError("contract check is not implemented")


check_c13_role_duplication = getattr(contract_checks_module, "check_c13_role_duplication", _missing_contract_check)
check_deck_spec_gates = getattr(contract_checks_module, "check_deck_spec_gates", _missing_contract_check)
check_page_visual_intent_preserved = getattr(
    contract_checks_module, "check_page_visual_intent_preserved", _missing_contract_check
)

RENDER_DECK_PATH = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "scripts" / "render_deck.py"
RENDER_DECK_SPEC = importlib.util.spec_from_file_location("render_deck", RENDER_DECK_PATH)
render_deck_module = importlib.util.module_from_spec(RENDER_DECK_SPEC)
RENDER_DECK_SPEC.loader.exec_module(render_deck_module)

FACTCHECK_DUMP_PATH = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "scripts" / "factcheck_dump.py"
FACTCHECK_DUMP_SPEC = importlib.util.spec_from_file_location("factcheck_dump", FACTCHECK_DUMP_PATH)
factcheck_dump_module = importlib.util.module_from_spec(FACTCHECK_DUMP_SPEC)
FACTCHECK_DUMP_SPEC.loader.exec_module(factcheck_dump_module)

EXTERNAL_REVIEW_PATH = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "scripts" / "external_review.py"
EXTERNAL_REVIEW_SPEC = importlib.util.spec_from_file_location("external_review", EXTERNAL_REVIEW_PATH)
external_review_module = importlib.util.module_from_spec(EXTERNAL_REVIEW_SPEC)
EXTERNAL_REVIEW_SPEC.loader.exec_module(external_review_module)

PPTX_EXPORT_PATH = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "scripts" / "pptx_export.py"
PPTX_EXPORT_SPEC = importlib.util.spec_from_file_location("pptx_export", PPTX_EXPORT_PATH)
pptx_export_module = importlib.util.module_from_spec(PPTX_EXPORT_SPEC)
PPTX_EXPORT_SPEC.loader.exec_module(pptx_export_module)

QA_LINT_PATH = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "scripts" / "qa_lint.py"
QA_LINT_SPEC = importlib.util.spec_from_file_location("qa_lint", QA_LINT_PATH)
qa_lint_module = importlib.util.module_from_spec(QA_LINT_SPEC)
QA_LINT_SPEC.loader.exec_module(qa_lint_module)

QA_INK_PATH = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "scripts" / "qa_ink.py"
QA_INK_SPEC = importlib.util.spec_from_file_location("qa_ink", QA_INK_PATH)
qa_ink_module = importlib.util.module_from_spec(QA_INK_SPEC)
QA_INK_SPEC.loader.exec_module(qa_ink_module)

RUN_CONTRACTS_PATH = pathlib.Path(__file__).resolve().parent / "run_contracts.py"
RUN_CONTRACTS_SPEC = importlib.util.spec_from_file_location("run_contracts", RUN_CONTRACTS_PATH)
run_contracts_module = importlib.util.module_from_spec(RUN_CONTRACTS_SPEC)
RUN_CONTRACTS_SPEC.loader.exec_module(run_contracts_module)

R4_RUN_DIR = pathlib.Path(__file__).resolve().parents[4] / "_workspace" / "20260707_clo_report_r4"

RUN_DECK_SH = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "scripts" / "run_deck.sh"
CAPTURE_DECK_SH = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "scripts" / "capture_deck.sh"


VALID_DAG = {
    "nodes": [
        {"id": "thesis", "type": "thesis", "text": "AI execution shifts advantage toward proof assets"},
        {"id": "p1", "type": "claim", "text": "Discovery is moving from search to answer engines"},
        {"id": "p2", "type": "claim", "text": "Trust becomes a scarce asset"},
    ],
    "edges": [
        {"from": "thesis", "to": "p1"},
        {"from": "thesis", "to": "p2"},
    ],
}

VALID_INSIGHTS = [
    {
        "id": "i1",
        "claim": "Discovery moves from search pages to answer engines",
        "evidence_ids": ["src_a", "src_b"],
        "derivation_type": "synthesis",
        "counter_signal": "SEO spend remains resilient in some categories",
        "narrative_reason": "distribution control is shifting",
        "source_overlap_score": 0.24,
        "from_state": "keyword-led search",
        "to_state": "answer-led discovery",
        "mechanism": "agentic summaries collapse the funnel",
    },
    {
        "id": "i2",
        "claim": "Trust assets become more valuable as synthetic content rises",
        "evidence_ids": ["src_c", "src_d"],
        "derivation_type": "cross_source_inference",
        "counter_signal": "some low-consideration categories remain price-led",
        "narrative_reason": "synthetic abundance raises proof requirements",
        "source_overlap_score": 0.31,
        "from_state": "brand story accepted at face value",
        "to_state": "proof-backed credibility demanded",
        "mechanism": "audiences discount undifferentiated AI content",
    },
]

VALID_STAGE_LOG = [
    {"stage": "intake-director", "artifact": "workspace/00_intake.json"},
    {"stage": "collector", "artifact": "workspace/01_evidence_pool.json"},
    {"stage": "verifier", "artifact": "workspace/02_verified.json"},
    {"stage": "analyst", "artifact": "workspace/03_insights.json"},
    {"stage": "editorial-director", "artifact": "workspace/04_proposition_dag.json"},
    {"stage": "page-planner", "artifact": "workspace/05_page_plan.json"},
    {"stage": "designer", "artifact": "workspace/06_deck_spec.json"},
    {"stage": "qa-reviewer", "artifact": "workspace/07_qa.json"},
]

VALID_DECK = {
    "genre": "trend-report",
    "proposition_dag": VALID_DAG,
    "insights": VALID_INSIGHTS,
    "rendered_pages": [
        {"title": "Discovery Moves From Search To Answers", "body": "Answer engines compress comparison into one structured response."},
        {"title": "Trust Becomes The Scarce Asset", "body": "Proof-backed assets matter when synthetic content is cheap."},
    ],
    "stage_log": VALID_STAGE_LOG,
}

VALID_CONTENT_REGISTRY = {
    "sources": {
        "src_a": {"publisher": "Pew Research Center", "url": "https://example.com/pew"},
        "src_b": {"publisher": "IAB", "url": "https://example.com/iab"},
    },
    "metrics": {
        "metric_click_drop": {
            "value": "47",
            "unit": "%",
            "source_ids": ["src_a"],
            "scope": "Google searches with AI summaries",
        },
        "metric_measurement": {
            "value": "72",
            "unit": "%",
            "source_ids": ["src_b"],
            "scope": "cross-platform measurement adoption",
        },
    },
}

VALID_DECK_SPEC = {
    "pages": [
        {
            "page_id": "p01",
            "short_title": "Search Stops Sending Clicks",
            "layout": "hero_metric",
            "allowed_source_ids": ["src_a"],
            "allowed_metric_ids": ["metric_click_drop"],
            "content": [
                {"type": "headline", "text": "Search Stops Sending Clicks"},
                {"type": "metric", "metric_id": "metric_click_drop"},
                {"type": "citation", "src_id": "src_a"},
            ],
        },
        {
            "page_id": "p02",
            "short_title": "Measurement Moves To Owned Data",
            "layout": "stat_grid",
            "allowed_source_ids": ["src_b"],
            "allowed_metric_ids": ["metric_measurement"],
            "content": [
                {"type": "headline", "text": "Measurement Moves To Owned Data"},
                {"type": "metric", "metric_id": "metric_measurement"},
                {"type": "citation", "src_id": "src_b"},
            ],
        },
    ]
}

VALID_DECK_SPEC_WITH_EYEBROW = {
    "pages": [
        dict(
            VALID_DECK_SPEC["pages"][0],
            content=[
                {"type": "eyebrow", "text": "Chapter 1"},
                *VALID_DECK_SPEC["pages"][0]["content"],
            ],
        )
    ]
}

VALID_DECK_SPEC_WITH_VIZ = {
    "pages": [
        {
            "page_id": "p03",
            "short_title": "Search Stops Sending Clicks",
            "layout": "statement",
            "allowed_source_ids": ["src_a", "src_b"],
            "allowed_metric_ids": ["metric_click_drop", "metric_measurement"],
            "content": [
                {"type": "headline", "text": "Search Stops Sending Clicks"},
                {
                    "type": "viz",
                    "chart": "before_after",
                    "title": "Summary Collapse",
                    "series": [
                        {"label": "Baseline", "metric_id": "metric_measurement", "role": "baseline"},
                        {"label": "With Summary", "metric_id": "metric_click_drop", "role": "highlight"},
                    ],
                    "note": "Registry values only",
                },
            ],
        }
    ]
}

VALID_COVER_DECK_SPEC = {
    "pages": [
        {
            "page_id": "cover",
            "short_title": "표지",
            "layout": "cover",
            "allowed_source_ids": [],
            "allowed_metric_ids": [],
            "content": [
                {"type": "eyebrow", "text": "표지"},
                {"type": "headline", "text": "Market Shifts"},
                {"type": "summary", "text": "What changes when proof beats reach"},
            ],
        }
    ]
}

VALID_CLOSING_DECK_SPEC = {
    "theme": "tech",
    "pages": [
        {
            "page_id": "p15",
            "short_title": "전이를 가려내는 일이 분기점이다",
            "layout": "closing",
            "allowed_source_ids": ["src_a"],
            "allowed_metric_ids": [],
            "content": [
                {"type": "eyebrow", "text": "맺음 · 핵심 시사점"},
                {"type": "headline", "text": "전이를 가려내는 일이 분기점이다"},
                {
                    "type": "bullets",
                    "items": [
                        {"text": "도구 · 인프라 — 비가역적 투자에 올라탄다", "source_ids": ["src_a"]},
                        {"text": "자율화 — 통제된 파일럿으로 가둔다"},
                        {"text": "신흥 기술 — 관찰 대상으로 둔다"},
                    ],
                },
                {"type": "callout", "text": "승부는 어느 전이가 끝났느냐를 먼저 가려내는 데서 갈린다."},
                {"type": "citation", "src_id": "src_a"},
            ],
        }
    ],
}

VALID_RENDERED_HTML = """
<section class="slide">
  <h1>Search Stops Sending Clicks</h1>
  <span data-metric-id="metric_click_drop">47%</span>
  <cite data-src-id="src_a">Pew Research Center</cite>
  <span data-page-number>01 / 02</span>
</section>
"""


class HarnessContractTests(unittest.TestCase):
    def test_c1_accepts_connected_dag(self):
        self.assertEqual(validate_c1_proposition_dag(VALID_DAG), [])

    def test_c1_rejects_orphan_claim_nodes(self):
        broken = {
            "nodes": VALID_DAG["nodes"] + [{"id": "orphan", "type": "claim", "text": "route=all benchmark bucket"}],
            "edges": VALID_DAG["edges"],
        }
        violations = validate_c1_proposition_dag(broken)
        self.assertTrue(any("orphan" in str(v) for v in violations))

    def test_c2_rejects_validation_metadata_in_content(self):
        pages = [{"title": "단일출처 강등", "body": "정성근거로만 유지한 슬라이드"}]
        violations = validate_c2_no_validation_metadata(pages)
        self.assertEqual(len(violations), 1)
        self.assertIn("validation metadata", violations[0].message)

    def test_c3_requires_state_transition_for_trend_insights(self):
        invalid = [dict(VALID_INSIGHTS[0], mechanism="")]
        violations = validate_c3_trend_state_transition("trend-report", invalid)
        self.assertEqual(len(violations), 1)
        self.assertIn("from_state/to_state/mechanism", violations[0].message)

    def test_c4_requires_multi_source_low_overlap_insights(self):
        invalid = [dict(VALID_INSIGHTS[0], evidence_ids=["src_a"], source_overlap_score=0.12)]
        violations = validate_c4_citation_tracker(invalid)
        self.assertEqual(len(violations), 1)
        self.assertIn("evidence_ids", violations[0].message)

    def test_c4_rejects_high_source_overlap(self):
        invalid = [dict(VALID_INSIGHTS[0], evidence_ids=["src_a", "src_b"], source_overlap_score=0.91)]
        violations = validate_c4_citation_tracker(invalid)
        self.assertEqual(len(violations), 1)
        self.assertIn("source_overlap_score", violations[0].message)

    def test_c5_accepts_pipeline_order(self):
        self.assertEqual(validate_c5_stage_order(VALID_STAGE_LOG), [])

    def test_c5_rejects_designer_before_page_plan(self):
        invalid = [
            {"stage": "intake-director", "artifact": "workspace/00_intake.json"},
            {"stage": "designer", "artifact": "workspace/06_deck_spec.json"},
            {"stage": "page-planner", "artifact": "workspace/05_page_plan.json"},
        ]
        violations = validate_c5_stage_order(invalid)
        self.assertEqual(len(violations), 1)
        self.assertIn("designer", violations[0].message)

    def test_c6_accepts_deck_spec_references_inside_page_allowlists(self):
        self.assertEqual(validate_c6_content_authority(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, VALID_RENDERED_HTML), [])
        appendix_badge_html = """
        <section class="slide layout-source_appendix">
          <div class="verified-badge">모든 수치 출처 연결 검증 · 출처 2곳</div>
        </section>
        """
        self.assertEqual(validate_c6_content_authority(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, appendix_badge_html), [])

    def test_c6_accepts_metric_registry_sources_without_page_source_allowlist(self):
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    allowed_source_ids=[],
                    allowed_metric_ids=["metric_click_drop"],
                    content=[
                        {"type": "metric", "metric_id": "metric_click_drop"},
                        {"type": "citation", "src_id": "src_a"},
                    ],
                )
            ]
        }
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])

    def test_c6_accepts_singular_metric_source_alias_without_page_source_allowlist(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_click_drop": {
                    **VALID_CONTENT_REGISTRY["metrics"]["metric_click_drop"],
                    "source_ids": [],
                    "source_id": "src_a",
                },
            },
        }
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    allowed_source_ids=[],
                    allowed_metric_ids=["metric_click_drop"],
                    content=[
                        {"type": "metric", "metric_id": "metric_click_drop"},
                        {"type": "citation", "src_id": "src_a"},
                    ],
                )
            ]
        }

        self.assertEqual(validate_c6_content_authority(spec, registry, ""), [])

    def test_c6_still_rejects_disallowed_metric_unknown_source_and_untagged_number(self):
        cases = {
            "disallowed_metric": (
                {
                    "pages": [
                        dict(
                            VALID_DECK_SPEC["pages"][0],
                            allowed_metric_ids=[],
                            content=[{"type": "metric", "metric_id": "metric_click_drop"}],
                        )
                    ]
                },
                "",
                "metric_id not in page allowed_metric_ids: metric_click_drop",
            ),
            "unknown_source": (
                {
                    "pages": [
                        dict(
                            VALID_DECK_SPEC["pages"][0],
                            content=[{"type": "citation", "src_id": "src_missing"}],
                        )
                    ]
                },
                "",
                "unknown source id referenced: src_missing",
            ),
            "untagged_number": (
                VALID_DECK_SPEC,
                "<section><p>CTR fell 999% after AI summaries.</p></section>",
                "untagged number in rendered output",
            ),
        }
        for name, (spec, rendered_html, expected) in cases.items():
            with self.subTest(name=name):
                violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, rendered_html)
                self.assertTrue(any(expected in str(v) for v in violations))

    def test_c6_accepts_viz_blocks_with_metric_id_series(self):
        self.assertEqual(validate_c6_content_authority(VALID_DECK_SPEC_WITH_VIZ, VALID_CONTENT_REGISTRY, ""), [])

    def test_c6_accepts_r2_source_and_viz_options(self):
        registry = {
            "sources": {
                **VALID_CONTENT_REGISTRY["sources"],
                "src_a": {
                    **VALID_CONTENT_REGISTRY["sources"]["src_a"],
                    "short_name": "Pew",
                },
            },
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_click_drop": {
                    **VALID_CONTENT_REGISTRY["metrics"]["metric_click_drop"],
                    "period": "1Q26",
                },
            },
        }
        spec = {
            "meta": {"page_chrome": "title_band"},
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    content=[
                        {
                            "type": "viz",
                            "chart": "before_after",
                            "title": "Summary Collapse",
                            "source_caption": "off",
                            "title_style": "band",
                            "series": [
                                {"label": "Baseline", "metric_id": "metric_measurement", "role": "baseline"},
                                {"label": "With Summary", "metric_id": "metric_click_drop", "role": "highlight"},
                            ],
                        }
                    ],
                )
            ],
        }

        self.assertEqual(validate_c6_content_authority(spec, registry, ""), [])

    def test_c6_rejects_invalid_r2_optional_fields_and_viz_options(self):
        registry = {
            "sources": {
                "src_a": {
                    **VALID_CONTENT_REGISTRY["sources"]["src_a"],
                    "short_name": ["Pew"],
                },
                "src_b": VALID_CONTENT_REGISTRY["sources"]["src_b"],
            },
            "metrics": {
                "metric_click_drop": {
                    **VALID_CONTENT_REGISTRY["metrics"]["metric_click_drop"],
                    "period": {"label": "1Q26"},
                },
                "metric_measurement": VALID_CONTENT_REGISTRY["metrics"]["metric_measurement"],
            },
        }
        spec = {
            "meta": {"page_chrome": "poster_band"},
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    content=[
                        {
                            "type": "viz",
                            "chart": "before_after",
                            "source_caption": "manual",
                            "title_style": "badge",
                            "series": [
                                {"label": "Baseline", "metric_id": "metric_measurement", "role": "baseline"},
                                {"label": "With Summary", "metric_id": "metric_click_drop", "role": "highlight"},
                            ],
                        }
                    ],
                )
            ],
        }

        violations = validate_c6_content_authority(spec, registry, "")

        self.assertTrue(any("source short_name must be text" in str(v) for v in violations))
        self.assertTrue(any("metric period must be text" in str(v) for v in violations))
        self.assertTrue(any("unsupported page_chrome: poster_band" in str(v) for v in violations))
        self.assertTrue(any("unsupported viz source_caption: manual" in str(v) for v in violations))
        self.assertTrue(any("unsupported viz title_style: badge" in str(v) for v in violations))

    def test_c6_rejects_running_head_title_band_collision(self):
        spec = {
            "meta": {
                "page_chrome": "title_band",
                "running_head": True,
            },
            "pages": [dict(VALID_DECK_SPEC["pages"][0])],
        }

        violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, "")

        self.assertTrue(any("page_chrome title_band cannot combine with running_head" in str(v) for v in violations))

    def test_c6_rejects_title_band_title_over_contract_limit(self):
        spec = {
            "meta": {"page_chrome": "title_band"},
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    short_title="아주 긴 밴드 제목 " * 12,
                )
            ],
        }

        violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, "")

        self.assertTrue(any("title_band title exceeds" in str(v) for v in violations))

    def test_c6_accepts_derived_metric_contract_and_inherited_sources(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_cagr": {
                    "label": "훈련 데이터 연평균 성장률",
                    "value": "8.9",
                    "unit": "%/년",
                    "derivation": "cagr",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                    "period": "2020~2025",
                    "formula_note": "(72/47)^(1/5)-1",
                },
            },
        }
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    allowed_source_ids=[],
                    allowed_metric_ids=["metric_cagr"],
                    content=[{"type": "metric", "metric_id": "metric_cagr"}, {"type": "citation", "src_id": "src_a"}],
                )
            ]
        }

        self.assertEqual(validate_c6_content_authority(spec, registry, ""), [])

    def test_c6_rejects_invalid_derived_metric_contracts(self):
        cases = {
            "bad_enum": (
                {
                    "metric_bad": {
                        "value": "260",
                        "unit": "%",
                        "derivation": "ratio",
                        "derived_from": ["metric_click_drop", "metric_measurement"],
                        "source_ids": [],
                        "status": "derived",
                    }
                },
                "unsupported metric derivation: ratio",
            ),
            "unknown_ref": (
                {
                    "metric_bad": {
                        "value": "260",
                        "unit": "%",
                        "derivation": "cagr",
                        "derived_from": ["metric_click_drop", "metric_missing"],
                        "source_ids": [],
                        "status": "derived",
                    }
                },
                "derived_from references unknown metric: metric_missing",
            ),
            "single_ref": (
                {
                    "metric_bad": {
                        "value": "260",
                        "unit": "%",
                        "derivation": "cagr",
                        "derived_from": ["metric_click_drop"],
                        "source_ids": [],
                        "status": "derived",
                    }
                },
                "derived metric must include at least 2 derived_from refs",
            ),
            "non_empty_sources": (
                {
                    "metric_bad": {
                        "value": "260",
                        "unit": "%",
                        "derivation": "cagr",
                        "derived_from": ["metric_click_drop", "metric_measurement"],
                        "source_ids": ["src_a"],
                        "status": "derived",
                    }
                },
                "derived metric source_ids must be empty",
            ),
            "non_empty_singular_source": (
                {
                    "metric_bad": {
                        "value": "260",
                        "unit": "%",
                        "derivation": "cagr",
                        "derived_from": ["metric_click_drop", "metric_measurement"],
                        "source_id": "src_a",
                        "status": "derived",
                    }
                },
                "derived metric source_ids must be empty",
            ),
            "cycle": (
                {
                    "metric_a": {
                        "value": "1",
                        "unit": "%",
                        "derivation": "delta_pct",
                        "derived_from": ["metric_b", "metric_click_drop"],
                        "source_ids": [],
                        "status": "derived",
                    },
                    "metric_b": {
                        "value": "2",
                        "unit": "%",
                        "derivation": "delta_pct",
                        "derived_from": ["metric_a", "metric_measurement"],
                        "source_ids": [],
                        "status": "derived",
                    },
                },
                "derived metric cycle",
            ),
        }
        for name, (extra_metrics, expected) in cases.items():
            with self.subTest(name=name):
                registry = {
                    "sources": VALID_CONTENT_REGISTRY["sources"],
                    "metrics": {**VALID_CONTENT_REGISTRY["metrics"], **extra_metrics},
                }
                violations = validate_c6_content_authority(VALID_DECK_SPEC, registry, "")
                self.assertTrue(any(expected in str(v) for v in violations), [str(v) for v in violations])

    def _series_registry(self, n_series=3, n_points=4, unit="%"):
        metrics = dict(VALID_CONTENT_REGISTRY["metrics"])
        for s in range(n_series):
            for p in range(n_points):
                metrics[f"metric_s{s}_p{p}"] = {
                    "value": str(10 + p),
                    "unit": unit,
                    "source_ids": ["src_a"],
                    "series_id": f"series_{s}",
                    "series_key": f"202{p}",
                }
        return {"sources": VALID_CONTENT_REGISTRY["sources"], "metrics": metrics}

    def test_c6_rejects_series_with_duplicate_keys_or_mixed_units(self):
        registry = self._series_registry(n_series=1)
        registry["metrics"]["metric_dup"] = {
            "value": "99", "unit": "%", "source_ids": ["src_a"],
            "series_id": "series_0", "series_key": "2020",
        }
        registry["metrics"]["metric_unit"] = {
            "value": "5", "unit": "억", "source_ids": ["src_a"],
            "series_id": "series_0", "series_key": "2029",
        }
        violations = validate_c6_content_authority({"pages": []}, registry, "")
        self.assertTrue(any("duplicate series_key" in str(v) for v in violations), [str(v) for v in violations])
        self.assertTrue(any("mixes units" in str(v) for v in violations), [str(v) for v in violations])

    def test_c6_report_tone_requires_chartable_series(self):
        spec = {"meta": {"tone": "report"}, "pages": []}
        violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("REPORT_TONE_DATA_THIN" in str(v) for v in violations), [str(v) for v in violations])
        ok = validate_c6_content_authority(spec, self._series_registry(), "")
        self.assertFalse(any("REPORT_TONE_DATA_THIN" in str(v) for v in ok), [str(v) for v in ok])

    def test_c6_report_tone_requires_chart_hero_composition(self):
        hero = {
            "page_id": "hero", "short_title": "h", "layout": "statement",
            "allowed_source_ids": ["src_a"], "allowed_metric_ids": ["metric_s0_p0"],
            "content": [
                {"type": "headline", "text": "차트 주인공"},
                {"type": "viz", "chart": "multi_line", "title": "Trend",
                 "series": [{"label": "L", "metric_id": "metric_s0_p0", "role": "highlight"}]},
            ],
        }
        texty = {
            "page_id": "texty", "short_title": "t", "layout": "stack",
            "allowed_source_ids": ["src_a"], "allowed_metric_ids": [],
            "content": [{"type": "headline", "text": "글"}, {"type": "note", "text": "글"}, {"type": "note", "text": "글"}],
        }
        registry = self._series_registry()
        bad = {"meta": {"tone": "report"}, "pages": [texty, dict(texty, page_id="t2"), dict(texty, page_id="t3"), hero]}
        violations = validate_c6_content_authority(bad, registry, "")
        self.assertTrue(any("REPORT_TONE_COMPOSITION" in str(v) for v in violations), [str(v) for v in violations])
        good = {"meta": {"tone": "report"}, "pages": [hero, dict(hero, page_id="h2"), texty]}
        ok = validate_c6_content_authority(good, registry, "")
        self.assertFalse(any("REPORT_TONE_COMPOSITION" in str(v) for v in ok), [str(v) for v in ok])

    def test_c6_no_report_tone_checks_without_tone(self):
        spec = {"pages": []}
        violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, "")
        self.assertFalse(any("REPORT_TONE" in str(v) for v in violations), [str(v) for v in violations])

    def test_c11_requires_source_coverage_and_big3_axes(self):
        import tempfile, pathlib
        from contract_checks import check_c11_source_coverage
        with tempfile.TemporaryDirectory() as tmp:
            run = pathlib.Path(tmp)
            # coverage 부재
            (run / "01_evidence_pool.json").write_text(json.dumps({"items": []}), encoding="utf-8")
            vs = check_c11_source_coverage(run)
            self.assertTrue(any("source_coverage 부재" in str(v) for v in vs))
            # none 판정에 검색어 1개 = 조기 포기 위반 + 축 누락
            (run / "01_evidence_pool.json").write_text(json.dumps({
                "source_coverage": [
                    {"axis": "kpmg", "queries": ["kpmg ai jobs"], "found": [], "verdict": "none"},
                ]
            }), encoding="utf-8")
            vs = check_c11_source_coverage(run)
            self.assertTrue(any("조기 포기" in str(v) for v in vs), [str(v) for v in vs])
            self.assertTrue(any("미탐색: pwc" in str(v) for v in vs))
            # 정상 케이스
            rows = [{"axis": a, "queries": ["q1", "q2"], "found": ["src_001"], "verdict": "found"}
                    for a in ("kpmg", "pwc", "deloitte", "government_stats", "academic")]
            (run / "01_evidence_pool.json").write_text(json.dumps({"source_coverage": rows}), encoding="utf-8")
            self.assertEqual(check_c11_source_coverage(run), [])

    def test_c6_rejects_derived_metric_with_wrong_recomputed_value(self):
        # 검산 게이트 (7/7 제대리 리뷰): 원천값 재계산과 등재값이 안 맞으면 위반
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_bad_delta": {
                    "label": "YoY 변화",
                    "value": "18",
                    "unit": "%",
                    "derivation": "delta_pct",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                    "formula_note": "(72-47)/47*100",
                },
            },
        }
        violations = validate_c6_content_authority({"pages": []}, registry, "")
        self.assertTrue(any("!= recomputed" in str(v) for v in violations), [str(v) for v in violations])

    def test_c6_rejects_derived_metric_without_formula_note(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_no_note": {
                    "label": "배수",
                    "value": "1.53",
                    "unit": "배",
                    "derivation": "multiple",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                },
            },
        }
        violations = validate_c6_content_authority({"pages": []}, registry, "")
        self.assertTrue(any("formula_note" in str(v) for v in violations), [str(v) for v in violations])

    def test_c6_accepts_r3_section_nav_annotations_and_metric_commentary_schema(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_delta_yoy": {
                    "label": "YoY 변화",
                    "value": "53.2",
                    "unit": "%",
                    "derivation": "delta_pct",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                    "formula_note": "(72-47)/47*100",
                },
                "metric_delta_qoq": {
                    "label": "QoQ 변화",
                    "value": "53.2",
                    "unit": "%",
                    "derivation": "delta_pct",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                    "formula_note": "(72-47)/47*100",
                },
            },
        }
        spec = {
            "meta": {"section_nav": "chips"},
            "pages": [
                {
                    "page_id": "r3_schema",
                    "short_title": "Metric Commentary",
                    "layout": "metric_commentary",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": ["metric_click_drop", "metric_measurement", "metric_delta_yoy", "metric_delta_qoq"],
                    "rows": [
                        {
                            "heading_metric_id": "metric_measurement",
                            "headline_metric_id": "metric_delta_yoy",
                            "bullets": [
                                {"label": "YoY", "metric_id": "metric_delta_yoy"},
                                {"label": "QoQ", "metric_id": "metric_delta_qoq"},
                            ],
                            "chart": {
                                "chart": "quarterly_bars",
                                "series": [
                                    {"metric_id": "metric_click_drop", "role": "baseline"},
                                    {"metric_id": "metric_measurement", "role": "highlight"},
                                ],
                            },
                        }
                    ],
                    "content": [],
                },
                {
                    "page_id": "r3_viz",
                    "short_title": "Annotated Trend",
                    "layout": "statement",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": ["metric_click_drop", "metric_measurement", "metric_delta_yoy"],
                    "content": [
                        {
                            "type": "viz",
                            "chart": "multi_line",
                            "series": [
                                {"label": "Start", "metric_id": "metric_click_drop", "role": "baseline"},
                                {"label": "End", "metric_id": "metric_measurement", "role": "highlight"},
                            ],
                            "annotations": [
                                {"kind": "callout", "metric_id": "metric_delta_yoy", "anchor_series": 1, "shape": "ellipse"},
                                {"kind": "endpoint_value", "series": 1},
                                {"kind": "trend_arrow", "series": 0},
                                {"kind": "event_band", "label": "COVID", "from_key": "Start", "to_key": "End"},
                            ],
                        }
                    ],
                },
            ],
        }

        self.assertEqual(validate_c6_content_authority(spec, registry, ""), [])

    def test_c6_rejects_invalid_r3_section_nav_annotations_and_metric_commentary(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_delta_yoy": {
                    "label": "YoY 변화",
                    "value": "18",
                    "unit": "%",
                    "derivation": "delta_abs",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                },
            },
        }
        spec = {
            "meta": {"section_nav": "freeform"},
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    content=[
                        {
                            "type": "viz",
                            "chart": "donut",
                            "series": [{"label": "Now", "metric_id": "metric_click_drop"}],
                            "annotations": [
                                {"kind": "callout", "metric_id": "metric_click_drop", "anchor_series": 0},
                                {"kind": "event_band", "label": "2020 shock", "from_key": "A", "to_key": "B"},
                            ],
                        }
                    ],
                ),
                {
                    "page_id": "bad_commentary",
                    "short_title": "Bad Commentary",
                    "layout": "metric_commentary",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": ["metric_click_drop", "metric_delta_yoy"],
                    "rows": [
                        {
                            "heading_metric_id": "metric_click_drop",
                            "headline_metric_id": "metric_delta_yoy",
                            "bullets": [{"label": "YoY", "metric_id": "metric_delta_yoy"}],
                            "chart": {"chart": "multi_line", "series": [{"metric_id": "metric_click_drop"}]},
                        }
                    ],
                    "content": [],
                },
            ],
        }

        violations = validate_c6_content_authority(spec, registry, "")

        self.assertTrue(any("unsupported section_nav: freeform" in str(v) for v in violations))
        self.assertTrue(any("annotations are supported only for" in str(v) for v in violations))
        self.assertTrue(any("callout annotation metric must reference a derived metric" in str(v) for v in violations))
        self.assertTrue(any("event_band label contains raw number" in str(v) for v in violations))
        self.assertTrue(any("event_band from_key must reference a series label" in str(v) for v in violations))
        self.assertTrue(any("metric_commentary headline_metric_id must reference delta_pct" in str(v) for v in violations))
        self.assertTrue(any("metric_commentary chart must be quarterly_bars" in str(v) for v in violations))

    def test_c6_rejects_unknown_viz_series_role(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    content=[
                        {
                            "type": "viz",
                            "chart": "before_after",
                            "series": [
                                {"label": "Baseline", "metric_id": "metric_measurement", "role": "decorative"},
                                {"label": "With Summary", "metric_id": "metric_click_drop", "role": "highlight"},
                            ],
                        }
                    ],
                )
            ]
        }

        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")

        self.assertTrue(any("unsupported viz series role: decorative" in str(v) for v in violations))

    def test_c6_accepts_semantic_color_roles_and_new_finance_charts(self):
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    content=[
                        {
                            "type": "viz",
                            "chart": "quarterly_bars",
                            "series": [
                                {"metric_id": "metric_measurement", "role": "baseline"},
                                {"metric_id": "metric_click_drop", "role": "positive"},
                                {"metric_id": "metric_measurement", "role": "negative"},
                                {"metric_id": "metric_click_drop", "role": "brand", "color": "#F7931A"},
                            ],
                        },
                        {
                            "type": "viz",
                            "chart": "fin_table",
                            "columns": ["이전", "현재", "변화"],
                            "series": [
                                {
                                    "label": "매출",
                                    "row_role": "group",
                                    "cells": [
                                        {"metric_id": "metric_measurement"},
                                        {"metric_id": "metric_click_drop"},
                                        {"text": "흑자전환"},
                                    ],
                                }
                            ],
                        },
                    ],
                )
            ]
        }

        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])

    def test_c6_rejects_fin_table_raw_number_text_cells(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    content=[
                        {
                            "type": "viz",
                            "chart": "fin_table",
                            "columns": ["이전", "현재"],
                            "series": [
                                {
                                    "label": "매출",
                                    "cells": [
                                        {"metric_id": "metric_measurement"},
                                        {"text": "999억"},
                                    ],
                                }
                            ],
                        }
                    ],
                )
            ]
        }

        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")

        self.assertTrue(any("fin_table text cell contains raw number" in str(v) for v in violations))

    def test_c6_accepts_text_table_without_numeric_cells(self):
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    layout="statement",
                    allowed_source_ids=[],
                    allowed_metric_ids=[],
                    content=[
                        {"type": "headline", "text": "Player Comparison"},
                        {
                            "type": "text_table",
                            "title": "브랜드 포지션",
                            "columns": ["브랜드", "주력 제품", "가격대", "포지션", "채널"],
                            "rows": [
                                ["CLO", "넥앤프로 에어코즈 메디소닉", "중가", "헬스 뷰티 멀티부위", "자사몰 SNS 없음"],
                                ["TheraFace", "얼굴 케어 기기", "프리미엄", "뷰티 테크", "자사몰 리테일"],
                            ],
                            "highlight_row": 0,
                        },
                    ],
                )
            ]
        }
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Text Table Fixture")
        self.assertIn("<table", html)
        self.assertIn("text-table", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_c6_rejects_text_table_numeric_cells_as_untagged_numbers(self):
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    layout="statement",
                    allowed_source_ids=[],
                    allowed_metric_ids=[],
                    content=[
                        {"type": "headline", "text": "Player Comparison"},
                        {
                            "type": "text_table",
                            "columns": ["브랜드", "주력 제품", "가격대", "포지션"],
                            "rows": [["CLO", "넥앤프로", "중가", "인지도 999%"]],
                        },
                    ],
                )
            ]
        }
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Text Table Numeric Fixture")
        violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html)
        self.assertTrue(any("untagged number in rendered output" in str(v) for v in violations))

    def test_c6_backed_number_in_body_passes(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_launch_sales": {
                    "value": "300",
                    "unit": "만개",
                    "source_ids": ["src_a"],
                }
            },
        }
        rendered = "<section><p>검증된 성과는 300만 개였다.</p></section>"

        self.assertEqual(validate_c6_content_authority(VALID_DECK_SPEC, registry, rendered), [])

    def test_c6_unbacked_number_in_body_fails(self):
        rendered = "<section><p>근거 없는 매출 999억이 본문에 들어갔다.</p></section>"

        violations = validate_c6_content_authority(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, rendered)

        self.assertTrue(any("untagged number in rendered output" in str(v) for v in violations))

    def test_c6_year_in_body_passes(self):
        rendered = "<section><p>2018~19년 흐름과 2024년 기준 변화가 이어졌다.</p></section>"

        self.assertEqual(validate_c6_content_authority(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, rendered), [])

    def test_c6_manual_source_label_still_fails(self):
        rendered = "<section><p>출처: 어쩌구</p></section>"

        violations = validate_c6_content_authority(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, rendered)

        self.assertTrue(any("manual source label" in str(v) for v in violations))

    def test_c6_enclosed_numeral_still_fails(self):
        rendered = "<section><p>① 첫 번째 근거를 본문에 썼다.</p></section>"

        violations = validate_c6_content_authority(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, rendered)

        self.assertTrue(any("enclosed numeral" in str(v) for v in violations))

    def test_c6_accepts_closing_layout_in_supported_layout_enum(self):
        self.assertIn("closing", contract_checks_module.SUPPORTED_LAYOUTS)
        self.assertEqual(validate_c6_content_authority(VALID_CLOSING_DECK_SPEC, VALID_CONTENT_REGISTRY, ""), [])

    def test_c6_accepts_new_signature_layouts_in_supported_layout_enum(self):
        for layout in ("mosaic_tiles", "split_status", "scenario_cards"):
            with self.subTest(layout=layout):
                self.assertIn(layout, contract_checks_module.SUPPORTED_LAYOUTS)
                spec = {"pages": [dict(VALID_DECK_SPEC["pages"][0], layout=layout)]}
                self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])

    def test_c6_accepts_pricing_cards_layout_in_supported_layout_enum(self):
        self.assertIn("pricing_cards", contract_checks_module.SUPPORTED_LAYOUTS)
        spec = {"pages": [dict(VALID_DECK_SPEC["pages"][0], layout="pricing_cards")]}
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])

    def test_c6_accepts_swot_quad_viz_without_metric_ids(self):
        self.assertIn("swot_quad", contract_checks_module.SUPPORTED_VIZ_CHART_TYPES)
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=[],
                    content=[
                        {
                            "type": "viz",
                            "chart": "swot_quad",
                            "title": "Qualitative SWOT",
                            "series": [
                                {"label": "Strengths", "items": ["Owned proof assets"], "role": "highlight"},
                                {"label": "Weaknesses", "items": ["Low aided awareness"]},
                                {"label": "Opportunities", "items": ["Partner distribution"]},
                                {"label": "Threats", "items": ["Platform policy shifts"]},
                            ],
                        }
                    ],
                )
            ]
        }
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])

    def test_c6_rejects_swot_quad_items_with_raw_numbers(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=[],
                    content=[
                        {
                            "type": "viz",
                            "chart": "swot_quad",
                            "series": [
                                {"label": "Strengths", "items": ["Awareness rose 47%"]},
                                {"label": "Weaknesses", "items": ["Manual process"]},
                                {"label": "Opportunities", "items": ["Retail partner"]},
                                {"label": "Threats", "items": ["Policy change"]},
                            ],
                        }
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("swot_quad item contains raw number" in str(v) for v in violations))

    def test_c6_accepts_gantt_and_pictograph_without_metric_ids(self):
        for chart, series in (
            (
                "gantt",
                [
                    {"label": "Discovery", "start": "2026-08", "end": "2026-09"},
                    {"label": "Launch", "start": "2026-10", "end": "2026-10", "milestone": True},
                ],
            ),
            ("pictograph", [{"label": "Adoption", "total": 14, "filled": 9}]),
        ):
            with self.subTest(chart=chart):
                spec = {
                    "pages": [
                        dict(
                            VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                            allowed_metric_ids=[],
                            content=[{"type": "viz", "chart": chart, "series": series}],
                        )
                    ]
                }
                self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])

    def test_c6_rejects_invalid_gantt_ranges_and_oversized_pictograph(self):
        invalid_specs = (
            (
                "gantt",
                [{"label": "Discovery", "start": "2026-10", "end": "2026-08"}],
                "gantt series item end must not precede start",
            ),
            (
                "pictograph",
                [{"label": "Adoption", "total": 21, "filled": 9}],
                "pictograph total must be at most 20",
            ),
        )
        for chart, series, message in invalid_specs:
            with self.subTest(chart=chart):
                spec = {
                    "pages": [
                        dict(
                            VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                            allowed_metric_ids=[],
                            content=[{"type": "viz", "chart": chart, "series": series}],
                        )
                    ]
                }
                violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, "")
                self.assertTrue(any(message in str(v) for v in violations))

    def test_c6_rejects_missing_labels_for_batch5_qualitative_charts(self):
        for chart, item in (
            ("gantt", {"start": "2026-08", "end": "2026-09"}),
            ("pictograph", {"total": 14, "filled": 9}),
        ):
            with self.subTest(chart=chart):
                spec = {
                    "pages": [
                        dict(
                            VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                            allowed_metric_ids=[],
                            content=[{"type": "viz", "chart": chart, "series": [item]}],
                        )
                    ]
                }
                violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, "")
                self.assertTrue(any(f"{chart} series item must include label" in str(v) for v in violations))

    def test_c6_rejects_gantt_bounds_and_multiple_pictograph_series(self):
        invalid_specs = (
            (
                "gantt",
                [{"label": "Long plan", "start": "2026-01", "end": "2029-02"}],
                "gantt schedule must span at most 36 months",
            ),
            (
                "gantt",
                [{"label": "Hidden lane", "start": "2026-01", "end": "2026-02", "lane": 8}],
                "gantt series item lane must be within 0..7",
            ),
            (
                "pictograph",
                [{"label": "A", "total": 4, "filled": 2}, {"label": "B", "total": 4, "filled": 3}],
                "pictograph must include exactly one series item",
            ),
        )
        for chart, series, message in invalid_specs:
            with self.subTest(message=message):
                spec = {
                    "pages": [
                        dict(
                            VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                            allowed_metric_ids=[],
                            content=[{"type": "viz", "chart": chart, "series": series}],
                        )
                    ]
                }
                violations = validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, "")
                self.assertTrue(any(message in str(v) for v in violations))

    # ── R5 논증 어휘(2026-07-08 DESIGN_R5_argument_diagrams.md) — pyramid/causal_chain/
    #    two_by_two/tradeoff. 4종 각각 정상 렌더 1 + 슬롯 미달 위반 1.

    def test_c6_accepts_pyramid_with_two_evidence_layers_and_renders(self):
        self.assertIn("pyramid", contract_checks_module.SUPPORTED_VIZ_CHART_TYPES)
        content = [
            {
                "type": "viz",
                "chart": "pyramid",
                "title": "실행 역량이 우위를 가른다",
                "series": [
                    {"role": "claim", "label": "조직적 실행 역량이 M&A 성패를 가른다"},
                    {"role": "evidence", "label": "카브아웃 실행 준비도", "metric_id": "metric_click_drop"},
                    {"role": "evidence", "label": "통합 체계 관리 우선순위"},
                ],
            }
        ]
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=["metric_click_drop"],
                    content=content,
                )
            ]
        }
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Pyramid Fixture")
        self.assertIn("<svg", html)
        self.assertIn('data-metric-id="metric_click_drop"', html)
        self.assertIn("47%", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_c6_rejects_pyramid_with_fewer_than_two_evidence_layers(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=[],
                    content=[
                        {
                            "type": "viz",
                            "chart": "pyramid",
                            "series": [
                                {"role": "claim", "label": "주장"},
                                {"role": "evidence", "label": "근거 하나뿐"},
                            ],
                        }
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("ARG_DIAGRAM_THIN" in str(v) and "pyramid" in str(v) for v in violations))

    def test_c6_accepts_causal_chain_with_three_nodes_and_renders(self):
        content = [
            {
                "type": "viz",
                "chart": "causal_chain",
                "title": "인과 사슬",
                "series": [
                    {"label": "검색 요약 확대"},
                    {"label": "클릭 감소", "evidence": "요약이 클릭을 대체", "metric_id": "metric_click_drop"},
                    {"label": "측정 체계 전환"},
                ],
            }
        ]
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=["metric_click_drop"],
                    content=content,
                )
            ]
        }
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Causal Chain Fixture")
        self.assertIn("<svg", html)
        self.assertIn('data-metric-id="metric_click_drop"', html)
        self.assertIn("47%", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_c6_rejects_causal_chain_with_fewer_than_three_nodes(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=[],
                    content=[
                        {
                            "type": "viz",
                            "chart": "causal_chain",
                            "series": [{"label": "A"}, {"label": "B"}],
                        }
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("ARG_DIAGRAM_THIN" in str(v) and "causal_chain" in str(v) for v in violations))

    def test_c6_accepts_two_by_two_with_three_items_and_renders(self):
        content = [
            {
                "type": "viz",
                "chart": "two_by_two",
                "title": "포지셔닝 맵",
                "x_axis": {"low": "저가", "high": "고가"},
                "y_axis": {"low": "로컬", "high": "글로벌"},
                "series": [
                    {"label": "플레이어 A", "x": 0.2, "y": 0.8},
                    {"label": "플레이어 B", "x": 0.6, "y": 0.3, "emphasis": True},
                    {"label": "플레이어 C", "x": 0.9, "y": 0.1},
                ],
            }
        ]
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=[],
                    content=content,
                )
            ]
        }
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Two By Two Fixture")
        self.assertIn("<svg", html)
        self.assertIn("플레이어 A", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_c6_rejects_two_by_two_with_fewer_than_three_items(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=[],
                    content=[
                        {
                            "type": "viz",
                            "chart": "two_by_two",
                            "x_axis": {"low": "저가", "high": "고가"},
                            "y_axis": {"low": "로컬", "high": "글로벌"},
                            "series": [
                                {"label": "플레이어 A", "x": 0.2, "y": 0.8},
                                {"label": "플레이어 B", "x": 0.6, "y": 0.3},
                            ],
                        }
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("ARG_DIAGRAM_THIN" in str(v) and "two_by_two" in str(v) for v in violations))

    def test_c6_rejects_two_by_two_item_xy_out_of_range(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=[],
                    content=[
                        {
                            "type": "viz",
                            "chart": "two_by_two",
                            "x_axis": {"low": "저가", "high": "고가"},
                            "y_axis": {"low": "로컬", "high": "글로벌"},
                            "series": [
                                {"label": "플레이어 A", "x": 0.2, "y": 0.8},
                                {"label": "플레이어 B", "x": 1.4, "y": 0.3},
                                {"label": "플레이어 C", "x": 0.5, "y": 0.5},
                            ],
                        }
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("two_by_two item x must be a number within 0..1" in str(v) for v in violations))

    def test_c6_accepts_tradeoff_with_both_sides_and_renders(self):
        content = [
            {
                "type": "viz",
                "chart": "tradeoff",
                "title": "빌드 vs 매수",
                "left_title": "자체 구축",
                "right_title": "인수",
                "series": [
                    {"side": "left", "label": "속도 통제"},
                    {"side": "right", "label": "즉시 확보", "metric_id": "metric_click_drop"},
                ],
            }
        ]
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=["metric_click_drop"],
                    content=content,
                )
            ]
        }
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, ""), [])
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Tradeoff Fixture")
        self.assertIn("<svg", html)
        self.assertIn('data-metric-id="metric_click_drop"', html)
        self.assertIn("47%", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_c6_rejects_tradeoff_missing_one_side(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=[],
                    content=[
                        {
                            "type": "viz",
                            "chart": "tradeoff",
                            "left_title": "자체 구축",
                            "right_title": "인수",
                            "series": [
                                {"side": "left", "label": "속도 통제"},
                                {"side": "left", "label": "문화 적합성"},
                            ],
                        }
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("ARG_DIAGRAM_THIN" in str(v) and "tradeoff" in str(v) for v in violations))

    def test_c6_rejects_unsupported_layouts(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    layout="unknown_layout",
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertEqual(len(violations), 1)
        self.assertIn("unsupported layout: unknown_layout", violations[0].message)

    def test_c6_rejects_viz_raw_numbers_in_designer_fields(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    content=[
                        {
                            "type": "viz",
                            "chart": "before_after",
                            "title": "CTR fell 47%",
                            "series": [
                                {"label": "Before 72%", "metric_id": "metric_measurement", "role": "baseline"},
                                {"label": "After", "metric_id": "metric_click_drop", "role": "highlight"},
                            ],
                            "note": "Gap is 25pp",
                        }
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("raw number" in str(v) for v in violations))

    def test_c6_rejects_viz_metric_ids_outside_page_allowlist(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    allowed_metric_ids=["metric_click_drop"],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertTrue(any("metric_id not in page allowed_metric_ids: metric_measurement" in str(v) for v in violations))

    def test_c6_rejects_unsupported_content_block_types_before_render(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    content=[
                        {"type": "headline", "text": "Search Stops Sending Clicks"},
                        {"type": "sparkline", "metric_id": "metric_click_drop"},
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertEqual(len(violations), 1)
        self.assertIn("unsupported content block type: sparkline", violations[0].message)

    def test_c6_rejects_unknown_or_disallowed_references(self):
        invalid = {
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    allowed_source_ids=["src_a"],
                    allowed_metric_ids=[],
                    content=[
                        {"type": "metric", "metric_id": "metric_click_drop"},
                        {"type": "citation", "src_id": "src_missing"},
                    ],
                )
            ]
        }
        violations = validate_c6_content_authority(invalid, VALID_CONTENT_REGISTRY, "")
        self.assertEqual(len(violations), 3)
        self.assertTrue(any("allowed_metric_ids" in str(v) for v in violations))
        self.assertTrue(any("unknown source id" in str(v) for v in violations))

    def test_c6_rejects_untagged_numbers_and_manual_source_labels_in_render(self):
        rendered = """
        <section class="slide">
          <h1>Manual Render</h1>
          <p>CTR fell 999% after AI summaries.</p>
          <p>출처: Pew Research Center</p>
        </section>
        """
        violations = validate_c6_content_authority(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, rendered)
        self.assertEqual(len(violations), 2)
        self.assertTrue(any("untagged number" in str(v) for v in violations))
        self.assertTrue(any("manual source label" in str(v) for v in violations))

    def test_validate_all_contracts_passes_valid_deck(self):
        self.assertEqual(validate_all_contracts(VALID_DECK), [])

    def test_validate_all_contracts_runs_c6_when_deck_spec_is_present(self):
        deck = dict(
            VALID_DECK,
            deck_spec=VALID_DECK_SPEC,
            content_registry=VALID_CONTENT_REGISTRY,
            rendered_html=VALID_RENDERED_HTML,
        )
        self.assertEqual(validate_all_contracts(deck), [])

    def test_c8_accepts_market_research_required_artifacts(self):
        intake = {"genre": "market-research"}
        evidence_pool = {"items": [{"source_type": "observation"} for _ in range(5)]}
        page_plan = {"pages": [{"genre_artifact": "taxonomy"}, {"genre_artifact": "player_table"}]}
        self.assertEqual(check_c8_genre_artifacts(intake, evidence_pool, page_plan), [])

    def test_c8_market_research_aliases_include_spaced_korean_and_brand_research(self):
        evidence_pool = {"items": [{"source_type": "observation"} for _ in range(5)]}
        page_plan = {"pages": [{"genre_artifact": "taxonomy"}]}
        for genre in ("시장 조사", "brand-research", "경쟁 분석"):
            with self.subTest(genre=genre):
                violations = check_c8_genre_artifacts({"genre": genre}, evidence_pool, page_plan)
                self.assertEqual(len(violations), 1)
                self.assertEqual(violations[0].contract_id, "C8")
                self.assertIn("player_table", violations[0].message)

    def test_c8_rejects_market_research_without_player_table(self):
        intake = {"genre": "market_research"}
        evidence_pool = {"items": [{"source_type": "observation"} for _ in range(5)]}
        page_plan = {"pages": [{"genre_artifact": "taxonomy"}]}
        violations = check_c8_genre_artifacts(intake, evidence_pool, page_plan)
        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C8")
        self.assertIn("player_table", violations[0].message)

    def test_c8_rejects_market_research_with_too_few_observations(self):
        intake = {"genre": "시장조사"}
        evidence_pool = {"items": [{"source_type": "observation"} for _ in range(3)]}
        page_plan = {"pages": [{"genre_artifact": "taxonomy"}, {"genre_artifact": "player_table"}]}
        violations = check_c8_genre_artifacts(intake, evidence_pool, page_plan)
        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C8")
        self.assertIn("requires at least 5 observation evidence items (found 3)", violations[0].message)

    def test_c8_skips_non_market_research_genres(self):
        intake = {"genre": "trend-report"}
        evidence_pool = {"items": []}
        page_plan = {"pages": []}
        self.assertEqual(check_c8_genre_artifacts(intake, evidence_pool, page_plan), [])

    def test_c8_skips_missing_genre(self):
        evidence_pool = {"items": []}
        page_plan = {"pages": []}
        self.assertEqual(check_c8_genre_artifacts({}, evidence_pool, page_plan), [])

    def _write_c10_run(self, run_dir: pathlib.Path, genre: str, source_registry):
        (run_dir / "00_intake.json").write_text(json.dumps({"genre": genre}), encoding="utf-8")
        (run_dir / "02_verified.json").write_text(
            json.dumps({"source_registry": source_registry}),
            encoding="utf-8",
        )

    def _check_c10(self, run_dir: pathlib.Path):
        self.assertTrue(
            hasattr(contract_checks_module, "check_c10_collection_evidence"),
            "check_c10_collection_evidence must exist",
        )
        return contract_checks_module.check_c10_collection_evidence(run_dir)

    def test_c10_market_research_rejects_legacy_registry_without_doc_type(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            self._write_c10_run(
                run_dir,
                "market-research",
                {"src_1": {"tier": "Tier-A", "local_path": "pdf/source.pdf"}},
            )

            violations = self._check_c10(run_dir)

        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C10")
        self.assertIn("doc_type", violations[0].message)
        self.assertEqual(violations[0].path, "02_verified.json.source_registry")

    def test_c10_market_research_rejects_fewer_than_five_canonical_sources(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            for index in range(4):
                source_file = run_dir / "pdf" / f"source_{index}.pdf"
                source_file.parent.mkdir(exist_ok=True)
                source_file.write_text("pdf", encoding="utf-8")
            self._write_c10_run(
                run_dir,
                "market-research",
                {
                    f"src_{index}": {
                        "tier": "Tier-A",
                        "doc_type": "pdf",
                        "local_path": f"pdf/source_{index}.pdf",
                        "cited_pages": [1],
                    }
                    for index in range(4)
                },
            )

            violations = self._check_c10(run_dir)

        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C10")
        self.assertIn("at least 5 canonical document sources (found 4)", violations[0].message)

    def test_c10_market_research_accepts_five_pdf_sources_with_files_and_pages(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            for index in range(5):
                source_file = run_dir / "pdf" / f"source_{index}.pdf"
                source_file.parent.mkdir(exist_ok=True)
                source_file.write_text("pdf", encoding="utf-8")
            self._write_c10_run(
                run_dir,
                "market-research",
                {
                    f"src_{index}": {
                        "tier": "Tier-A",
                        "doc_type": "pdf",
                        "local_path": f"pdf/source_{index}.pdf",
                        "cited_pages": [1, "12-15"],
                    }
                    for index in range(5)
                },
            )

            violations = self._check_c10(run_dir)

        self.assertEqual(violations, [])

    def test_c10_market_research_rejects_missing_local_path_file_per_source(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            for index in range(4):
                source_file = run_dir / "pdf" / f"source_{index}.pdf"
                source_file.parent.mkdir(exist_ok=True)
                source_file.write_text("pdf", encoding="utf-8")
            self._write_c10_run(
                run_dir,
                "market-research",
                {
                    f"src_{index}": {
                        "tier": "Tier-A",
                        "doc_type": "pdf",
                        "local_path": f"pdf/source_{index}.pdf",
                        "cited_pages": [1],
                    }
                    for index in range(5)
                },
            )

            violations = self._check_c10(run_dir)

        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C10")
        self.assertIn("local_path file missing", violations[0].message)
        self.assertEqual(violations[0].path, "02_verified.json.source_registry.src_4.local_path")

    def test_c10_market_research_requires_pdf_pages_but_accepts_extract_note(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            for index in range(5):
                source_file = run_dir / "pdf" / f"source_{index}.pdf"
                source_file.parent.mkdir(exist_ok=True)
                source_file.write_text("source", encoding="utf-8")
            self._write_c10_run(
                run_dir,
                "market-research",
                {
                    "src_pdf": {
                        "tier": "Tier-A",
                        "doc_type": "pdf",
                        "local_path": "pdf/source_0.pdf",
                    },
                    **{
                        f"src_db_{index}": {
                            "tier": "Tier-A",
                            "doc_type": "official_db_extract",
                            "local_path": f"pdf/source_{index}.pdf",
                            "extract_note": "official database extract",
                        }
                        for index in range(1, 5)
                    },
                },
            )

            violations = self._check_c10(run_dir)

        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C10")
        self.assertIn("cited_pages", violations[0].message)
        self.assertEqual(violations[0].path, "02_verified.json.source_registry.src_pdf.cited_pages")

    def test_c10_skips_non_market_research_genre(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            self._write_c10_run(run_dir, "trend-report", {})

            violations = self._check_c10(run_dir)

        self.assertEqual(violations, [])

    def test_c12_accepts_seed_sources_registered_and_appendix_split(self):
        deck_spec = {
            "pages": [
                {
                    "page_id": "appendix",
                    "short_title": "출처",
                    "layout": "source_appendix",
                    "allowed_source_ids": ["src_seed", "src_research"],
                    "allowed_metric_ids": [],
                    "content": [{"type": "headline", "text": "출처"}],
                }
            ]
        }
        registry = {
            "source_registry": {
                "src_seed": {
                    "publisher": "Client",
                    "title": "Client Deck 2026",
                    "local_path": "seed/client_deck.pdf",
                    "provenance": "seed",
                },
                "src_research": {
                    "publisher": "Pew",
                    "title": "Pew Report",
                    "url": "https://example.com/report",
                    "provenance": "research",
                },
            },
            "metric_registry": VALID_CONTENT_REGISTRY["metrics"],
        }
        html = render_deck_module.render_deck(deck_spec, registry, title="Seed Appendix Fixture")
        deck = dict(
            VALID_DECK,
            intake={"provided_sources": [{"kind": "file", "ref": "uploads/client_deck.pdf"}]},
            deck_spec=deck_spec,
            content_registry=registry,
            rendered_html=html,
        )

        self.assertIn("제공하신 자료", html)
        self.assertIn("추가 조사", html)
        self.assertEqual(validate_all_contracts(deck), [])

    def test_c12_rejects_provided_source_missing_from_registry(self):
        deck_spec = {
            "pages": [
                {
                    "page_id": "appendix",
                    "short_title": "출처",
                    "layout": "source_appendix",
                    "allowed_source_ids": ["src_research"],
                    "allowed_metric_ids": [],
                    "content": [{"type": "headline", "text": "출처"}],
                }
            ]
        }
        registry = {
            "source_registry": {
                "src_research": {
                    "publisher": "Pew",
                    "title": "Pew Report",
                    "url": "https://example.com/report",
                    "provenance": "research",
                },
            },
            "metric_registry": VALID_CONTENT_REGISTRY["metrics"],
        }
        deck = dict(
            VALID_DECK,
            intake={"provided_sources": [{"kind": "url", "ref": "https://example.com/client"}]},
            deck_spec=deck_spec,
            content_registry=registry,
            rendered_html="<section><div class='eyebrow'>제공하신 자료 (1건)</div><div class='eyebrow'>추가 조사 (1건)</div></section>",
        )

        violations = validate_all_contracts(deck)

        self.assertEqual(len([violation for violation in violations if violation.contract_id == "C12"]), 1)
        self.assertIn("provided source missing", str(violations[-1]))

    def test_c9_rejects_missing_external_review_file(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            (run_dir / "deck.html").write_text("<section>Deck</section>", encoding="utf-8")

            violations = check_c9_final_review(run_dir)

        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C9")
        self.assertIn("final external review missing", violations[0].message)

    def test_c9_accepts_matching_hash_with_codex_success(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            html_path = run_dir / "deck.html"
            html_path.write_text("<section>Deck</section>", encoding="utf-8")
            deck_hash = contract_checks_module.sha256_file(html_path)
            (run_dir / "review_codex.txt").write_text("지적 1. " + ("충분한 리뷰 본문입니다. " * 20), encoding="utf-8")
            (run_dir / "08_external_review.json").write_text(
                json.dumps(
                    {
                        "deck_html_sha256": deck_hash,
                        "codex": {"ok": True, "file": "review_codex.txt"},
                        "gemini": {"ok": False, "file": "review_gemini.txt"},
                    }
                ),
                encoding="utf-8",
            )

            violations = check_c9_final_review(run_dir)

        self.assertEqual(violations, [])

    def test_c9_rejects_deck_html_changed_after_external_review(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            html_path = run_dir / "deck.html"
            html_path.write_text("<section>Current</section>", encoding="utf-8")
            reviewed_hash = "0" * 64
            (run_dir / "08_external_review.json").write_text(
                json.dumps(
                    {
                        "deck_html_sha256": reviewed_hash,
                        "codex": {"ok": True, "file": "review_codex.txt"},
                        "gemini": {"ok": False, "file": "review_gemini.txt"},
                    }
                ),
                encoding="utf-8",
            )

            violations = check_c9_final_review(run_dir)

        self.assertEqual(len(violations), 1)
        self.assertIn("deck.html changed after external review", violations[0].message)
        self.assertIn(reviewed_hash[:12], violations[0].message)

    def test_c9_rejects_when_both_external_reviewers_failed(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            html_path = run_dir / "deck.html"
            html_path.write_text("<section>Deck</section>", encoding="utf-8")
            deck_hash = contract_checks_module.sha256_file(html_path)
            (run_dir / "08_external_review.json").write_text(
                json.dumps(
                    {
                        "deck_html_sha256": deck_hash,
                        "codex": {"ok": False, "file": "review_codex.txt"},
                        "gemini": {"ok": False, "file": "review_gemini.txt"},
                    }
                ),
                encoding="utf-8",
            )

            violations = check_c9_final_review(run_dir)

        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C9")
        self.assertIn("both external reviewers failed", violations[0].message)

    def test_c9_rejects_ok_reviewers_when_output_files_are_sentinel_or_too_short(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            html_path = run_dir / "deck.html"
            html_path.write_text("<section>Deck</section>", encoding="utf-8")
            deck_hash = contract_checks_module.sha256_file(html_path)
            (run_dir / "review_codex.txt").write_text("REVIEWER_TIMEOUT (600s)\n", encoding="utf-8")
            (run_dir / "review_gemini.txt").write_text("too short\n", encoding="utf-8")
            (run_dir / "08_external_review.json").write_text(
                json.dumps(
                    {
                        "deck_html_sha256": deck_hash,
                        "codex": {"ok": True, "file": "review_codex.txt"},
                        "gemini": {"ok": True, "file": "review_gemini.txt"},
                    }
                ),
                encoding="utf-8",
            )

            violations = check_c9_final_review(run_dir)

        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C9")
        self.assertIn("both external reviewers failed", violations[0].message)

    def test_c9_accepts_when_one_ok_reviewer_has_substantive_output(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            html_path = run_dir / "deck.html"
            html_path.write_text("<section>Deck</section>", encoding="utf-8")
            deck_hash = contract_checks_module.sha256_file(html_path)
            (run_dir / "review_codex.txt").write_text("REVIEWER_FAILED_TO_START\nboom\n", encoding="utf-8")
            (run_dir / "review_gemini.txt").write_text("지적 1. " + ("충분한 리뷰 본문입니다. " * 20), encoding="utf-8")
            (run_dir / "08_external_review.json").write_text(
                json.dumps(
                    {
                        "deck_html_sha256": deck_hash,
                        "codex": {"ok": True, "file": "review_codex.txt"},
                        "gemini": {"ok": True, "file": "review_gemini.txt"},
                    }
                ),
                encoding="utf-8",
            )

            violations = check_c9_final_review(run_dir)

        self.assertEqual(violations, [])

    def test_render_deck_injects_metric_values_and_generated_citations(self):
        html = render_deck_module.render_deck(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, title="C6 Fixture")
        self.assertIn('data-metric-id="metric_click_drop"', html)
        self.assertIn("47%", html)
        self.assertIn('data-src-id="src_a"', html)
        self.assertIn("Pew Research Center", html)
        self.assertNotIn("출처:", html)
        self.assertEqual(validate_c6_content_authority(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, html), [])

        appendix_spec = {
            "pages": [
                dict(VALID_DECK_SPEC["pages"][0]),
                {
                    "page_id": "appendix",
                    "short_title": "출처",
                    "layout": "source_appendix",
                    "allowed_source_ids": ["src_a", "src_b"],
                    "allowed_metric_ids": [],
                    "content": [{"type": "headline", "text": "출처"}],
                },
            ]
        }
        appendix_registry = {
            "sources": {
                "src_a": {
                    "publisher": "Pew Research Center",
                    "title": "Pew & Partners 2026",
                    "url": "https://example.com/pew?q=ai&src=deck",
                },
                "src_b": {"publisher": "IAB", "title": "IAB Report", "url": ""},
            },
            "metrics": VALID_CONTENT_REGISTRY["metrics"],
        }
        appendix_html = render_deck_module.render_deck(appendix_spec, appendix_registry, title="Appendix Fixture")
        # 7/7 후추님: 출처는 텍스트만 — 링크·화살표 미표기 (URL은 registry에 보존)
        self.assertNotIn('<a class="appendix-link"', appendix_html)
        self.assertIn('Pew &amp; Partners 2026', appendix_html)
        self.assertIn('<span class="appendix-title" data-src-id="src_b">IAB Report</span>', appendix_html)
        # 7/22 후추님: 검증 배지("모든 수치 출처 연결 검증 · 출처 N곳") 제거 — 독자용 아님. 부재를 검증.
        self.assertNotIn("모든 수치 출처 연결 검증", appendix_html)
        self.assertEqual(validate_c6_content_authority(appendix_spec, appendix_registry, appendix_html), [])

    def test_render_deck_resolves_inline_metric_tokens_across_text_paths(self):
        registry = {
            "sources": {
                "src_inline": {"publisher": "Inline Source", "url": "https://example.com/inline"},
            },
            "metrics": {
                # verified registry 실물은 source_id 단수형을 쓴다. 인라인 토큰만으로도 출처가 잡혀야 한다.
                "metric_inline": {
                    "value": "47",
                    "unit": "%",
                    "source_ids": [],
                    "source_id": "src_inline",
                },
            },
        }
        token = "{{metric_inline}}"
        spec = {
            "pages": [
                {
                    "page_id": "matrix_inline",
                    "short_title": "Matrix",
                    "layout": "matrix",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [
                        {"type": "headline", "text": f"기준 {token}"},
                        {"type": "body", "text": f"왼쪽 — 점유 {token}"},
                        {"type": "body", "text": f"오른쪽 {token}"},
                    ],
                },
                {
                    "page_id": "poster_inline",
                    "short_title": "Poster",
                    "layout": "poster",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [{"type": "body", "text": f"포스터 {token}"}],
                },
                {
                    "page_id": "index_inline",
                    "short_title": "Index",
                    "layout": "index",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [{"type": "bullets", "items": [f"Part — 설명 {token}"]}],
                },
                {
                    "page_id": "divider_inline",
                    "short_title": "Divider",
                    "layout": "divider",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [{"type": "headline", "text": f"간지 {token}"}],
                },
                {
                    "page_id": "cover_inline",
                    "short_title": "Cover",
                    "layout": "cover",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [
                        {"type": "headline", "text": f"표지 {token}"},
                        {"type": "body", "text": f"부제 {token}"},
                    ],
                },
                {
                    "page_id": "svg_inline",
                    "short_title": "Two By Two",
                    "layout": "statement",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [
                        {
                            "type": "viz",
                            "chart": "two_by_two",
                            "title": f"지도 {token}",
                            "x_axis": {"low": f"낮음 {token}", "high": f"높음 {token}"},
                            "y_axis": {"low": "낮음", "high": "높음"},
                            "series": [
                                {"label": f"A {token}", "x": 0.2, "y": 0.2},
                                {"label": "B", "x": 0.5, "y": 0.5},
                                {"label": "C", "x": 0.8, "y": 0.8},
                            ],
                        },
                        {"type": "footnote", "definition": f"각주 {token}"},
                    ],
                },
                {
                    "page_id": "dashboard_inline",
                    "short_title": "Dashboard",
                    "layout": "dashboard",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [{"type": "body", "text": f"대시보드 {token}"}],
                },
                {
                    "page_id": "stepper_inline",
                    "short_title": "Stepper",
                    "layout": "stepper",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [{"type": "body", "text": f"단계 {token}"}],
                },
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="Inline Metric Fixture")

        self.assertNotIn(token, html)
        self.assertGreaterEqual(html.count("47%"), 14)
        matrix_page_match = re.search(
            r'data-page-id="matrix_inline"(?P<body>.*?)(?=<section class="slide|</body>)',
            html,
            flags=re.DOTALL,
        )
        self.assertIsNotNone(matrix_page_match)
        matrix_page = matrix_page_match.group("body")
        self.assertIn('data-metric-id="metric_inline"', matrix_page)
        self.assertIn('data-src-id="src_inline"', matrix_page)

    def test_format_metric_value_rounds_trillion_display_without_mutating_registry(self):
        cases = (
            ({"value": "33.3425", "unit": "조원"}, "33.34 조원"),
            ({"value": "33.3000", "unit": "조 원"}, "33.3 조 원"),
            ({"value": "33.999", "unit": "조원"}, "34.0 조원"),
            ({"value": "-0.004", "unit": "조원"}, "0.0 조원"),
            ({"value": "4.0", "unit": "조원"}, "4.0 조원"),
            ({"value": "0.20", "unit": "조원"}, "0.20 조원"),
            ({"value": "33", "unit": "조"}, "33조"),
            ({"value": "102961", "unit": "억원"}, "10조 2,961억 원"),
        )
        for metric, expected in cases:
            with self.subTest(metric=metric):
                original = dict(metric)
                self.assertEqual(render_deck_module._format_metric_value(metric), expected)
                self.assertEqual(metric, original)

    def test_svg_inline_metric_keeps_authority_context_after_display_rounding(self):
        registry = {
            "sources": {"src_inline": {"publisher": "Inline", "url": "https://example.com/inline"}},
            "metrics": {
                "metric_inline": {
                    "value": "33.3425",
                    "unit": "조원",
                    "source_id": "src_inline",
                }
            },
        }
        token = "{{metric_inline}}"
        spec = {
            "pages": [
                {
                    "page_id": "svg_rounded_inline",
                    "short_title": "SVG metric",
                    "layout": "statement",
                    "allowed_source_ids": ["src_inline"],
                    "allowed_metric_ids": ["metric_inline"],
                    "content": [
                        {
                            "type": "viz",
                            "chart": "two_by_two",
                            "title": f"시장 {token}",
                            "note": f"주석 {token}",
                            "x_axis": {"low": f"낮음 {token}", "high": "높음"},
                            "y_axis": {"low": "낮음", "high": "높음"},
                            "series": [
                                {"label": f"A {token}", "x": 0.2, "y": 0.2},
                                {"label": "B", "x": 0.5, "y": 0.5},
                                {"label": "C", "x": 0.8, "y": 0.8},
                            ],
                        }
                    ],
                }
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="SVG Metric")

        self.assertNotIn(token, html)
        self.assertIn('<tspan data-metric-id="metric_inline">33.34 조원</tspan>', html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_svg_wrapping_measures_metric_display_value_without_dropping_token(self):
        registry = {
            "sources": {
                "src_main": {"publisher": "Main", "url": "https://example.com/main"},
                "src_inline": {"publisher": "Inline", "url": "https://example.com/inline"},
            },
            "metrics": {
                "metric_main": {"value": "60", "unit": "%", "source_id": "src_main"},
                "metric_inline": {"value": "47", "unit": "%", "source_id": "src_inline"},
            },
        }
        token = "{{metric_inline}}"
        spec = {
            "pages": [
                {
                    "page_id": "svg_wrap_inline",
                    "short_title": "SVG wrap",
                    "layout": "statement",
                    "allowed_source_ids": ["src_main", "src_inline"],
                    "allowed_metric_ids": ["metric_main", "metric_inline"],
                    "content": [
                        {
                            "type": "viz",
                            "chart": "donut",
                            "series": [
                                {
                                    "metric_id": "metric_main",
                                    "label": f"AAAA BBBB CCCC {token}",
                                    "role": "highlight",
                                }
                            ],
                        }
                    ],
                }
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="SVG Wrap")

        self.assertNotIn(token, html)
        self.assertIn('<tspan data-metric-id="metric_inline">47%</tspan>', html)
        self.assertIn('data-src-id="src_inline"', html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_render_deck_telemetry_absent_by_default(self):
        # R6 — meta.telemetry 없으면 gtag 문자열 0회 (기본 OFF·출력 바이트 불변 보호)
        html = render_deck_module.render_deck(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, title="No Telemetry Fixture")
        self.assertEqual(html.count("gtag"), 0)

    def test_render_deck_telemetry_injects_ga4_snippet_when_enabled(self):
        # R6 — meta.telemetry.ga_id 있으면 이벤트 3종 + ga_id 삽입
        telemetry_spec = dict(VALID_DECK_SPEC, meta={"telemetry": {"ga_id": "G-TESTID123"}})
        html = render_deck_module.render_deck(telemetry_spec, VALID_CONTENT_REGISTRY, title="Telemetry Fixture")
        self.assertIn("gtag", html)
        self.assertIn("deck_page", html)
        self.assertIn("deck_source_click", html)
        self.assertIn("deck_read_end", html)
        self.assertIn("G-TESTID123", html)

    def test_render_deck_source_appendix_links_only_http_schemes(self):
        appendix_spec = {
            "pages": [
                {
                    "page_id": "appendix",
                    "short_title": "출처",
                    "layout": "source_appendix",
                    "allowed_source_ids": ["src_a", "src_b"],
                    "allowed_metric_ids": [],
                    "content": [{"type": "headline", "text": "출처"}],
                }
            ]
        }
        appendix_registry = {
            "sources": {
                "src_a": {"publisher": "Pew", "title": "Safe", "url": "http://example.com/report"},
                "src_b": {"publisher": "Bad", "title": "Unsafe", "url": "javascript:alert(1)"},
            },
            "metrics": VALID_CONTENT_REGISTRY["metrics"],
        }

        appendix_html = render_deck_module.render_deck(appendix_spec, appendix_registry, title="Appendix Fixture")

        self.assertNotIn('<a class="appendix-link"', appendix_html)
        self.assertIn('Safe', appendix_html)
        self.assertIn('<span class="appendix-title" data-src-id="src_b">Unsafe</span>', appendix_html)
        self.assertNotIn('href="javascript:alert(1)"', appendix_html)

    def test_render_deck_text_table_cells_use_rich_text_markup(self):
        spec = {
            "pages": [
                {
                    "page_id": "table_fixture",
                    "short_title": "Table",
                    "layout": "statement",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": [],
                    "content": [
                        {
                            "type": "text_table",
                            "columns": ["구분", "해석"],
                            "rows": [["A", "==중요== 신호"]],
                        }
                    ],
                }
            ]
        }

        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Text Table Rich Fixture")

        self.assertIn('<td><b class="kw">중요</b> 신호</td>', html)
        self.assertNotIn("==중요==", html)

    def test_render_deck_outputs_svg_for_each_supported_viz_chart(self):
        for chart in ("before_after", "dumbbell", "flow", "big_number", "gap_map", "shift"):
            with self.subTest(chart=chart):
                spec = {
                    "pages": [
                        dict(
                            VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                            content=[
                                dict(
                                    VALID_DECK_SPEC_WITH_VIZ["pages"][0]["content"][1],
                                    chart=chart,
                                )
                            ],
                        )
                    ]
                }
                html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Viz Fixture")
                self.assertIn("<svg", html)
                self.assertIn('data-metric-id="metric_click_drop"', html)
                self.assertIn("47%", html)
                self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_cover_hides_cover_role_and_uses_light_background(self):
        html = render_deck_module.render_deck(VALID_COVER_DECK_SPEC, VALID_CONTENT_REGISTRY, title="Cover Fixture")
        rendered_body = html.split("</style>", 1)[1]
        self.assertIn("Market Shifts", html)
        self.assertNotIn("표지", rendered_body)
        self.assertNotIn('<footer class="slide-foot">', html)
        self.assertNotIn("background: #111", html)
        self.assertIn("axis-strip", html)

    def test_render_deck_keeps_footer_in_flow_so_page_number_stays_visible(self):
        # 푸터를 정상 흐름(flex 아이템)에 두고 본문은 overflow:hidden로 잘리게 해서
        # 과밀 슬라이드에서도 본문이 푸터/페이지번호 위로 겹치지 못하게 한다(겹침 버그의 근본 차단).
        html = render_deck_module.render_deck(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, title="Footer Fixture")
        self.assertIn(".slide-foot {", html)
        self.assertIn("flex: 0 0 auto;", html)
        self.assertIn("overflow: hidden;", html)
        self.assertIn('class="page-number" data-page-number>01 / 02</span>', html)

    def test_render_deck_wraps_and_abbreviates_long_source_rows(self):
        registry = {
            "sources": {
                f"src_{index}": {"publisher": f"Very Long Publisher Name {index}", "url": f"https://example.com/{index}"}
                for index in range(1, 7)
            },
            "metrics": VALID_CONTENT_REGISTRY["metrics"],
        }
        spec = {
            "pages": [
                {
                    "page_id": "source_fixture",
                    "short_title": "Source Fixture",
                    "layout": "statement",
                    "allowed_source_ids": [f"src_{index}" for index in range(1, 7)],
                    "allowed_metric_ids": [],
                    "content": [
                        {"type": "headline", "text": "Source Fixture"},
                        *({"type": "citation", "src_id": f"src_{index}"} for index in range(1, 7)),
                    ],
                }
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="Source Fixture")
        source_css = html.split(".source-row {", 1)[1].split("}", 1)[0]

        self.assertIn("flex-wrap: wrap;", source_css)
        self.assertIn("white-space: normal;", source_css)
        self.assertNotIn("overflow: hidden;", source_css)
        self.assertIn('class="source-more"', html)
        self.assertIn("+2", html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_capture_deck_detects_hidden_source_row_clips(self):
        script = CAPTURE_DECK_SH.read_text(encoding="utf-8")

        self.assertIn("FIT_SOURCE_CLIP", script)
        self.assertIn("sourceClips", script)

    def test_capture_deck_uses_true_type_webfont_and_waits_before_pdf_print(self):
        script = CAPTURE_DECK_SH.read_text(encoding="utf-8")

        self.assertIn("Pretendard-Regular.woff2", script)
        self.assertIn("format('woff2')", script)
        self.assertIn("document.fonts.ready", script)
        self.assertIn("document.fonts.check", script)
        self.assertIn("--virtual-time-budget=", script)
        self.assertIn("PDF_FONT_TYPE3", script)
        self.assertIn("PDF_FONT_READY_ERROR", script)
        self.assertIn("Bad bounding box in Type 3 glyph", script)
        self.assertIn("Promise.resolve(window.__tickdeckFontsReady).then", script)
        self.assertRegex(script, r'--virtual-time-budget=\d+\s+--dump-dom')
        self.assertIn("FIT_FONT_ERROR", script)
        self.assertIn("FIT_CHECK_ERROR", script)
        self.assertRegex(script, r"FITREPORT\(_ERROR\)\?\\\|")
        self.assertLess(script.index('rm -f "$OUT"'), script.index('FONT_VERSION='))
        self.assertLess(script.index("FIT_CHECK_ERROR"), script.index('mv "$PDF_TMP" "$OUT"'))

        injector = script.split("<<'PY'\n", 1)[1].split("\nPY\n", 1)[0]
        with tempfile.TemporaryDirectory() as tmp:
            root = pathlib.Path(tmp)
            source = root / "deck.html"
            capture = root / "capture.html"
            source.write_text(
                '<html><head><style>:root{--font-body:"Pretendard", "Apple SD Gothic Neo", sans-serif;'
                '--font-head:ui-monospace,"Pretendard","Apple SD Gothic Neo",monospace;}</style></head>'
                '<body>본문에서 "Pretendard"를 인용한다.</body></html>',
                encoding="utf-8",
            )
            for name in (
                "Thin",
                "ExtraLight",
                "Light",
                "Regular",
                "Medium",
                "SemiBold",
                "Bold",
                "ExtraBold",
                "Black",
            ):
                (root / f"Pretendard-{name}.woff2").write_bytes(f"fake-{name}".encode())
            original_argv = sys.argv
            try:
                sys.argv = ["capture-font-inject", str(source), str(capture), str(root)]
                exec(compile(injector, "capture-font-inject", "exec"), {})
            finally:
                sys.argv = original_argv

            captured = capture.read_text(encoding="utf-8")
            font_style = re.search(r'<style id="tickdeck-pdf-font">(.*?)</style>', captured, re.DOTALL)
            self.assertIsNotNone(font_style)
            self.assertEqual(font_style.group(1).count("@font-face"), 9)
            self.assertIn("font-weight:100", font_style.group(1))
            self.assertIn("font-weight:300", font_style.group(1))
            source_style = re.search(r'<head><style>(.*?)</style>', captured, re.DOTALL)
            self.assertIsNotNone(source_style)
            self.assertNotIn('"Pretendard"', source_style.group(1))
            self.assertNotIn('"Apple SD Gothic Neo"', source_style.group(1))
            self.assertIn('--font-head:ui-monospace,"TickDeck PDF Pretendard"', captured)
            self.assertIn('본문에서 "Pretendard"를 인용한다.', captured)

    def test_render_main_fails_when_pdf_capture_fails(self):
        capture_failure = subprocess.CompletedProcess(
            args=["bash", str(CAPTURE_DECK_SH)],
            returncode=7,
            stdout="",
            stderr="PDF_FONT_EMBED_ERROR: fixture",
        )
        with tempfile.TemporaryDirectory() as tmp:
            root = pathlib.Path(tmp)
            spec = root / "06_deck_spec.json"
            registry = root / "02_verified.json"
            output = root / "deck.html"
            stale_pdf = root / "deck.pdf"
            spec.write_text('{"pages": []}', encoding="utf-8")
            registry.write_text('{"sources": {}, "metrics": {}}', encoding="utf-8")
            stale_pdf.write_bytes(b"stale")
            argv = ["render_deck.py", str(spec), str(registry), "-o", str(output), "--unattested"]
            with (
                mock.patch.object(render_deck_module, "render_deck", return_value="<html></html>"),
                mock.patch.object(render_deck_module.subprocess, "run", return_value=capture_failure),
                mock.patch.object(sys, "argv", argv),
            ):
                with self.assertRaises(SystemExit) as raised:
                    render_deck_module.main()

            self.assertFalse(stale_pdf.exists())

        self.assertEqual(raised.exception.code, 7)

    def test_render_main_html_only_writes_html_without_capture(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = pathlib.Path(tmp)
            spec = root / "probe_spec.json"
            registry = root / "probe_registry.json"
            output = root / "probe.html"
            spec.write_text('{"pages": []}', encoding="utf-8")
            registry.write_text('{"sources": {}, "metrics": {}}', encoding="utf-8")
            argv = [
                "render_deck.py", str(spec), str(registry), "-o", str(output),
                "--html-only", "--unattested",
            ]
            with (
                mock.patch.object(render_deck_module, "render_deck", return_value="<html>probe</html>"),
                mock.patch.object(render_deck_module.subprocess, "run") as capture,
                mock.patch.object(sys, "argv", argv),
            ):
                render_deck_module.main()

            self.assertEqual(output.read_text(encoding="utf-8"), "<html>probe</html>")
            capture.assert_not_called()

    def test_render_main_fails_when_capture_script_is_missing(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = pathlib.Path(tmp)
            spec = root / "06_deck_spec.json"
            registry = root / "02_verified.json"
            output = root / "deck.html"
            stale_pdf = root / "deck.pdf"
            spec.write_text('{"pages": []}', encoding="utf-8")
            registry.write_text('{"sources": {}, "metrics": {}}', encoding="utf-8")
            stale_pdf.write_bytes(b"stale")
            argv = ["render_deck.py", str(spec), str(registry), "-o", str(output), "--unattested"]
            with (
                mock.patch.object(render_deck_module, "render_deck", return_value="<html></html>"),
                mock.patch.object(render_deck_module.Path, "exists", return_value=False),
                mock.patch.object(sys, "argv", argv),
            ):
                with self.assertRaises(FileNotFoundError):
                    render_deck_module.main()

            self.assertFalse(stale_pdf.exists())

    def test_render_deck_outputs_closing_layout(self):
        html = render_deck_module.render_deck(VALID_CLOSING_DECK_SPEC, VALID_CONTENT_REGISTRY, title="Closing Fixture")
        self.assertIn("layout-closing", html)
        self.assertIn("closing-body", html)
        self.assertIn("closing-point", html)
        self.assertIn("closing-label", html)
        self.assertIn("도구 · 인프라", html)
        self.assertIn("closing-callout", html)
        self.assertIn("어느 전이가 끝났느냐", html)
        self.assertIn('class="page-number" data-page-number>01 / 01</span>', html)
        self.assertEqual(validate_c6_content_authority(VALID_CLOSING_DECK_SPEC, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_outputs_new_signature_layouts(self):
        specs = {
            "mosaic_tiles": [
                {"type": "headline", "text": "Mosaic Lead"},
                {"type": "body", "text": "First tile copy"},
                {"type": "metric", "metric_id": "metric_click_drop"},
                {"type": "bullets", "items": ["One", "Two"]},
                {"type": "body", "text": "Last tile copy"},
            ],
            "split_status": [
                {"type": "headline", "text": "Status Lead"},
                {"type": "body", "text": "Narrative first"},
                {"type": "bullets", "items": ["Qualitative signal"]},
                {"type": "metric", "metric_id": "metric_click_drop"},
                {"type": "metric", "metric_id": "metric_measurement"},
            ],
            "scenario_cards": [
                {"type": "headline", "text": "Base Case"},
                {"type": "body", "text": "Scenario body"},
                {"type": "metric", "metric_id": "metric_click_drop"},
                {"type": "headline", "text": "Upside Case"},
                {"type": "body", "text": "Scenario body"},
                {"type": "metric", "metric_id": "metric_measurement"},
            ],
        }
        expected_markers = {
            "mosaic_tiles": ("mosaic-body", "mosaic-grid", "mosaic-tile-large"),
            "split_status": ("split-status-body", "status-copy", "status-chip"),
            "scenario_cards": ("scenario-body", "scenario-grid", "scenario-card"),
        }
        for layout, content in specs.items():
            with self.subTest(layout=layout):
                spec = {
                    "pages": [
                        {
                            "page_id": f"{layout}_fixture",
                            "short_title": layout,
                            "layout": layout,
                            "allowed_source_ids": ["src_a", "src_b"],
                            "allowed_metric_ids": ["metric_click_drop", "metric_measurement"],
                            "content": content,
                        }
                    ]
                }
                html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Signature Fixture")
                for marker in expected_markers[layout]:
                    self.assertIn(marker, html)
                self.assertIn('data-metric-id="metric_click_drop"', html)
                self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_outputs_pricing_cards_with_configurable_emphasis(self):
        spec = {
            "pages": [
                {
                    "page_id": "pricing_fixture",
                    "short_title": "Choose The Operating Model",
                    "layout": "pricing_cards",
                    "emphasis_style": "border",
                    "allowed_source_ids": ["src_a", "src_b"],
                    "allowed_metric_ids": ["metric_click_drop", "metric_measurement"],
                    "content": [
                        {"type": "headline", "text": "Starter"},
                        {"type": "bullets", "items": ["Manual reporting", "Email support"]},
                        {"type": "metric", "metric_id": "metric_click_drop"},
                        {"type": "headline", "text": "Operator", "emphasis": True},
                        {"type": "bullets", "items": ["Live dashboard", "Priority review"]},
                        {"type": "metric", "metric_id": "metric_measurement"},
                    ],
                }
            ]
        }
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Pricing Fixture")
        self.assertIn("pricing-body", html)
        self.assertIn("pricing-grid", html)
        self.assertIn("pricing-card-emphasis", html)
        self.assertIn("pricing-emphasis-border", html)
        self.assertIn('data-metric-id="metric_click_drop"', html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_outputs_swot_quad_without_metric_ids(self):
        spec = {
            "pages": [
                {
                    "page_id": "swot_fixture",
                    "short_title": "Strategic Read",
                    "layout": "statement",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": [],
                    "content": [
                        {
                            "type": "viz",
                            "chart": "swot_quad",
                            "title": "Market Position",
                            "series": [
                                {"label": "Strengths", "items": ["Proof library"], "role": "highlight"},
                                {"label": "Weaknesses", "items": ["Small sales surface"]},
                                {"label": "Opportunities", "items": ["Agency bundles"]},
                                {"label": "Threats", "items": ["Platform lock-in"]},
                            ],
                        }
                    ],
                }
            ]
        }
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="SWOT Fixture")
        self.assertIn("visual-swot-quad", html)
        self.assertIn("swot-cell-highlight", html)
        self.assertIn("Proof library", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_outputs_swot_axis_labels_only_when_requested(self):
        content = {
            "type": "viz",
            "chart": "swot_quad",
            "series": [
                {"label": "Strengths", "items": ["Proof library"]},
                {"label": "Weaknesses", "items": ["Small sales surface"]},
                {"label": "Opportunities", "items": ["Agency bundles"]},
                {"label": "Threats", "items": ["Platform lock-in"]},
            ],
        }
        plain_spec = {"pages": [dict(VALID_DECK_SPEC_WITH_VIZ["pages"][0], allowed_metric_ids=[], content=[content])]}
        plain_html = render_deck_module.render_deck(plain_spec, VALID_CONTENT_REGISTRY, title="SWOT Plain")
        self.assertNotIn("swot-axis-arrow", plain_html)

        axis_content = dict(content, axis_labels={"x": ["Low", "High"], "y": ["High", "Low"]})
        axis_spec = {"pages": [dict(VALID_DECK_SPEC_WITH_VIZ["pages"][0], allowed_metric_ids=[], content=[axis_content])]}
        axis_html = render_deck_module.render_deck(axis_spec, VALID_CONTENT_REGISTRY, title="SWOT Axis")
        self.assertIn("swot-axis-arrow", axis_html)
        self.assertIn("High", axis_html)

    def test_render_deck_outputs_hierarchy_pyramid_without_changing_argument_mode(self):
        hierarchy = {
            "type": "viz",
            "chart": "pyramid",
            "pyramid_style": "hierarchy",
            "series": [
                {"label": "Direction"},
                {"label": "Portfolio"},
                {"label": "Execution"},
            ],
        }
        hierarchy_spec = {"pages": [dict(VALID_DECK_SPEC_WITH_VIZ["pages"][0], allowed_metric_ids=[], content=[hierarchy])]}
        hierarchy_html = render_deck_module.render_deck(hierarchy_spec, VALID_CONTENT_REGISTRY, title="Hierarchy")
        self.assertIn("pyramid-hierarchy-rule", hierarchy_html)

        argument = dict(hierarchy)
        argument.pop("pyramid_style")
        argument["series"] = [
            {"role": "claim", "label": "Direction"},
            {"role": "evidence", "label": "Portfolio"},
            {"role": "evidence", "label": "Execution"},
        ]
        argument_spec = {"pages": [dict(VALID_DECK_SPEC_WITH_VIZ["pages"][0], allowed_metric_ids=[], content=[argument])]}
        argument_html = render_deck_module.render_deck(argument_spec, VALID_CONTENT_REGISTRY, title="Argument")
        self.assertNotIn("pyramid-hierarchy-rule", argument_html)
        self.assertIn("<polygon", argument_html)

    def test_render_deck_outputs_gantt_grid_and_milestone_variants(self):
        grid = {
            "type": "viz",
            "chart": "gantt",
            "series": [
                {"label": "Discovery", "start": "2026-08", "end": "2026-09"},
                {"label": "Build", "start": "2026-09", "end": "2026-11"},
            ],
        }
        milestone = {
            "type": "viz",
            "chart": "gantt",
            "series": [
                {"label": "Kickoff", "start": "2026-08", "end": "2026-08", "milestone": True},
                {"label": "Launch", "start": "2026-10", "end": "2026-10", "milestone": True},
            ],
        }
        for content, marker in ((grid, "gantt-grid-bar"), (milestone, "gantt-milestone-diamond")):
            with self.subTest(marker=marker):
                spec = {"pages": [dict(VALID_DECK_SPEC_WITH_VIZ["pages"][0], allowed_metric_ids=[], content=[content])]}
                html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Gantt")
                self.assertIn(marker, html)
                self.assertIn("gantt-grid-rule", html)
                self.assertNotIn('data-metric-id=""', html)
                self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_outputs_pictograph_people_ratio(self):
        content = {
            "type": "viz",
            "chart": "pictograph",
            "series": [{"label": "Adoption", "total": 14, "filled": 9}],
        }
        spec = {"pages": [dict(VALID_DECK_SPEC_WITH_VIZ["pages"][0], allowed_metric_ids=[], content=[content])]}
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Pictograph")
        self.assertEqual(html.count('class="pictograph-person'), 14)
        self.assertIn("9 / 14", html)
        self.assertIn("viz-structured-number", html)
        self.assertNotIn('data-metric-id=""', html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_outputs_semantic_color_roles(self):
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    content=[
                        {
                            "type": "viz",
                            "chart": "before_after",
                            "series": [
                                {"label": "Baseline", "metric_id": "metric_measurement", "role": "baseline"},
                                {"label": "Bad", "metric_id": "metric_click_drop", "role": "negative"},
                                {"label": "Good", "metric_id": "metric_measurement", "role": "positive"},
                                {"label": "Brand", "metric_id": "metric_click_drop", "role": "brand", "color": "#F7931A"},
                            ],
                        }
                    ],
                )
            ]
        }

        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Semantic Fixture")

        self.assertIn("var(--semantic-negative)", html)
        self.assertIn("var(--semantic-positive)", html)
        self.assertIn("#F7931A", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_outputs_quarterly_bars_with_onbar_labels(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                "quarter_a": {"value": "120", "unit": "억", "label": "Q1", "source_ids": ["src_a"]},
                "quarter_b": {"value": "160", "unit": "억", "label": "Q2", "source_ids": ["src_a"]},
                "quarter_c": {"value": "142", "unit": "억", "label": "Q3", "source_ids": ["src_a"]},
                "quarter_d": {"value": "210", "unit": "억", "label": "Q4", "source_ids": ["src_a"]},
            },
        }
        spec = {
            "pages": [
                {
                    "page_id": "quarterly_fixture",
                    "short_title": "Quarterly Bars",
                    "layout": "statement",
                    "allowed_source_ids": ["src_a"],
                    "allowed_metric_ids": ["quarter_a", "quarter_b", "quarter_c", "quarter_d"],
                    "content": [
                        {"type": "viz", "chart": "quarterly_bars", "axis": "hidden", "axis_break": True, "series": [
                            {"metric_id": "quarter_a", "role": "baseline"},
                            {"metric_id": "quarter_b", "role": "baseline"},
                            {"metric_id": "quarter_c", "role": "highlight"},
                            {"metric_id": "quarter_d", "role": "highlight"},
                        ]},
                    ],
                }
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="Quarterly Fixture")

        self.assertIn("visual-quarterly-bars", html)
        self.assertIn("quarter-value-onbar", html)
        self.assertIn("≈", html)
        self.assertIn('data-metric-id="quarter_d"', html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_render_deck_outputs_fin_table_options(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                "revenue_prev": {"value": "120", "unit": "억", "source_ids": ["src_a"]},
                "revenue_now": {"value": "180", "unit": "억", "source_ids": ["src_a"]},
                "profit_prev": {"value": "-12", "unit": "억", "source_ids": ["src_a"]},
                "profit_now": {"value": "24", "unit": "억", "source_ids": ["src_a"]},
                "margin_prev": {"value": "-10", "unit": "%", "source_ids": ["src_a"]},
                "margin_now": {"value": "13", "unit": "%", "source_ids": ["src_a"]},
            },
        }
        spec = {
            "pages": [
                {
                    "page_id": "fin_fixture",
                    "short_title": "Financial Table",
                    "layout": "statement",
                    "allowed_source_ids": ["src_a"],
                    "allowed_metric_ids": [
                        "revenue_prev",
                        "revenue_now",
                        "profit_prev",
                        "profit_now",
                        "margin_prev",
                        "margin_now",
                    ],
                    "content": [
                        {
                            "type": "viz",
                            "chart": "fin_table",
                            "columns": ["이전", "현재", "변화"],
                            "accent_column": 1,
                            "negative_style": "paren",
                            "series": [
                                {"label": "매출", "row_role": "group", "cells": [{"metric_id": "revenue_prev"}, {"metric_id": "revenue_now"}, {"text": "성장"}]},
                                {"label": "영업이익", "row_role": "sub", "cells": [{"metric_id": "profit_prev"}, {"metric_id": "profit_now"}, {"text": "흑자전환"}]},
                                {"label": "영업이익률", "row_role": "ratio", "cells": [{"metric_id": "margin_prev"}, {"metric_id": "margin_now"}, {"text": "개선"}]},
                            ],
                        },
                    ],
                }
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="Financial Fixture")

        self.assertIn("visual-fin-table", html)
        self.assertIn("fin-accent-column", html)
        self.assertIn("fin-row-group", html)
        self.assertIn("fin-row-sub", html)
        self.assertIn("font-style=\"italic\"", html)
        self.assertIn("(12억)", html)
        self.assertIn("var(--semantic-negative)", html)
        self.assertIn("흑자전환", html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_render_deck_outputs_running_head_chrome_for_body_pages(self):
        spec = {
            "meta": {"page_chrome": "running_head", "short_title": "Proof OS"},
            "pages": [
                dict(VALID_COVER_DECK_SPEC["pages"][0]),
                dict(VALID_DECK_SPEC["pages"][0], page_id="body_a"),
                dict(VALID_DECK_SPEC["pages"][1], page_id="body_b", content=[{"type": "eyebrow", "text": "Signal"}, *VALID_DECK_SPEC["pages"][1]["content"]]),
                dict(VALID_CLOSING_DECK_SPEC["pages"][0]),
            ],
        }
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Running Head Fixture")
        body_a = html.split('data-page-id="body_a"', 1)[1].split("</section>", 1)[0]
        body_b = html.split('data-page-id="body_b"', 1)[1].split("</section>", 1)[0]
        cover = html.split('data-page-id="cover"', 1)[1].split("</section>", 1)[0]
        self.assertIn("running-head", body_a)
        # kicker는 명시적 eyebrow만 — 내부 role/layout 명 폴백은 크롬 노출 금지 (7/4 클차장 수정)
        self.assertNotIn('running-head-kicker">HERO_METRIC', body_a)
        self.assertIn('running-head-kicker"></span>', body_a)
        self.assertIn('running-head-kicker">Signal<', body_b)
        self.assertIn("Proof OS", body_a)
        self.assertIn("01 / 02", body_a)
        self.assertIn("PREV", body_b)
        self.assertNotIn("NEXT", body_b)
        self.assertNotIn("page-number", body_a)
        self.assertNotIn("running-head", cover)

    def test_render_deck_outputs_title_band_chrome_for_body_and_divider_pages(self):
        spec = {
            "meta": {"page_chrome": "title_band", "short_title": "Proof OS"},
            "pages": [
                dict(VALID_COVER_DECK_SPEC["pages"][0]),
                dict(VALID_DECK_SPEC["pages"][0], page_id="body_a", short_title="Signal One"),
                {
                    "page_id": "divider_a",
                    "short_title": "Part One",
                    "layout": "divider",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": [],
                    "content": [{"type": "headline", "text": "중앙 스테이트먼트"}],
                },
                dict(VALID_CLOSING_DECK_SPEC["pages"][0]),
            ],
        }
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Title Band Fixture")
        body_a = html.split('data-page-id="body_a"', 1)[1].split("</section>", 1)[0]
        divider_a = html.split('data-page-id="divider_a"', 1)[1].split("</section>", 1)[0]
        cover = html.split('data-page-id="cover"', 1)[1].split("</section>", 1)[0]
        closing = html.split('data-page-id="p15"', 1)[1].split("</section>", 1)[0]

        self.assertIn("title-band", body_a)
        self.assertIn("Signal One", body_a)
        divider_section = re.search(r'<section class="([^"]*)" data-page-id="divider_a">', html)
        self.assertIsNotNone(divider_section)
        self.assertNotIn("page-title-band", divider_section.group(1))
        self.assertNotIn("title-band", divider_a)
        self.assertIn("중앙 스테이트먼트", divider_a)
        self.assertNotIn("title-band", cover)
        self.assertNotIn("title-band", closing)
        self.assertEqual(render_deck_module._title_band_html(""), "")
        self.assertEqual(render_deck_module._title_band_html(" \t\n"), "")
        self.assertEqual(render_deck_module._title_band_html(None), "")

    def test_render_deck_source_caption_is_opt_in_and_uses_registry_short_name(self):
        registry = {
            "sources": {
                "src_a": {
                    **VALID_CONTENT_REGISTRY["sources"]["src_a"],
                    "short_name": "Pew",
                },
                "src_b": VALID_CONTENT_REGISTRY["sources"]["src_b"],
            },
            "metrics": {
                "metric_click_drop": {
                    **VALID_CONTENT_REGISTRY["metrics"]["metric_click_drop"],
                    "period": "1Q26",
                },
                "metric_measurement": VALID_CONTENT_REGISTRY["metrics"]["metric_measurement"],
            },
        }
        spec = {
            "theme": "tech",
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    page_id="caption_fixture",
                    content=[
                        {
                            "type": "viz",
                            "chart": "before_after",
                            "title": "Summary Collapse",
                            "source_caption": "on",
                            "series": [
                                {"label": "Before", "metric_id": "metric_measurement", "role": "baseline"},
                                {"label": "After", "metric_id": "metric_click_drop", "role": "highlight"},
                            ],
                        }
                    ],
                )
            ],
        }

        html = render_deck_module.render_deck(spec, registry, title="Source Caption Fixture")

        self.assertIn("visual-source-caption", html)
        self.assertIn("— 1Q26 · 출처: IAB, Pew", html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_render_deck_source_caption_falls_back_to_source_publisher_prefix(self):
        registry = {
            "sources": {
                "src_a": {
                    "publisher": "식품의약품안전처 (SNUH 미러 게시)",
                    "url": "https://example.com/mfds",
                },
            },
            "metrics": {
                "metric_click_drop": {
                    **VALID_CONTENT_REGISTRY["metrics"]["metric_click_drop"],
                    "source_ids": ["src_a"],
                },
            },
        }
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    page_id="fallback_caption",
                    allowed_source_ids=["src_a"],
                    allowed_metric_ids=["metric_click_drop"],
                    content=[
                        {
                            "type": "viz",
                            "chart": "big_number",
                            "title": "Fallback Caption",
                            "source_caption": "on",
                            "series": [{"metric_id": "metric_click_drop", "role": "highlight"}],
                        }
                    ],
                )
            ],
        }

        html = render_deck_module.render_deck(spec, registry, title="Fallback Caption Fixture")

        self.assertIn("출처: 식품의약품안전처", html)

    def test_render_deck_viz_title_style_band_is_block_local(self):
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    page_id="viz_title_band",
                    content=[
                        {
                            "type": "viz",
                            "chart": "before_after",
                            "title": "Summary Collapse",
                            "title_style": "band",
                            "series": [
                                {"label": "Before", "metric_id": "metric_measurement", "role": "baseline"},
                                {"label": "After", "metric_id": "metric_click_drop", "role": "highlight"},
                            ],
                        }
                    ],
                )
            ],
        }

        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Viz Band Fixture")

        self.assertIn("visual-title-band", html)
        self.assertNotIn("chrome-title-band", html)

    def test_render_deck_renders_derived_metric_value_but_inherits_sources(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_cagr": {
                    "label": "훈련 데이터 연평균 성장률",
                    "value": "8.9",
                    "unit": "%/년",
                    "derivation": "cagr",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                    "period": "2020~2025",
                    "formula_note": "(72/47)^(1/5)-1",
                },
            },
        }
        spec = {
            "pages": [
                {
                    "page_id": "derived_render",
                    "short_title": "Derived Metric",
                    "layout": "statement",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": ["metric_cagr"],
                    "content": [{"type": "metric", "metric_id": "metric_cagr"}],
                }
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="Derived Fixture")

        self.assertIn('data-metric-id="metric_cagr"', html)
        self.assertIn("8.9%/년", html)
        self.assertIn("Pew Research Center", html)
        self.assertIn("IAB", html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_render_deck_outputs_r3_annotations_from_registry_metrics(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_growth": {
                    "label": "성장률",
                    "value": "8.9",
                    "unit": "%/년",
                    "derivation": "cagr",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                    "period": "2020~2025",
                    "formula_note": "(72/47)^(1/5)-1",
                },
            },
        }
        spec = {
            "pages": [
                dict(
                    VALID_DECK_SPEC_WITH_VIZ["pages"][0],
                    page_id="annotated_fixture",
                    allowed_source_ids=[],
                    allowed_metric_ids=["metric_click_drop", "metric_measurement", "metric_growth"],
                    content=[
                        {
                            "type": "viz",
                            "chart": "multi_line",
                            "series": [
                                {"label": "Start", "metric_id": "metric_click_drop", "role": "baseline"},
                                {"label": "End", "metric_id": "metric_measurement", "role": "highlight"},
                            ],
                            "annotations": [
                                {"kind": "callout", "metric_id": "metric_growth", "anchor_series": 1, "shape": "ellipse"},
                                {"kind": "endpoint_value", "series": 1},
                                {"kind": "trend_arrow", "series": 0},
                                {"kind": "event_band", "label": "COVID", "from_key": "Start", "to_key": "End"},
                            ],
                        }
                    ],
                )
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="Annotated Fixture")

        self.assertIn('data-annotation-kind="callout"', html)
        self.assertIn('data-annotation-kind="endpoint_value"', html)
        self.assertIn('data-annotation-kind="trend_arrow"', html)
        self.assertIn('data-annotation-kind="event_band"', html)
        self.assertIn('data-metric-id="metric_growth"', html)
        self.assertIn("8.9%/년", html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_render_deck_outputs_section_nav_from_divider_labels(self):
        spec = {
            "meta": {"section_nav": "chips"},
            "pages": [
                dict(VALID_COVER_DECK_SPEC["pages"][0]),
                {
                    "page_id": "divider_a",
                    "short_title": "증거 — 시장 신호",
                    "layout": "divider",
                    "section_index": 1,
                    "section_label": "증거",
                    "section_nav_label": "designer free text must not render",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": [],
                    "content": [{"type": "headline", "text": "증거 — 시장 신호"}],
                },
                dict(VALID_DECK_SPEC["pages"][0], page_id="body_a"),
                {
                    "page_id": "divider_b",
                    "short_title": "행동 — 운영 전환",
                    "layout": "divider",
                    "section_index": 2,
                    "section_label": "행동",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": [],
                    "content": [{"type": "headline", "text": "행동 — 운영 전환"}],
                },
            ],
        }

        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Section Nav Fixture")

        self.assertIn("section-nav section-nav-chips", html)
        self.assertIn("증거", html)
        self.assertIn("행동", html)
        self.assertNotIn("designer free text must not render", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_outputs_metric_commentary_layout(self):
        registry = {
            "sources": VALID_CONTENT_REGISTRY["sources"],
            "metrics": {
                **VALID_CONTENT_REGISTRY["metrics"],
                "metric_delta_yoy": {
                    "label": "YoY 변화",
                    "value": "53.2",
                    "unit": "%",
                    "derivation": "delta_pct",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                    "formula_note": "(72-47)/47*100",
                },
                "metric_delta_qoq": {
                    "label": "QoQ 변화",
                    "value": "53.2",
                    "unit": "%",
                    "derivation": "delta_pct",
                    "derived_from": ["metric_click_drop", "metric_measurement"],
                    "source_ids": [],
                    "status": "derived",
                    "formula_note": "(72-47)/47*100",
                },
            },
        }
        spec = {
            "pages": [
                {
                    "page_id": "metric_commentary_fixture",
                    "short_title": "Revenue Momentum",
                    "layout": "metric_commentary",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": ["metric_click_drop", "metric_measurement", "metric_delta_yoy", "metric_delta_qoq"],
                    "rows": [
                        {
                            "heading_metric_id": "metric_measurement",
                            "headline_metric_id": "metric_delta_yoy",
                            "bullets": [
                                {"label": "YoY", "metric_id": "metric_delta_yoy"},
                                {"label": "QoQ", "metric_id": "metric_delta_qoq"},
                            ],
                            "chart": {
                                "chart": "quarterly_bars",
                                "series": [
                                    {"metric_id": "metric_click_drop", "role": "baseline"},
                                    {"metric_id": "metric_measurement", "role": "highlight"},
                                ],
                            },
                        }
                    ],
                    "content": [],
                }
            ]
        }

        html = render_deck_module.render_deck(spec, registry, title="Metric Commentary Fixture")

        self.assertIn("metric-commentary-body", html)
        self.assertIn("metric-commentary-headline", html)
        self.assertIn("commentary-bullet-label", html)
        self.assertIn("visual-quarterly-bars", html)
        self.assertIn("18%", html)
        self.assertEqual(validate_c6_content_authority(spec, registry, html), [])

    def test_render_deck_outputs_side_wordmark_from_context(self):
        spec = {
            "meta": {"short_title": "Proof OS"},
            "pages": [
                dict(
                    VALID_DECK_SPEC["pages"][0],
                    page_id="wordmark_fixture",
                    decor="side_wordmark",
                    section_label="Evidence",
                )
            ],
        }
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Side Wordmark Fixture")
        self.assertIn("decor-side-wordmark", html)
        self.assertIn("side-wordmark", html)
        self.assertIn("Evidence", html)
        self.assertEqual(validate_c6_content_authority(spec, VALID_CONTENT_REGISTRY, html), [])

    def test_render_deck_omits_eyebrow_when_no_explicit_block(self):
        # 명시 eyebrow 없으면 role/layout 내부명 폴백 노출 금지(C2 계열) — div 자체를 렌더하지 않는다(7/4).
        html = render_deck_module.render_deck(VALID_DECK_SPEC, VALID_CONTENT_REGISTRY, title="No Eyebrow Fixture")
        self.assertNotIn('<div class="eyebrow', html)
        self.assertNotIn("HERO_METRIC", html)
        self.assertNotIn("STAT_GRID", html)

    def test_render_deck_supports_eyebrow_blocks_as_section_labels(self):
        html = render_deck_module.render_deck(VALID_DECK_SPEC_WITH_EYEBROW, VALID_CONTENT_REGISTRY, title="Eyebrow Fixture")
        self.assertIn('<div class="eyebrow">Chapter 1</div>', html)
        self.assertNotIn("unsupported content block type: eyebrow", html)
        self.assertEqual(validate_c6_content_authority(VALID_DECK_SPEC_WITH_EYEBROW, VALID_CONTENT_REGISTRY, html), [])

    def test_renderer_and_contracts_share_supported_content_block_types(self):
        self.assertTrue(hasattr(contract_checks_module, "SUPPORTED_CONTENT_BLOCK_TYPES"))
        self.assertIs(
            render_deck_module.SUPPORTED_CONTENT_BLOCK_TYPES,
            contract_checks_module.SUPPORTED_CONTENT_BLOCK_TYPES,
        )
        self.assertIn("eyebrow", contract_checks_module.SUPPORTED_CONTENT_BLOCK_TYPES)
        self.assertIn("viz", contract_checks_module.SUPPORTED_CONTENT_BLOCK_TYPES)
        self.assertIn("text_table", contract_checks_module.SUPPORTED_CONTENT_BLOCK_TYPES)

    def test_every_supported_viz_chart_has_a_renderer(self):
        # 계약 enum에 있는데 렌더러 분기가 없으면 "조용한 no-op"이 된다 — 1:1을 코드로 강제.
        self.assertEqual(
            set(render_deck_module._CHART_RENDERERS),
            set(contract_checks_module.SUPPORTED_VIZ_CHART_TYPES),
        )

    def test_external_review_lenses_metadata_matches_prompt_lenses(self):
        prompt = external_review_module.build_review_prompt([{"page_id": "p01", "text": "본문"}])
        self.assertEqual(len(external_review_module.LENSES), 5)
        for lens in external_review_module.LENSES:
            self.assertIn(lens, prompt)
        self.assertTrue(any("독자 패널" in lens for lens in external_review_module.LENSES))

    def test_external_review_resolves_gemini_env_overrides_and_checks_existence(self):
        with tempfile.TemporaryDirectory() as tmp:
            wrapper = pathlib.Path(tmp) / "gemini_call_wrapper.py"
            python = pathlib.Path(tmp) / "python"
            wrapper.write_text("# wrapper\n", encoding="utf-8")
            python.write_text("# python\n", encoding="utf-8")
            with mock.patch.dict(os.environ, {"GEMINI_WRAPPER": str(wrapper), "GEMINI_PY": str(python)}):
                resolved_python, resolved_wrapper = external_review_module.resolve_gemini_paths()
            self.assertEqual(resolved_python, python)
            self.assertEqual(resolved_wrapper, wrapper)

        with mock.patch.dict(os.environ, {"GEMINI_WRAPPER": "/missing/wrapper.py", "GEMINI_PY": "/missing/python"}):
            with self.assertRaises(FileNotFoundError) as ctx:
                external_review_module.resolve_gemini_paths()
        self.assertIn("GEMINI", str(ctx.exception))

    def test_external_review_gemini_command_reads_prompt_file_without_prompt_in_argv(self):
        with tempfile.TemporaryDirectory() as tmp:
            prompt_path = pathlib.Path(tmp) / "prompt.txt"
            prompt_text = "매우 긴 프롬프트 본문"
            prompt_path.write_text(prompt_text, encoding="utf-8")

            command = external_review_module.build_gemini_command(
                pathlib.Path("/bin/python3"),
                pathlib.Path("/tmp/gemini_call_wrapper.py"),
                prompt_path,
            )

        self.assertIn(str(prompt_path), command)
        self.assertIn("-c", command)
        self.assertNotIn(prompt_text, "\0".join(command))
        self.assertNotIn("--prompt", command)

    def test_external_review_codex_command_keeps_argv_but_enforces_size_limit(self):
        command = external_review_module.build_codex_command("짧은 프롬프트")
        self.assertEqual(command[:3], ["codex", "exec", "--skip-git-repo-check"])
        self.assertIn("짧은 프롬프트", command)
        with self.assertRaises(ValueError):
            external_review_module.build_codex_command("x" * (external_review_module.CODEX_PROMPT_ARGV_LIMIT + 1))

    @unittest.skipUnless(R4_RUN_DIR.is_dir(), f"R4 실물 런 없음: {R4_RUN_DIR}")
    def test_factcheck_dump_builds_crosswalk_table_from_real_r4_run(self):
        # R4는 읽기 전용 — 08 산출은 tempdir 사본에만 쓴다.
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            shutil.copy(R4_RUN_DIR / "06_deck_spec.json", run_dir / "06_deck_spec.json")
            shutil.copy(R4_RUN_DIR / "02_verified.json", run_dir / "02_verified.json")

            exit_code = factcheck_dump_module.main(["--run", str(run_dir)])
            self.assertEqual(exit_code, 0)

            table_path = run_dir / "08_factcheck_table.json"
            self.assertTrue(table_path.exists())
            rows = json.loads(table_path.read_text(encoding="utf-8"))

        self.assertGreater(len(rows), 0)
        for row in rows:
            self.assertIn("metric_id", row)
            self.assertIn("page_no", row)
            self.assertTrue(
                row.get("source_url") or row.get("local_path") or row.get("source_missing"),
                f"row without a source pointer or source_missing flag: {row}",
            )

    @unittest.skipUnless(R4_RUN_DIR.is_dir(), f"R4 실물 런 없음: {R4_RUN_DIR}")
    def test_external_review_spec_stage_prompt_includes_real_page_text(self):
        # LLM 호출 없이 프롬프트 구성까지만 — R4 06_deck_spec.json은 읽기만 한다.
        pages = external_review_module.extract_page_texts_from_spec(R4_RUN_DIR)
        self.assertGreater(len(pages), 0)
        prompt = external_review_module.build_review_prompt(pages)

        deck_spec = json.loads((R4_RUN_DIR / "06_deck_spec.json").read_text(encoding="utf-8"))
        cover_title = deck_spec["pages"][0]["short_title"]
        self.assertIn(cover_title, prompt)
        # 06_deck_spec.json에는 metric 값이 없다 — 02_verified.json metric_registry 해석까지 됐는지 확인.
        self.assertIn("2823억원", prompt)  # metric_127 (세라젬 연결 매출액, 2018)

    def test_pptx_export_run_checked_reports_timeout_with_configured_seconds(self):
        with mock.patch.object(
            pptx_export_module.subprocess,
            "run",
            side_effect=subprocess.TimeoutExpired(["pip"], 300),
        ) as run:
            with self.assertRaises(RuntimeError) as ctx:
                pptx_export_module.run_checked(["pip"], "pip install", timeout=300)

        self.assertEqual(run.call_args.kwargs["timeout"], 300)
        self.assertIn("pip install timed out after 300s", str(ctx.exception))

    def test_pptx_export_run_chrome_reports_timeout_with_configured_seconds(self):
        with mock.patch.object(
            pptx_export_module.subprocess,
            "run",
            side_effect=subprocess.TimeoutExpired(["chrome"], 180),
        ) as run:
            with self.assertRaises(SystemExit) as ctx:
                pptx_export_module.run_chrome("chrome", ["--dump-dom"], "LAYOUT_DUMP")

        self.assertEqual(run.call_args.kwargs["timeout"], 180)
        self.assertIn("LAYOUT_DUMP_TIMEOUT: Chrome timed out after 180s", str(ctx.exception))

    def test_pptx_export_maps_viz_to_native_bar_column_and_line_charts(self):
        self.assertEqual(pptx_export_module.native_chart_kind("dumbbell"), "bar")
        self.assertEqual(pptx_export_module.native_chart_kind("rising_columns"), "column")
        self.assertEqual(pptx_export_module.native_chart_kind("multi_line"), "line")

    def test_pptx_export_hides_only_native_chart_cards_from_background(self):
        script = pptx_export_module.hide_picked_text_script()
        self.assertIn(".visual-multi-line", script)
        self.assertIn(".visual-dumbbell", script)
        self.assertNotIn("querySelectorAll('.visual-card')", script)
        self.assertNotIn(".visual-donut", script)

    def test_pptx_export_uses_qa_safe_font(self):
        self.assertEqual(pptx_export_module.FONT_NAME, "Arial")

    def test_pptx_export_builds_customer_safe_sources_note_from_page_allowlist(self):
        page = {"allowed_source_ids": ["src_a"]}
        registry = {
            "source_registry": {
                "src_a": {
                    "publisher": "Pew Research Center",
                    "title": "AI at Work",
                    "url": "https://example.com/report",
                    "local_path": "_workspace/run/collector/src_a.pdf",
                    "conditions": "Survey evidence only.",
                }
            }
        }
        note = pptx_export_module.sources_note_text(page, registry)
        self.assertTrue(note.startswith("[Sources]\n"))
        self.assertIn("Pew Research Center — AI at Work", note)
        self.assertIn("https://example.com/report", note)
        self.assertNotIn("_workspace/", note)
        self.assertNotIn("Survey evidence only.", note)

    def test_pptx_export_returns_no_sources_note_for_empty_allowlist(self):
        self.assertEqual(pptx_export_module.sources_note_text({"allowed_source_ids": []}, {}), "")
        internal_registry = {"source_registry": {"src_a": {"title": "verifier pool src_a"}}}
        self.assertEqual(
            pptx_export_module.sources_note_text({"allowed_source_ids": ["src_a"]}, internal_registry),
            "",
        )

    def test_pptx_export_metric_number_uses_first_numeric_token_and_rejects_missing(self):
        self.assertEqual(pptx_export_module._metric_number({"value": "2026-08-02"}, "metric_date"), 2026.0)
        with self.assertRaises(SystemExit) as ctx:
            pptx_export_module._metric_number({}, "metric_missing")
        self.assertEqual(str(ctx.exception), "METRIC_NOT_NUMERIC: metric_missing")

    def test_pptx_export_rejects_mixed_units_in_one_chart(self):
        block = {
            "chart": "multi_line",
            "series": [
                {"label": "A", "metric_id": "metric_a", "role": "highlight"},
                {"label": "B", "metric_id": "metric_b", "role": "baseline"},
            ],
        }
        registry = {
            "metric_registry": {
                "metric_a": {"value": "10%", "unit": "%"},
                "metric_b": {"value": "20억원", "unit": "억원"},
            }
        }
        with self.assertRaises(SystemExit) as ctx:
            pptx_export_module._native_chart_data(block, registry)
        self.assertIn("METRIC_UNIT_MISMATCH", str(ctx.exception))

    def test_pptx_export_native_chart_data_splits_multi_line_by_role(self):
        block = {
            "chart": "multi_line",
            "series": [
                {"label": "2024", "metric_id": "metric_a", "role": "highlight"},
                {"label": "2025", "metric_id": "metric_b", "role": "highlight"},
                {"label": "2024", "metric_id": "metric_c", "role": "baseline"},
                {"label": "2025", "metric_id": "metric_d", "role": "baseline"},
            ],
        }
        registry = {
            "metric_registry": {
                key: {"value": str(value), "unit": "%"}
                for key, value in {
                    "metric_a": 10,
                    "metric_b": 20,
                    "metric_c": 30,
                    "metric_d": 40,
                }.items()
            }
        }
        chart_data, series_count = pptx_export_module._native_chart_data(block, registry)
        self.assertEqual([category.label for category in chart_data.categories], ["2024", "2025"])
        self.assertEqual(series_count, 2)
        self.assertEqual([series.name for series in chart_data], ["highlight", "baseline"])

    def test_pptx_export_verify_rejects_wrong_multi_line_series_and_unsafe_notes(self):
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        deck_spec = {
            "pages": [
                {
                    "page_id": "p01",
                    "allowed_source_ids": ["src_a"],
                    "content": [
                        {
                            "type": "viz",
                            "chart": "multi_line",
                            "series": [
                                {"metric_id": "metric_a", "role": "highlight"},
                                {"metric_id": "metric_b", "role": "baseline"},
                            ],
                        }
                    ],
                }
            ]
        }
        layout = {"slides": [{"page_id": "p01", "boxes": [], "viz_boxes": [{"chart": "multi_line"}]}]}
        with tempfile.TemporaryDirectory() as td:
            path = pathlib.Path(td) / "bad.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            data = ChartData()
            data.categories = ["A", "B"]
            data.add_series("merged", [1, 2])
            slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1), Inches(1), Inches(5), Inches(3), data)
            slide.notes_slide.notes_text_frame.text = "[Sources]\n\n_workspace/run/src_a.pdf"
            prs.save(path)

            with self.assertRaises(SystemExit) as ctx:
                pptx_export_module.verify_pptx(path, layout, deck_spec, require_background_and_text=False)
        self.assertIn("multi_line series=1, expected 2", str(ctx.exception))

    def test_pptx_export_verify_rejects_internal_path_and_header_only_notes(self):
        from pptx import Presentation

        layout = {"slides": [{"page_id": "p01", "boxes": [], "viz_boxes": []}]}
        for note, expected in (
            ("[Sources]\n\n_workspace/run/src_a.pdf", "notes contain _workspace/"),
            ("[Sources]", "notes contain header only"),
        ):
            with self.subTest(note=note), tempfile.TemporaryDirectory() as td:
                path = pathlib.Path(td) / "bad-notes.pptx"
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.notes_slide.notes_text_frame.text = note
                prs.save(path)

                with self.assertRaises(SystemExit) as ctx:
                    pptx_export_module.verify_pptx(path, layout, require_background_and_text=False)
            self.assertIn(expected, str(ctx.exception))

    def test_render_deck_embeds_local_image_as_data_uri_on_cover(self):
        spec = {
            "pages": [
                {
                    "page_id": "p01",
                    "short_title": "Cover",
                    "layout": "cover",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": [],
                    "content": [
                        {"type": "headline", "text": "Image demo"},
                        {"type": "image", "asset": "hero.png", "alt": "Abstract office"},
                    ],
                }
            ]
        }
        with tempfile.TemporaryDirectory() as td:
            asset = pathlib.Path(td) / "hero.png"
            asset.write_bytes(b"\x89PNG\r\n\x1a\nfixture")
            html = render_deck_module.render_deck(
                spec,
                VALID_CONTENT_REGISTRY,
                asset_base=pathlib.Path(td),
            )
        self.assertIn('class="local-image"', html)
        self.assertIn('src="data:image/png;base64,', html)
        self.assertIn('alt="Abstract office"', html)

    def test_render_deck_rejects_local_image_outside_cover_or_divider(self):
        spec = {
            "pages": [
                {
                    "page_id": "p02",
                    "short_title": "Body",
                    "layout": "statement",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": [],
                    "content": [{"type": "image", "asset": "hero.png", "alt": "No"}],
                }
            ]
        }
        with self.assertRaisesRegex(ValueError, "image blocks are only allowed"):
            render_deck_module.render_deck(
                spec,
                VALID_CONTENT_REGISTRY,
                asset_base=pathlib.Path("."),
            )

    def test_qa_lint_treats_text_table_rows_as_body_text(self):
        spec = {
            "archetype": "selfcheck",
            "pages": [
                {
                    "page_id": "table_rows",
                    "short_title": "표",
                    "layout": "statement",
                    "allowed_metric_ids": [],
                    "content": [
                        {
                            "type": "text_table",
                            "columns": ["구분", "해석"],
                            "rows": [["단, 관찰되지 않는다", "성장 47%"]],
                        }
                    ],
                }
            ],
        }

        codes = {defect["code"] for defect in qa_lint_module.lint_deck(spec, {})}

        self.assertIn("READER_FIRST_CAVEAT", codes)
        self.assertIn("READER_FIRST_EPISTEMIC", codes)
        self.assertIn("RAW_NUMBER_IN_LABEL", codes)

    def test_run_contracts_defaults_to_deck_html_not_latest_html(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            deck_html = run_dir / "deck.html"
            later_html = run_dir / "z_later.html"
            deck_html.write_text("<section>deck</section>", encoding="utf-8")
            later_html.write_text("<section>later</section>", encoding="utf-8")
            os.utime(deck_html, (1, 1))
            os.utime(later_html, (2, 2))

            selected = run_contracts_module.select_html_path(run_dir, None)

        self.assertEqual(selected, deck_html)

    def test_run_contracts_applies_c9_to_explicit_html(self):
        with tempfile.TemporaryDirectory() as tmp:
            run_dir = pathlib.Path(tmp)
            custom_html = run_dir / "custom.html"
            custom_html.write_text("<section>custom</section>", encoding="utf-8")
            with (
                mock.patch.object(sys, "argv", ["run_contracts.py", str(run_dir), str(custom_html)]),
                mock.patch.object(run_contracts_module, "validate_all_contracts", return_value=[]),
                mock.patch.object(run_contracts_module, "check_c9_final_review", return_value=[]) as check_c9,
                mock.patch.object(run_contracts_module, "check_c11_source_coverage", return_value=[]),
                mock.patch.object(run_contracts_module, "check_c14_viz_intent_preserved", return_value=[]),
                mock.patch.object(run_contracts_module, "check_c15_page_count_ceiling", return_value=[]),
                mock.patch("builtins.print") as output,
            ):
                result = run_contracts_module.main()

        self.assertEqual(result, 0)
        check_c9.assert_called_once_with(run_dir)
        self.assertNotIn(
            "N/A C9 final_review 검사",
            "\n".join(" ".join(map(str, call.args)) for call in output.call_args_list),
        )

    def test_run_deck_parses_run_id_from_claude_result_before_workspace_fallback(self):
        command = (
            f'TICKDECK_RUN_DECK_TEST=1 source "{RUN_DECK_SH}"; '
            "printf '%s\\n' '중간 로그' '20260705_clo_market' | parse_result_run_id"
        )
        completed = subprocess.run(["bash", "-c", command], capture_output=True, text=True, check=False)
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertEqual(completed.stdout.strip(), "20260705_clo_market")

    def test_capture_deck_includes_fit_band_overflow_gate(self):
        script = CAPTURE_DECK_SH.read_text(encoding="utf-8")

        self.assertIn("FIT_BAND_OVERFLOW", script)
        self.assertIn("bandovf", script)

    def test_capture_deck_includes_fit_annotation_overlap_gate(self):
        script = CAPTURE_DECK_SH.read_text(encoding="utf-8")

        self.assertIn("FIT_ANNOTATION_OVERLAP", script)
        self.assertIn("annotationOverlap", script)

    def test_validate_all_contracts_can_raise(self):
        broken = dict(VALID_DECK, rendered_pages=[{"title": "신뢰도 강등", "body": "본문"}])
        with self.assertRaises(ContractViolation):
            validate_all_contracts(broken, raise_on_error=True)

    # ── 20260810 소형 게이트 6종 (codex_gates_task.md) ──────────────────────

    def test_gate1_qa_ink_distribution_fails_and_exempts_non_body_layouts(self):
        # 게이트 1 — 어제 사고 재현: cover/closing은 잉크가 적어도 정상이라 본문 판정에서 뺀다.
        # 본문(statement) 3장은 전부 저밀도(1%)라 분포 기준·중앙값 기준 둘 다 FAIL이어야 한다.
        with tempfile.TemporaryDirectory() as td:
            run_dir = pathlib.Path(td)
            pdf_path = run_dir / "deck.pdf"
            pdf_path.write_bytes(b"%PDF-fake")
            spec = {
                "pages": [
                    {"layout": "cover"},
                    {"layout": "statement"},
                    {"layout": "statement"},
                    {"layout": "statement"},
                    {"layout": "closing"},
                ]
            }
            (run_dir / "06_deck_spec.json").write_text(json.dumps(spec), encoding="utf-8")

            fake_ratios = {1: 0.50, 2: 0.01, 3: 0.01, 4: 0.01, 5: 0.50}

            def fake_run(cmd, capture_output=True, text=True):
                prefix = pathlib.Path(cmd[-1])
                for idx in fake_ratios:
                    (prefix.parent / f"{prefix.name}-{idx}.png").write_bytes(b"fake")
                return subprocess.CompletedProcess(cmd, 0, "", "")

            def fake_ink_ratio(path):
                return fake_ratios[int(re.search(r"-(\d+)\.png$", path.name).group(1))]

            with (
                mock.patch.object(qa_ink_module.subprocess, "run", side_effect=fake_run),
                mock.patch.object(qa_ink_module, "ink_ratio", side_effect=fake_ink_ratio),
            ):
                lines, is_fail = qa_ink_module.check_pdf(pdf_path)

        self.assertTrue(is_fail)
        joined = "\n".join(lines)
        self.assertIn("INK_SPARSE_FAIL", joined)
        self.assertIn("p2", joined)
        self.assertIn("p3", joined)
        self.assertIn("p4", joined)
        # cover(p1)/closing(p5)는 본문이 아니므로 걸린 페이지 목록에 없어야 한다.
        self.assertNotIn("p1(", joined)
        self.assertNotIn("p5(", joined)

    def test_gate2_c14_flags_lost_chart_intent_deck_level(self):
        # 게이트 2 — page-plan이 예고한 차트 의도(p05·p09) 대비 spec viz=0이면 C14.
        page_plan = {
            "pages": [
                {"page_id": "p05", "layout_hint": "split", "content_notes": "차트는 한 개만"},
                {"page_id": "p09", "layout_hint": "statement", "content_notes": "매출 3분할 chart"},
                {"page_id": "p20", "layout_hint": "cover", "content_notes": "표지. 수치·출처 없음."},
            ]
        }
        deck_spec_no_viz = {"pages": [{"content": [{"type": "body", "text": "x"}]}]}
        violations = check_c14_viz_intent_preserved(page_plan, deck_spec_no_viz)
        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C14")
        self.assertIn("p05", violations[0].message)
        self.assertIn("p09", violations[0].message)

        deck_spec_with_viz = {
            "pages": [{"content": [{"type": "viz", "chart": "donut"}, {"type": "viz", "chart": "bar"}]}]
        }
        self.assertEqual(check_c14_viz_intent_preserved(page_plan, deck_spec_with_viz), [])
        # page_plan 없음(빈 입력) → 비활성, 기존 워크스페이스 호환
        self.assertEqual(check_c14_viz_intent_preserved({}, deck_spec_no_viz), [])

        # 2026-08-10 재설계: 부분 대체(의도 2장 중 viz 1개뿐, 나머지는 표)는 통과한다.
        # 데이터가 적은 페이지는 표/큰 숫자 카드가 차트보다 더 정확할 수 있어, 개수 미달
        # 자체를 위반으로 잡지 않는다 — 잡는 것은 전면 소실(0개)뿐이다.
        deck_spec_partial_viz = {
            "pages": [{"content": [{"type": "viz", "chart": "donut"}, {"type": "body", "text": "표로 대체"}]}]
        }
        self.assertEqual(check_c14_viz_intent_preserved(page_plan, deck_spec_partial_viz), [])

    def test_gate3_c15_flags_page_inflation_over_ceiling(self):
        # 게이트 3 — 어제 사고 수치 그대로: plan 28장 → spec 41장(상한 ceil(28*1.2)=34) 초과.
        page_plan = {"pages": [{"page_id": f"p{i}"} for i in range(28)]}
        deck_spec_41 = {"pages": [{"page_id": f"p{i}"} for i in range(41)]}
        violations = check_c15_page_count_ceiling(page_plan, deck_spec_41)
        self.assertEqual(len(violations), 1)
        self.assertEqual(violations[0].contract_id, "C15")
        self.assertIn("28", violations[0].message)
        self.assertIn("41", violations[0].message)
        self.assertIn("34", violations[0].message)

        deck_spec_ok = {"pages": [{"page_id": f"p{i}"} for i in range(34)]}
        self.assertEqual(check_c15_page_count_ceiling(page_plan, deck_spec_ok), [])
        self.assertEqual(check_c15_page_count_ceiling({}, deck_spec_41), [])

    def test_c13_role_duplication_passes_distinct_slots_and_rejects_v1_duplicates(self):
        passing = {
            "pages": [{
                "page_id": "p01",
                "short_title": "시장 전환",
                "content": [
                    {"type": "headline", "text": "구매 기준이 가격에서 증거로 이동한다"},
                    {"type": "body", "text": "후기와 검증 자료가 선택을 뒷받침한다."},
                    {"type": "callout", "text": "증거 자산을 먼저 쌓아야 한다"},
                    {"type": "text_table", "title": "채널별 실행 우선순위", "columns": ["채널", "순위"], "rows": []},
                ],
            }]
        }
        self.assertEqual(check_c13_role_duplication(passing), [])

        failing = {
            "pages": [{
                "page_id": "p01",
                "short_title": "시장 전환",
                "content": [
                    {"type": "headline", "text": "시장 전환"},
                    {"type": "body", "text": "후기를 먼저 확보한다. 채널을 넓힌다."},
                    {"type": "callout", "text": "후기를 먼저 확보한다"},
                    {"type": "text_table", "title": "시장 전환", "columns": ["채널", "순위"], "rows": []},
                ],
            }]
        }
        violations = check_c13_role_duplication(failing)
        self.assertEqual(len(violations), 3)
        self.assertTrue(all(v.contract_id == "C13" for v in violations))

        long_headline = {
            "pages": [{
                "page_id": "p01",
                "short_title": "개요",
                "content": [{"type": "headline", "text": "사업 개요 및 향후 비전과 전략 로드맵 전반"}],
            }]
        }
        self.assertEqual(check_c13_role_duplication(long_headline, threshold=0.95), [])

    def test_page_visual_intent_groups_split_children_and_rejects_self_approval(self):
        plan = {"pages": [{"page_id": "plan-1", "visual_intent": {"intent": "chart", "desc": "추이"}}]}
        split_spec = {
            "meta": {"spec_author": "designer-agent"},
            "pages": [
                {"page_id": "p01", "plan_id": "plan-1", "split_seq": 1, "split_total": 2, "content": [{"type": "body", "text": "설명"}]},
                {"page_id": "p02", "plan_id": "plan-1", "split_seq": 2, "split_total": 2, "content": [{"type": "viz", "chart": "donut"}]},
            ],
        }
        self.assertEqual(check_page_visual_intent_preserved(plan, split_spec), [])

        self_approved = {
            "meta": {"spec_author": "designer-agent"},
            "pages": [{
                "page_id": "p01",
                "plan_id": "plan-1",
                "content": [{"type": "body", "text": "표로 강등"}],
                "visual_downgrade": {"reason": "값이 적음", "approved_by": "designer-agent", "date": "2026-08-11"},
            }],
        }
        violations = check_page_visual_intent_preserved(plan, self_approved)
        self.assertEqual(len(violations), 1)
        self.assertIn("자기 승인", violations[0].message)

        self_approved["pages"][0]["visual_downgrade"]["approved_by"] = "본부"
        self.assertEqual(check_page_visual_intent_preserved(plan, self_approved), [])
        del self_approved["meta"]["spec_author"]
        violations = check_page_visual_intent_preserved(plan, self_approved)
        self.assertEqual(len(violations), 1)
        self.assertIn("[미규명]", violations[0].message)

        unknown_authors = [None, "", "   ", 7, ["designer-agent"], {"name": "designer-agent"}]
        for unknown_author in unknown_authors:
            with self.subTest(spec_author=unknown_author):
                self_approved["meta"]["spec_author"] = unknown_author
                violations = check_page_visual_intent_preserved(plan, self_approved)
                self.assertEqual(len(violations), 1)
                self.assertIn("[미규명]", violations[0].message)

        split_downgrade = {
            "meta": {"spec_author": "designer-agent"},
            "pages": [
                {
                    "page_id": "p01", "plan_id": "plan-1", "split_seq": 1, "split_total": 2,
                    "content": [{"type": "body", "text": "설명"}],
                    "visual_downgrade": {"reason": "값이 적음", "approved_by": "본부", "date": "2026-08-11"},
                },
                {
                    "page_id": "p02", "plan_id": "plan-1", "split_seq": 2, "split_total": 2,
                    "content": [{"type": "body", "text": "결론"}],
                },
            ],
        }
        violations = check_page_visual_intent_preserved(plan, split_downgrade)
        self.assertEqual(len(violations), 1)
        self.assertIn("승인이 없음", violations[0].message)
        split_downgrade["pages"][1]["visual_downgrade"] = {
            "reason": "값이 적음", "approved_by": "본부", "date": "2026-08-11",
        }
        self.assertEqual(check_page_visual_intent_preserved(plan, split_downgrade), [])
        split_downgrade["pages"][1]["content"].append({"type": "viz", "chart": "donut"})
        del split_downgrade["pages"][0]["visual_downgrade"]
        del split_downgrade["pages"][1]["visual_downgrade"]
        self.assertEqual(check_page_visual_intent_preserved(plan, split_downgrade), [])

        malformed_plan = {"pages": [{"page_id": "plan-1", "visual_intent": {"intent": "graph", "desc": "오타"}}]}
        self.assertTrue(check_page_visual_intent_preserved(malformed_plan, split_spec))
        self.assertTrue(check_page_visual_intent_preserved({}, split_spec))

        invalid_date = {
            "meta": {"spec_author": "본부"},
            "pages": [{
                "page_id": "p01",
                "plan_id": "plan-1",
                "content": [{"type": "body", "text": "강등"}],
                "visual_downgrade": {"reason": "값이 적음", "approved_by": "headquarters", "date": "2026-99-99"},
            }],
        }
        violations = check_page_visual_intent_preserved(plan, invalid_date)
        self.assertTrue(any("date" in v.message or "자기 승인" in v.message for v in violations))

    def test_deck_gates_each_have_pass_and_fail_cases(self):
        def page(page_id, *, visual=True, chars=100, layout="stack"):
            content = [{"type": "body", "text": "가" * chars}]
            if visual:
                content.append({"type": "metric", "metric_id": "m1"})
            return {"page_id": page_id, "layout": layout, "content": content}

        plan = {"pages": [{"page_id": f"plan-{i}"} for i in range(10)]}
        intake = {}
        registry = {
            "sources": {"s1": {"title": "근거", "url": "https://example.com"}},
            "metrics": {"m1": {"value": 1, "unit": "개", "source_ids": ["s1"]}},
        }
        passing_pages = [page(f"p{i}", visual=i % 3 != 0, chars=100 + i) for i in range(1, 11)]
        passing_pages[0]["content"].append({"type": "metric_grid"})
        passing_budgets = [
            {"page_id": f"p{i}", "verdict": "FIT", "height_px": 400, "capacity_px": 600}
            for i in range(1, 11)
        ]
        passing = check_deck_spec_gates(
            plan, {"pages": passing_pages}, intake, registry, calibration={}, layout_results=passing_budgets
        )
        self.assertEqual(passing.violations, [])
        self.assertEqual(passing.warnings, [])

        boundary_pages = [page(f"p{i}", visual=i <= 5) for i in range(1, 11)]
        boundary_budgets = [
            {
                "page_id": f"p{i}",
                "verdict": "SPARSE" if i <= 2 else "FIT",
                "height_px": 360,
                "capacity_px": 600,
            }
            for i in range(1, 11)
        ]
        boundary = check_deck_spec_gates(
            plan, {"pages": boundary_pages}, intake, registry, {}, layout_results=boundary_budgets
        )
        self.assertFalse(any("SPARSE 비율" in v.message for v in boundary.violations))
        self.assertFalse(any("밀도 중앙값" in v.message for v in boundary.violations))
        self.assertFalse(any("시각 포함 비율" in v.message for v in boundary.violations))

        inflated = check_deck_spec_gates(
            plan,
            {"pages": passing_pages + [page("p11"), page("p12"), page("p13")]},
            intake,
            registry,
            calibration={},
            layout_results=passing_budgets + [
                {"page_id": f"p{i}", "verdict": "FIT", "height_px": 400, "capacity_px": 600}
                for i in range(11, 14)
            ],
        )
        self.assertTrue(any(v.contract_id == "C15" for v in inflated.violations))

        sparse_budgets = [dict(row) for row in passing_budgets]
        for row in sparse_budgets[:3]:
            row.update(verdict="FIT", height_px=300)
        sparse = check_deck_spec_gates(plan, {"pages": passing_pages}, intake, registry, {}, layout_results=sparse_budgets)
        self.assertTrue(any("SPARSE 비율" in v.message for v in sparse.violations))

        low_median_budgets = [dict(row, height_px=350) for row in passing_budgets]
        low_median = check_deck_spec_gates(plan, {"pages": passing_pages}, intake, registry, {}, layout_results=low_median_budgets)
        self.assertTrue(any("밀도 중앙값" in v.message for v in low_median.violations))

        text_only_pages = [page(f"p{i}", visual=False) for i in range(1, 5)] + [page("p5")]
        text_budgets = [
            {"page_id": f"p{i}", "verdict": "FIT", "height_px": 400, "capacity_px": 600}
            for i in range(1, 6)
        ]
        text_only = check_deck_spec_gates(plan, {"pages": text_only_pages}, intake, registry, {}, layout_results=text_budgets)
        self.assertTrue(any("연속 텍스트-온리" in v.message for v in text_only.violations))
        self.assertTrue(any("시각 포함 비율" in v.message for v in text_only.violations))

        varied_pages = [page("p1", chars=1), page("p2", chars=1000), page("p3", chars=1)]
        varied_budgets = [
            {"page_id": f"p{i}", "verdict": "FIT", "height_px": 400, "capacity_px": 600}
            for i in range(1, 4)
        ]
        varied = check_deck_spec_gates(plan, {"pages": varied_pages}, intake, registry, {}, layout_results=varied_budgets)
        self.assertEqual(len(varied.warnings), 1)
        self.assertIn("정보량 편차", varied.warnings[0].message)

        exempt = check_deck_spec_gates(
            plan,
            {"pages": text_only_pages},
            {"intentional_text_deck": True},
            registry,
            {},
            layout_results=text_budgets,
        )
        self.assertFalse(any("시각 포함 비율" in v.message for v in exempt.violations))
        self.assertTrue(any("연속 텍스트-온리" in v.message for v in exempt.violations))

        not_exempt = check_deck_spec_gates(
            plan,
            {"pages": text_only_pages},
            {"intentional_text_deck": "false"},
            registry,
            {},
            layout_results=text_budgets,
        )
        self.assertTrue(any("시각 포함 비율" in v.message for v in not_exempt.violations))

        alias_only_pages = [
            {"page_id": f"p{i}", "layout": "stack", "content": [{"type": "metrics", "metric_ids": ["m1"]}]}
            for i in range(1, 5)
        ]
        alias_only = check_deck_spec_gates(
            plan, {"pages": alias_only_pages}, intake, registry, {}, layout_results=text_budgets[:4]
        )
        self.assertTrue(any("시각 포함 비율" in v.message for v in alias_only.violations))
        self.assertTrue(any("연속 텍스트-온리" in v.message for v in alias_only.violations))

        no_body = check_deck_spec_gates(
            plan,
            {"pages": [{"page_id": "p1", "layout": "cover", "content": []}]},
            intake,
            registry,
            {},
            layout_results=[{"page_id": "p1", "verdict": "FIT", "height_px": 400, "capacity_px": 600}],
        )
        self.assertTrue(any("본문 페이지가 없음" in v.message for v in no_body.violations))

    def test_deck_gates_reject_viz_metric_note_but_only_warn_for_viz_metric_and_missing_subtitle(self):
        def run(content):
            deck_spec = {"pages": [{
                "page_id": "p1", "layout": "stack",
                "content": [{"type": "headline", "text": "결론"}, *content],
            }]}
            return check_deck_spec_gates(
                {"pages": [{"page_id": "plan-1"}]}, deck_spec, {}, {}, {},
                layout_results=[{"page_id": "p1", "verdict": "FIT", "height_px": 400, "capacity_px": 600}],
            )

        b1 = run([{"type": "viz", "subtitle": "응답자 비율, %"}, {"type": "metric"}, {"type": "note"}])
        self.assertTrue(any("차트와 숫자 카드를 나누거나, note를 다음 장으로 옮기세요" in item.message for item in b1.violations))

        for evidence_type, metric_type in (("viz", "metric"), ("viz", "metric_grid"), ("text_table", "metric")):
            with self.subTest(evidence_type=evidence_type, metric_type=metric_type):
                evidence = {"type": evidence_type}
                if evidence_type == "viz":
                    evidence["subtitle"] = "응답자 비율, %"
                result = run([evidence, {"type": metric_type}])
                self.assertFalse(any("차트와 숫자 카드를" in item.message for item in result.violations))
                self.assertTrue(any("시각물 2종" in item.message for item in result.warnings))

        aliases = run([{"type": "viz", "subtitle": "응답자 비율, %"}, {"type": "metrics"}, {"type": "stat_grid"}])
        self.assertFalse(any("시각물 2종" in item.message for item in aliases.warnings))

        c = run([{"type": "viz"}])
        self.assertFalse(any("subtitle" in item.message for item in c.violations))
        self.assertTrue(any("subtitle" in item.message for item in c.warnings))

    def test_render_viz_three_level_title_and_preserves_legacy_html_bytes(self):
        legacy_html = render_deck_module.render_deck(
            VALID_DECK_SPEC_WITH_VIZ, VALID_CONTENT_REGISTRY, title="Legacy Viz"
        )
        self.assertEqual(
            hashlib.sha256(legacy_html.encode("utf-8")).hexdigest(),
            "fb1db570c87ba331ec93081f812bd997fcbb89b845bfa734d28498dd403bc854",
        )

        viz = dict(
            VALID_DECK_SPEC_WITH_VIZ["pages"][0]["content"][1],
            exhibit="Exhibit 3",
            title="AI를 쓰는 조직은 늘었지만 수익까지 간 곳은 적다.",
            subtitle="AI 도입 단계별 비율, % of respondents",
        )
        spec = {"pages": [dict(VALID_DECK_SPEC_WITH_VIZ["pages"][0], content=[viz])]}
        html = render_deck_module.render_deck(spec, VALID_CONTENT_REGISTRY, title="Three Level Viz")
        self.assertRegex(html, r'class="visual-title-eyebrow"[^>]*>Exhibit 3</div>')
        self.assertRegex(html, r'class="visual-title-conclusion"[^>]*>AI를 쓰는 조직은 늘었지만 수익까지 간 곳은 적다\.</div>')
        self.assertRegex(html, r'class="visual-title-subtitle"[^>]*>AI 도입 단계별 비율, % of respondents</div>')

    def test_generated_schema_declares_optional_viz_title_fields_as_strings(self):
        schema_module_path = pathlib.Path(__file__).with_name("generate_deck_spec_schema.py")
        module_spec = importlib.util.spec_from_file_location("generate_deck_spec_schema_for_test", schema_module_path)
        schema_module = importlib.util.module_from_spec(module_spec)
        module_spec.loader.exec_module(schema_module)
        schema = json.loads(schema_module.schema_bytes())
        block_schema = schema["properties"]["pages"]["items"]["properties"]["content"]["items"]
        self.assertEqual(block_schema["properties"]["exhibit"], {"type": "string"})
        self.assertEqual(block_schema["properties"]["title"], {"type": "string"})
        self.assertEqual(block_schema["properties"]["subtitle"], {"type": "string"})
        self.assertNotIn("exhibit", block_schema.get("required", []))
        self.assertNotIn("subtitle", block_schema.get("required", []))

    def test_gate4_capture_deck_exits_2_on_fit_overflow_and_render_main_keeps_debug_pdf(self):
        # 게이트 4 — capture_deck.sh는 FIT_OVERFLOW를 잡으면 PDF를 만든 채로 exit 2해야 하고,
        # render_deck.py main()은 그 exit 2를 받으면 (다른 실패와 달리) PDF를 지우지 않고 전파해야 한다.
        script = CAPTURE_DECK_SH.read_text(encoding="utf-8")
        self.assertIn("FIT_OVERFLOW_HIT=1", script)
        self.assertIn('if [ "${FIT_OVERFLOW_HIT:-0}" = "1" ] || [ "${INK_FAIL_HIT:-0}" = "1" ]; then', script)
        self.assertLess(script.index('mv "$PDF_TMP" "$OUT"'), script.index("FIT_OVERFLOW_HIT:-0"))

        capture_overflow = subprocess.CompletedProcess(
            args=["bash", str(CAPTURE_DECK_SH)],
            returncode=2,
            stdout="FIT_OVERFLOW: p01 — 본문이 세로 공간을 초과해 잘림.",
            stderr="",
        )
        with tempfile.TemporaryDirectory() as tmp:
            root = pathlib.Path(tmp)
            spec = root / "06_deck_spec.json"
            registry = root / "02_verified.json"
            output = root / "deck.html"
            debug_pdf = root / "deck.pdf"
            spec.write_text('{"pages": []}', encoding="utf-8")
            registry.write_text('{"sources": {}, "metrics": {}}', encoding="utf-8")
            debug_pdf.write_bytes(b"debug-artifact")
            argv = ["render_deck.py", str(spec), str(registry), "-o", str(output)]
            with (
                mock.patch.object(render_deck_module, "render_deck", return_value="<html></html>"),
                mock.patch.object(render_deck_module.subprocess, "run", return_value=capture_overflow),
                mock.patch.object(sys, "argv", argv),
            ):
                with self.assertRaises(SystemExit) as raised:
                    render_deck_module.main()

            self.assertTrue(debug_pdf.exists())
            self.assertEqual(debug_pdf.read_bytes(), b"debug-artifact")

        self.assertEqual(raised.exception.code, 2)

    def test_gate5_capture_deck_font_fetch_prefers_local_repo_assets_before_cdn(self):
        # 게이트 5 — fetch_font는 CDN보다 저장소 내장 assets/fonts를 먼저 확인해야 하고,
        # 그 자산이 실제로 존재 + 해시가 일치해야 오프라인 폴백이 성립한다.
        script = CAPTURE_DECK_SH.read_text(encoding="utf-8")
        self.assertIn('LOCAL_FONT_DIR="$SCRIPT_DIR/../assets/fonts"', script)
        fetch_font_body = script.split("fetch_font() {", 1)[1].split("\ncurl -fsSL", 1)[0]
        self.assertIn("LOCAL_FONT_DIR", fetch_font_body)

        assets_dir = pathlib.Path(__file__).resolve().parents[2] / "deck-harness" / "assets" / "fonts"
        expected_sha = {
            "Pretendard-Regular.woff2": "fad853f7f47c6c8b103171e7193fa095708cdcd70850a71d93aa5379e8a61d63",
            "Pretendard-Bold.woff2": "4609c3356e536fafe38f4add0daeceb3d8595d3057bce13c428c33ddbd43d362",
        }
        for name, sha in expected_sha.items():
            font_path = assets_dir / name
            self.assertTrue(font_path.exists(), f"missing vendored font: {name}")
            self.assertEqual(hashlib.sha256(font_path.read_bytes()).hexdigest(), sha)

    def test_gate6_render_layout_body_rejects_unknown_layout_but_keeps_supported_generic(self):
        # 게이트 6 — SUPPORTED_LAYOUTS 화이트리스트 밖은 즉시 raise. statement/timeline처럼
        # 마지막 generic branch로 떨어지는 "지원되는" layout은 그대로 렌더돼야 한다(C-10 경고).
        registry = {
            "sources": {"src_a": {"publisher": "Test", "url": "https://example.com"}},
            "metrics": {"metric_a": {"value": "1", "unit": "%", "source_ids": ["src_a"]}},
        }
        bad_spec = {
            "pages": [
                {
                    "page_id": "p01",
                    "short_title": "Bad",
                    "layout": "totally_unknown_layout",
                    "allowed_source_ids": [],
                    "allowed_metric_ids": [],
                    "content": [{"type": "headline", "text": "x"}],
                }
            ]
        }
        with self.assertRaises(ValueError) as ctx:
            render_deck_module.render_deck(bad_spec, registry, title="Bad Layout Fixture")
        self.assertIn("unsupported layout", str(ctx.exception))
        self.assertIn("totally_unknown_layout", str(ctx.exception))

        good_spec = {
            "pages": [
                dict(bad_spec["pages"][0], layout="statement", short_title="S", content=[{"type": "headline", "text": "ok-statement"}]),
                dict(bad_spec["pages"][0], page_id="p02", layout="timeline", short_title="T", content=[{"type": "headline", "text": "ok-timeline"}]),
            ]
        }
        html = render_deck_module.render_deck(good_spec, registry, title="Good Layout Fixture")
        self.assertIn("ok-statement", html)
        self.assertIn("ok-timeline", html)

    def test_schema_check_fails_after_source_of_truth_changes(self):
        scripts_dir = pathlib.Path(__file__).resolve().parent
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = pathlib.Path(temp_dir)
            shutil.copy2(scripts_dir / "contract_checks.py", temp_path / "contract_checks.py")
            shutil.copy2(scripts_dir / "generate_deck_spec_schema.py", temp_path / "generate_deck_spec_schema.py")
            schema_path = temp_path / "deck_spec.schema.json"

            generated = subprocess.run(
                [sys.executable, str(temp_path / "generate_deck_spec_schema.py"), "--output", str(schema_path)],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertEqual(generated.returncode, 0, generated.stderr)

            source_path = temp_path / "contract_checks.py"
            source = source_path.read_text(encoding="utf-8")
            source_path.write_text(source.replace('        "text_table",', '        "text_table",\n        "future_block",'), encoding="utf-8")

            checked = subprocess.run(
                [
                    sys.executable,
                    str(temp_path / "generate_deck_spec_schema.py"),
                    "--check",
                    "--output",
                    str(schema_path),
                ],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertNotEqual(checked.returncode, 0)
            self.assertIn("out of date", checked.stderr)


if __name__ == "__main__":
    unittest.main(verbosity=2)
