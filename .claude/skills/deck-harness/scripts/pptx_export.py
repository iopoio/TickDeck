#!/usr/bin/env python3
"""Export TickDeck HTML to editable-text PPTX.

Usage:
  python3 pptx_export.py <deck.html> [-o deck.pptx] [--verify]
"""

from __future__ import annotations

import argparse
import html
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import unicodedata
from functools import lru_cache
from pathlib import Path


SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000
FONT_NAME = "맑은 고딕"
VERIFY_PHRASE = "흔들렸으나 무너지지 않았다"
PX_TO_PT = 0.75
FIT_MIN_RATIO = 0.70
CHROME_TIMEOUT_SECONDS = 180
PIP_INSTALL_TIMEOUT_SECONDS = 300
_FONT_METRIC_FALLBACK_WARNED = False


LAYOUT_DUMP_SCRIPT = r"""
<script>
(function(){
  function hasDirectText(el) {
    return Array.prototype.some.call(el.childNodes, function(n) {
      return n.nodeType === 3 && n.textContent.trim();
    });
  }
  function px(value, fallback) {
    var n = parseFloat(value);
    return Number.isFinite(n) ? n : fallback;
  }
  var slides = [];
  document.querySelectorAll('.slide').forEach(function(slide, idx) {
    var slideRect = slide.getBoundingClientRect();
    var pageId = slide.dataset.pageId || slide.id || ('p' + String(idx + 1).padStart(2, '0'));
    var boxes = [];
    slide.querySelectorAll('*').forEach(function(el) {
      if (el.closest('svg')) return;
      if (el.closest('[aria-hidden="true"]')) return;
      if (el.offsetParent === null) return;
      if (!hasDirectText(el)) return;
      if (el.closest('[data-pptx-picked]')) return;
      var text = (el.innerText || el.textContent || '').trim();
      if (!text) return;
      // Range rect = 실제 텍스트 내용만 — ::before 대시 장식 등 의사요소가 el rect에 포함돼
      // PPT 텍스트박스가 장식 위에 겹치던 결함(7/5 실측 p08 bullets). 실패 시 el rect 폴백.
      var rect;
      try {
        var range = document.createRange();
        range.selectNodeContents(el);
        rect = range.getBoundingClientRect();
      } catch (e) { rect = el.getBoundingClientRect(); }
      if (!rect || rect.width <= 0 || rect.height <= 0) rect = el.getBoundingClientRect();
      if (rect.width <= 0 || rect.height <= 0) return;
      // 줄 수 = Range 줄별 rect의 top 클러스터 수(정확) — el rect는 패딩 포함이라 오판(콜아웃 2줄 오판·7/5),
      // rect.height/lineHeight 반올림도 근사였다. 이후 글리프 top에서 반행간을 빼 줄박스 격자에 y를 정렬:
      // PPT 고정 줄간격(spcPts)의 줄격자가 브라우저 줄박스와 일치해야 아랫박스 침범이 없다(p04 실측).
      var lineTops = [];
      try {
        Array.prototype.forEach.call(range.getClientRects(), function(r){
          if (r.width <= 0 || r.height <= 0) return;
          for (var k = 0; k < lineTops.length; k++) { if (Math.abs(lineTops[k] - r.top) < 3) return; }
          lineTops.push(r.top);
        });
      } catch (e) {}
      var cs0 = getComputedStyle(el);
      var lh0 = px(cs0.lineHeight, px(cs0.fontSize, 16) * 1.2);
      var lcExact = Math.max(1, lineTops.length || Math.round(rect.height / lh0));
      var glyphH = rect.height - (lcExact - 1) * lh0;
      var halfLeading = Math.max(0, (lh0 - glyphH) / 2);
      rect = {left: rect.left, width: rect.width, top: rect.top - halfLeading, height: lcExact * lh0};
      el.setAttribute('data-pptx-picked', '1');
      var cs = getComputedStyle(el);
      var fontSize = px(cs.fontSize, 16);
      var lineHeight = px(cs.lineHeight, fontSize * 1.2);
      var letterSpacing = cs.letterSpacing === 'normal' ? 0 : px(cs.letterSpacing, 0);
      var metricEl = el.closest('[data-metric-id]');
      var srcEl = el.closest('[data-src-id]');
      // 브라우저 기준 줄 수 — PPT에서 폰트 대체(맑은고딕 등)로 재줄바꿈되며 겹치던 결함(7/5 LibreOffice 실측).
      // 1줄 요소는 word_wrap을 꺼 줄바꿈을 봉인한다(파이썬 쪽 분기).
      var lineCount = lcExact;
      boxes.push({
        line_count: lineCount,
        page_id: pageId,
        text: text,
        x: rect.left - slideRect.left,
        y: rect.top - slideRect.top,
        w: rect.width,
        h: rect.height,
        slide_w: slideRect.width,
        slide_h: slideRect.height,
        font_size_px: fontSize,
        font_weight: cs.fontWeight,
        color: cs.color,
        text_align: cs.textAlign,
        line_height_px: lineHeight,
        letter_spacing: letterSpacing,
        metric_id: metricEl ? metricEl.dataset.metricId : null,
        src_id: srcEl ? srcEl.dataset.srcId : null
      });
    });
    slides.push({page_id: pageId, slide_w: slideRect.width, slide_h: slideRect.height, boxes: boxes});
  });
  var out = document.createElement('script');
  out.type = 'application/json';
  out.id = '__layout_dump__';
  out.textContent = JSON.stringify({slides: slides});
  document.body.appendChild(out);
})();
</script>
"""


HIDE_PICKED_TEXT_SCRIPT = r"""
<script>
(function(){
  document.querySelectorAll('[data-pptx-picked="1"]').forEach(function(el) {
    var targets = [el].concat(Array.prototype.slice.call(el.querySelectorAll('*')));
    targets.forEach(function(target) {
      if (target.closest('svg')) return;
      target.style.setProperty('color', 'transparent', 'important');
      target.style.setProperty('-webkit-text-fill-color', 'transparent', 'important');
      target.style.setProperty('text-shadow', 'none', 'important');
    });
  });
})();
</script>
"""


def script_root() -> Path:
    return Path(__file__).resolve().parent


def harness_root() -> Path:
    return script_root().parent


def venv_dir() -> Path:
    return harness_root() / "pptx_venv"


def venv_python() -> Path:
    return venv_dir() / "bin" / "python"


def venv_pip() -> Path:
    return venv_dir() / "bin" / "pip"


def run_checked(cmd: list[str], error_label: str, timeout: int | None = None) -> subprocess.CompletedProcess[str]:
    try:
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
    except subprocess.TimeoutExpired as exc:
        limit = timeout if timeout is not None else exc.timeout
        raise RuntimeError(f"{error_label} timed out after {limit}s") from exc
    if proc.returncode != 0:
        detail = (proc.stderr or proc.stdout or "").strip()
        raise RuntimeError(f"{error_label} failed: {detail}")
    return proc


def packages_available(py: Path) -> bool:
    proc = subprocess.run(
        [str(py), "-c", "import pptx, fitz"],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    return proc.returncode == 0


def python_version_tag(py: Path) -> str:
    proc = run_checked(
        [str(py), "-c", "import sys; print(f'python{sys.version_info.major}.{sys.version_info.minor}')"],
        "python version check",
    )
    return proc.stdout.strip()


def site_packages_dir(py: Path) -> Path:
    proc = run_checked(
        [str(py), "-c", "import sysconfig; print(sysconfig.get_paths()['purelib'])"],
        "site-packages lookup",
    )
    return Path(proc.stdout.strip())


def local_site_package_candidates(env_value: str | None = None) -> list[Path]:
    candidates = []
    raw_env = env_value if env_value is not None else os.environ.get("PPTX_EXPORT_EXTRA_SITE_PACKAGES", "")
    for item in raw_env.split(os.pathsep):
        if item:
            candidates.append(Path(item).expanduser())

    automation = Path("/Users/hwa/Projects/Automation")
    candidates.extend(
        [
            automation / "Think/.venv/lib/python3.14/site-packages",
            harness_root() / "scripts/.venv/lib/python3.12/site-packages",
        ]
    )
    return candidates


def link_existing_site_packages(py: Path) -> Path | None:
    target_site = site_packages_dir(py)
    version_tag = python_version_tag(py)
    for candidate in local_site_package_candidates():
        if not candidate.exists():
            continue
        if version_tag not in str(candidate):
            continue
        if not (candidate / "pptx").exists() or not (candidate / "fitz").exists():
            continue
        target_site.mkdir(parents=True, exist_ok=True)
        pth = target_site / "tickdeck_pptx_extra_site_packages.pth"
        pth.write_text(str(candidate) + "\n", encoding="utf-8")
        if packages_available(py):
            return candidate
        try:
            pth.unlink()
        except FileNotFoundError:
            pass
    return None


def bootstrap_venv() -> None:
    root = venv_dir()
    py = venv_python()
    if not py.exists():
        try:
            run_checked([sys.executable, "-m", "venv", str(root)], "venv bootstrap")
        except Exception as exc:
            raise SystemExit(f"BOOTSTRAP_ERROR: could not create {root}: {exc}") from exc
    if not packages_available(py):
        try:
            run_checked(
                [str(venv_pip()), "install", "python-pptx", "pymupdf"],
                "pip install",
                timeout=PIP_INSTALL_TIMEOUT_SECONDS,
            )
        except Exception as exc:
            linked = link_existing_site_packages(py)
            if linked:
                return
            raise SystemExit(
                "BOOTSTRAP_ERROR: could not install python-pptx pymupdf into "
                f"{root}: {exc}"
            ) from exc


def ensure_runtime() -> None:
    target = venv_python()
    if not target.exists() or not packages_available(target):
        bootstrap_venv()
    # venv python이 시스템 python 심링크면 executable.resolve() 비교가 같은 바이너리로 접혀
    # re-exec을 건너뛰고 시스템 python으로 돌던 결함(7/5 실측) — venv 여부는 sys.prefix로 판정.
    if Path(sys.prefix).resolve() != venv_dir().resolve():
        os.execv(str(target), [str(target), str(Path(__file__).resolve()), *sys.argv[1:]])


def find_chrome() -> str:
    candidates = [
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        shutil.which("chromium"),
        shutil.which("google-chrome"),
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists() and os.access(candidate, os.X_OK):
            return str(candidate)
    raise SystemExit("NO_CHROME: Chrome/Chromium/Edge not found for layout dump and PDF capture")


def insert_before_body(source: str, injection: str) -> str:
    lower = source.lower()
    body = lower.rfind("</body>")
    if body != -1:
        return source[:body] + injection + source[body:]
    html_end = lower.rfind("</html>")
    if html_end != -1:
        return source[:html_end] + injection + source[html_end:]
    return source + injection


def sibling_temp_html(deck_html: Path, suffix: str, content: str) -> Path:
    handle = tempfile.NamedTemporaryFile(
        "w",
        encoding="utf-8",
        delete=False,
        dir=str(deck_html.parent),
        prefix=f"{deck_html.stem}.pptx_",
        suffix=suffix,
    )
    with handle:
        handle.write(content)
    return Path(handle.name)


def run_chrome(chrome: str, args: list[str], label: str) -> subprocess.CompletedProcess[str]:
    try:
        proc = subprocess.run(
            [chrome, *args],
            capture_output=True,
            text=True,
            timeout=CHROME_TIMEOUT_SECONDS,
        )
    except subprocess.TimeoutExpired as exc:
        limit = int(exc.timeout or CHROME_TIMEOUT_SECONDS)
        raise SystemExit(f"{label}_TIMEOUT: Chrome timed out after {limit}s") from exc
    if proc.returncode != 0:
        detail = (proc.stderr or proc.stdout or "").strip()
        raise SystemExit(f"{label}_ERROR: Chrome exited {proc.returncode}: {detail}")
    return proc


def extract_layout_json(dom: str) -> dict:
    match = re.search(
        r"<script(?=[^>]*\bid=[\"']__layout_dump__[\"'])[^>]*>(.*?)</script>",
        dom,
        flags=re.IGNORECASE | re.DOTALL,
    )
    if not match:
        raise SystemExit("LAYOUT_DUMP_ERROR: __layout_dump__ element not found in Chrome dump-dom output")
    raw = html.unescape(match.group(1).strip())
    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        raise SystemExit(f"LAYOUT_DUMP_ERROR: invalid JSON in __layout_dump__: {exc}") from exc


def dump_layout(deck_html: Path, chrome: str) -> tuple[dict, str, Path]:
    source = deck_html.read_text(encoding="utf-8")
    injected = sibling_temp_html(deck_html, ".layout.html", insert_before_body(source, LAYOUT_DUMP_SCRIPT))
    try:
        proc = run_chrome(
            chrome,
            ["--headless=new", "--disable-gpu", "--dump-dom", injected.resolve().as_uri()],
            "LAYOUT_DUMP",
        )
        return extract_layout_json(proc.stdout), proc.stdout, injected
    except BaseException:
        cleanup_paths([injected])
        raise


def print_hidden_pdf(deck_html: Path, dom: str, chrome: str, out_pdf: Path) -> Path:
    hidden_html = sibling_temp_html(deck_html, ".hidden.html", insert_before_body(dom, HIDE_PICKED_TEXT_SCRIPT))
    try:
        run_chrome(
            chrome,
            [
                "--headless=new",
                "--disable-gpu",
                "--no-pdf-header-footer",
                f"--print-to-pdf={out_pdf}",
                hidden_html.resolve().as_uri(),
            ],
            "PRINT_PDF",
        )
        if not out_pdf.exists():
            raise SystemExit(f"PRINT_PDF_ERROR: Chrome did not create {out_pdf}")
        return hidden_html
    except BaseException:
        cleanup_paths([hidden_html])
        raise


def render_pdf_pages(pdf: Path, layout: dict, png_dir: Path, keep_p03_path: Path | None):
    import fitz

    png_dir.mkdir(parents=True, exist_ok=True)
    slides = layout.get("slides", [])
    pages = []
    kept_p03 = None
    with fitz.open(str(pdf)) as doc:
        if len(doc) != len(slides):
            raise SystemExit(f"PDF_PAGE_ERROR: pdf pages={len(doc)} layout slides={len(slides)}")
        for idx, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            png = png_dir / f"bg_{idx + 1:03d}.png"
            pix.save(str(png))
            pages.append(png)
            if keep_p03_path and slides[idx].get("page_id") == "p03":
                shutil.copyfile(png, keep_p03_path)
                kept_p03 = keep_p03_path
    return pages, kept_p03


def parse_css_color(value: str) -> tuple[int, int, int]:
    value = (value or "").strip()
    if not value or value == "transparent":
        return (0, 0, 0)
    if value.startswith("#"):
        raw = value[1:]
        if len(raw) == 3:
            return tuple(int(ch * 2, 16) for ch in raw)  # type: ignore[return-value]
        if len(raw) >= 6:
            return (int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    nums = re.findall(r"[-+]?\d*\.?\d+%?", value)
    if len(nums) < 3:
        return (0, 0, 0)

    def component(token: str, srgb_unit: bool) -> int:
        if token.endswith("%"):
            n = float(token[:-1]) / 100 * 255
        else:
            n = float(token)
            if srgb_unit:
                n *= 255
        return max(0, min(255, int(round(n))))

    srgb_unit = value.startswith("color(") and all(
        not nums[i].endswith("%") and abs(float(nums[i])) <= 1 for i in range(3)
    )
    return tuple(component(nums[i], srgb_unit) for i in range(3))  # type: ignore[return-value]


def font_weight_number(value: object) -> int:
    text = str(value or "").strip().lower()
    if text == "bold":
        return 700
    if text in {"normal", "regular"}:
        return 400
    match = re.search(r"\d+", text)
    return int(match.group(0)) if match else 400


def px_to_emu(px: float, slide_w: float) -> int:
    scale = SLIDE_W_EMU / (slide_w or 1280)
    return int(round(px * scale))


def emu_to_px(emu: int, slide_w: float) -> float:
    scale = SLIDE_W_EMU / (slide_w or 1280)
    return emu / scale


def text_box_geometry_emu(box: dict) -> dict[str, int]:
    slide_w = float(box.get("slide_w") or 1280)
    width_px = float(box.get("w") or 0) * 1.02 + 2
    return {
        "x": px_to_emu(float(box.get("x") or 0), slide_w),
        "y": px_to_emu(float(box.get("y") or 0), slide_w),
        "w": max(1, px_to_emu(width_px, slide_w)),
        "h": max(1, px_to_emu(float(box.get("h") or 0), slide_w)),
    }


def alt_text_for_box(box: dict) -> str:
    parts = []
    if box.get("metric_id"):
        parts.append(f"metric_id={box['metric_id']}")
    if box.get("src_id"):
        parts.append(f"src_id={box['src_id']}")
    return "; ".join(parts)


def set_shape_descr(shape, descr: str) -> None:
    if descr:
        shape._element.nvSpPr.cNvPr.set("descr", descr)


def paragraph_alignment(value: str):
    from pptx.enum.text import PP_ALIGN

    normalized = (value or "left").strip().lower()
    return {
        "left": PP_ALIGN.LEFT,
        "start": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "end": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(normalized, PP_ALIGN.LEFT)


def _first_existing_path(paths: list[Path]) -> Path | None:
    for path in paths:
        expanded = path.expanduser()
        if expanded.exists():
            return expanded
    return None


@lru_cache(maxsize=1)
def resolve_malgun_fontfiles() -> tuple[Path | None, Path | None]:
    roots = [
        Path("/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts"),
        Path("~/Library/Fonts"),
        Path("/Library/Fonts"),
    ]
    regular = _first_existing_path([root / "malgun.ttf" for root in roots])
    bold = _first_existing_path([root / "malgunbd.ttf" for root in roots])
    return regular, bold


def warn_font_metric_fallback() -> None:
    global _FONT_METRIC_FALLBACK_WARNED
    if not _FONT_METRIC_FALLBACK_WARNED:
        print("FONT_METRIC_FALLBACK", file=sys.stderr)
        _FONT_METRIC_FALLBACK_WARNED = True


@lru_cache(maxsize=2)
def _metric_font(bold: bool):
    regular_path, bold_path = resolve_malgun_fontfiles()
    fontfile = bold_path if bold and bold_path else regular_path
    if not fontfile:
        return None
    try:
        import fitz

        return fitz.Font(fontfile=str(fontfile))
    except Exception:
        return None


def is_cjk_or_fullwidth(ch: str) -> bool:
    return unicodedata.east_asian_width(ch) in {"F", "W"}


def fallback_text_width_pt(text: str, font_pt: float, bold: bool) -> float:
    width_em = 0.0
    for ch in text:
        if ch.isspace():
            width_em += 0.33
        elif is_cjk_or_fullwidth(ch):
            width_em += 1.0
        else:
            width_em += 0.62
    if bold:
        width_em *= 1.04
    return width_em * font_pt


def measure_text_width_pt(text: str, font_pt: float, bold: bool, letter_spacing_pt: float = 0.0) -> float:
    font = _metric_font(bold)
    if font:
        try:
            width = font.text_length(text, fontsize=font_pt)
        except Exception:
            font = None
    if not font:
        warn_font_metric_fallback()
        width = fallback_text_width_pt(text, font_pt, bold)
    if text:
        width += letter_spacing_pt * max(0, len(text) - 1)
    return width


def letter_spacing_pt_from_box(box: dict) -> float:
    return float(box.get("letter_spacing") or 0) * PX_TO_PT


def floor_one_decimal(value: float) -> float:
    return math.floor(value * 10) / 10


def warn_fit_shrink_cap(box: dict) -> None:
    page = str(box.get("page_id") or "?")
    snippet = str(box.get("text") or "").replace("\n", " ")[:30]
    print(f"FIT_SHRINK_CAP: {page}/{snippet}", file=sys.stderr)


def fit_single_line_font_pt(
    lines: list[str],
    font_pt: float,
    box_w_pt: float,
    bold: bool,
    letter_spacing_pt: float,
    box: dict,
) -> float:
    text_w_pt = max((measure_text_width_pt(line, font_pt, bold, letter_spacing_pt) for line in lines), default=0)
    if not text_w_pt or text_w_pt <= box_w_pt * 1.01:
        return font_pt

    target_pt = floor_one_decimal(font_pt * box_w_pt / text_w_pt)
    min_pt = font_pt * FIT_MIN_RATIO
    if target_pt < min_pt:
        warn_fit_shrink_cap(box)
        return min_pt
    return target_pt


def _wrap_tokens(text: str) -> list[str]:
    tokens: list[str] = []
    word: list[str] = []

    def flush_word() -> None:
        if word:
            tokens.append("".join(word))
            word.clear()

    for ch in text:
        if ch.isspace():
            flush_word()
            tokens.append(ch)
        elif is_cjk_or_fullwidth(ch):
            flush_word()
            tokens.append(ch)
        else:
            word.append(ch)
    flush_word()
    return tokens


def wrapped_lines_needed(text: str, box_w_pt: float, font_pt: float, bold: bool, letter_spacing_pt: float) -> int:
    if box_w_pt <= 0:
        return 1

    def fits(value: str) -> bool:
        return measure_text_width_pt(value, font_pt, bold, letter_spacing_pt) <= box_w_pt

    def place_on_empty_line(token: str) -> tuple[str, int]:
        current = ""
        extra_lines = 0
        for ch in token:
            candidate = current + ch
            if not current or fits(candidate):
                current = candidate
            else:
                extra_lines += 1
                current = ch
        return current, extra_lines

    lines = 1
    current = ""
    for token in _wrap_tokens(text):
        if token.isspace():
            if current and fits(current + token):
                current += token
            continue
        if not current:
            current, extra = place_on_empty_line(token)
            lines += extra
        elif fits(current + token):
            current += token
        else:
            lines += 1
            current, extra = place_on_empty_line(token)
            lines += extra
    return max(1, lines)


def wrapped_text_fits(
    text: str,
    line_count: int,
    box_w_pt: float,
    box_h_pt: float,
    font_pt: float,
    line_spacing: float,
    bold: bool,
    letter_spacing_pt: float,
) -> bool:
    needed = wrapped_lines_needed(text, box_w_pt, font_pt, bold, letter_spacing_pt)
    if needed > line_count:
        return False
    if box_h_pt > 0 and needed * line_spacing * font_pt > box_h_pt * 1.05:
        return False
    return True


def fit_wrapped_font_pt(
    text: str,
    font_pt: float,
    line_count: int,
    box_w_pt: float,
    box_h_pt: float,
    line_spacing: float,
    bold: bool,
    letter_spacing_pt: float,
    box: dict,
) -> float:
    if wrapped_text_fits(text, line_count, box_w_pt, box_h_pt, font_pt, line_spacing, bold, letter_spacing_pt):
        return font_pt

    min_pt = font_pt * FIT_MIN_RATIO
    current = font_pt
    while current - 0.5 >= min_pt:
        current = round(current - 0.5, 3)
        if wrapped_text_fits(text, line_count, box_w_pt, box_h_pt, current, line_spacing, bold, letter_spacing_pt):
            return current
    if wrapped_text_fits(text, line_count, box_w_pt, box_h_pt, min_pt, line_spacing, bold, letter_spacing_pt):
        return min_pt

    warn_fit_shrink_cap(box)
    return min_pt


def fitted_font_pt(box: dict, font_pt: float, line_spacing: float, bold: bool) -> float:
    text = str(box.get("text") or "")
    lines = text.splitlines() or [""]
    box_w_pt = float(box.get("w") or 0) * PX_TO_PT
    box_h_pt = float(box.get("h") or 0) * PX_TO_PT
    line_count = max(1, int(box.get("line_count") or 1))
    word_wrap = line_count > 1
    letter_spacing_pt = letter_spacing_pt_from_box(box)

    if box_w_pt <= 0:
        return font_pt
    if len(lines) > 1 or not word_wrap:
        return fit_single_line_font_pt(lines, font_pt, box_w_pt, bold, letter_spacing_pt, box)
    return fit_wrapped_font_pt(
        text,
        font_pt,
        line_count,
        box_w_pt,
        box_h_pt,
        line_spacing,
        bold,
        letter_spacing_pt,
        box,
    )


def set_run_character_spacing(run, letter_spacing_px: float) -> None:
    spc = int(letter_spacing_px * PX_TO_PT * 100)
    if spc:
        run._r.get_or_add_rPr().set("spc", str(spc))


def apply_text_to_shape(shape, box: dict) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.util import Pt

    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    # 브라우저에서 1줄이던 요소는 word_wrap 봉인 — 폰트 대체(맑은고딕 등) 폭 차이로
    # 재줄바꿈돼 아래 박스와 겹치던 결함(7/5 LibreOffice 실측 p08 h1·p09 출처행).
    frame.word_wrap = int(box.get("line_count") or 1) > 1
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_ANCHOR.TOP

    lines = str(box.get("text") or "").splitlines() or [""]
    font_size_px = float(box.get("font_size_px") or 16)
    line_height_px = float(box.get("line_height_px") or font_size_px * 1.2)
    line_spacing = line_height_px / font_size_px if font_size_px else 1.2
    rgb = parse_css_color(str(box.get("color") or "rgb(0,0,0)"))
    alignment = paragraph_alignment(str(box.get("text_align") or "left"))
    bold = font_weight_number(box.get("font_weight")) >= 600
    font_pt = max(1, int(round(font_size_px * 0.75)))
    fit_pt = fitted_font_pt(box, font_pt, line_spacing, bold)
    letter_spacing_px = float(box.get("letter_spacing") or 0)

    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = alignment
        # 배수(multiple) 줄간격은 PPT가 폰트 고유 행높이(맑은고딕>CSS 1em)에 곱해 브라우저보다
        # 줄이 굵어짐(7/5 실측 p04 라벨 3째줄 침범) — 포인트 고정값으로 브라우저 기하와 일치시킨다.
        paragraph.line_spacing = Pt(line_height_px * PX_TO_PT)
        run = paragraph.add_run()
        run.text = line
        font = run.font
        font.name = FONT_NAME
        font.size = Pt(fit_pt)
        font.bold = bold
        font.color.rgb = RGBColor(*rgb)
        set_run_character_spacing(run, letter_spacing_px)


def build_pptx(layout: dict, background_pngs: list[Path], out: Path) -> int:
    from pptx import Presentation
    from pptx.util import Emu

    slides = layout.get("slides", [])
    if len(slides) != len(background_pngs):
        raise SystemExit(f"PPTX_BUILD_ERROR: slides={len(slides)} backgrounds={len(background_pngs)}")

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank = prs.slide_layouts[6]

    for slide_data, bg in zip(slides, background_pngs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(bg), 0, 0, width=Emu(SLIDE_W_EMU), height=Emu(SLIDE_H_EMU))
        for box in slide_data.get("boxes", []):
            geom = text_box_geometry_emu(box)
            shape = slide.shapes.add_textbox(
                Emu(geom["x"]),
                Emu(geom["y"]),
                Emu(geom["w"]),
                Emu(geom["h"]),
            )
            apply_text_to_shape(shape, box)
            set_shape_descr(shape, alt_text_for_box(box))

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return len(slides)


def pptx_shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def shape_counts(slide) -> tuple[int, int]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pictures = sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    text_boxes = sum(1 for shape in slide.shapes if pptx_shape_text(shape))
    return pictures, text_boxes


def find_layout_box(layout: dict, page_id: str, phrase: str) -> dict | None:
    for slide in layout.get("slides", []):
        if slide.get("page_id") != page_id:
            continue
        for box in slide.get("boxes", []):
            if phrase in str(box.get("text") or ""):
                return box
    return None


def verify_pptx(pptx_path: Path, layout: dict) -> dict:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides = layout.get("slides", [])
    if len(prs.slides) != len(slides):
        raise SystemExit(f"VERIFY_ERROR: pptx slides={len(prs.slides)} layout slides={len(slides)}")

    total_text_boxes = 0
    for idx, slide in enumerate(prs.slides):
        pictures, text_boxes = shape_counts(slide)
        if pictures != 1:
            raise SystemExit(f"VERIFY_ERROR: slide {idx + 1} background pictures={pictures}, expected 1")
        if text_boxes < 1:
            raise SystemExit(f"VERIFY_ERROR: slide {idx + 1} text boxes={text_boxes}, expected >=1")
        total_text_boxes += text_boxes

    phrase_box = find_layout_box(layout, "p03", VERIFY_PHRASE)
    phrase_checked = False
    max_delta_px = None
    if phrase_box:
        p03_index = next(
            (idx for idx, slide in enumerate(slides) if slide.get("page_id") == "p03"),
            None,
        )
        if p03_index is None:
            raise SystemExit("VERIFY_ERROR: p03 layout slide not found")
        match_shape = None
        for shape in prs.slides[p03_index].shapes:
            if VERIFY_PHRASE in pptx_shape_text(shape):
                match_shape = shape
                break
        if match_shape is None:
            raise SystemExit(f"VERIFY_ERROR: p03 phrase not found in PPTX text boxes: {VERIFY_PHRASE}")
        expected = text_box_geometry_emu(phrase_box)
        actual = {"x": match_shape.left, "y": match_shape.top, "w": match_shape.width, "h": match_shape.height}
        slide_w = float(phrase_box.get("slide_w") or 1280)
        deltas = {key: abs(emu_to_px(int(actual[key]) - expected[key], slide_w)) for key in expected}
        max_delta_px = max(deltas.values())
        if max_delta_px > 1:
            raise SystemExit(f"VERIFY_ERROR: p03 phrase geometry delta >1px: {deltas}")
        phrase_checked = True

    return {
        "slides": len(prs.slides),
        "text_boxes": total_text_boxes,
        "phrase_checked": phrase_checked,
        "max_delta_px": max_delta_px,
    }


def export_deck(deck_html: Path, out: Path, verify: bool) -> tuple[int, Path | None, dict | None]:
    chrome = find_chrome()
    temp_files = []
    try:
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            layout, dumped_dom, injected_html = dump_layout(deck_html, chrome)
            temp_files.append(injected_html)
            hidden_pdf = tmp / "text_hidden.pdf"
            hidden_html = print_hidden_pdf(deck_html, dumped_dom, chrome, hidden_pdf)
            temp_files.append(hidden_html)
            p03_png = out.with_name(f"{out.stem}.p03_hidden.png")
            background_pngs, kept_p03 = render_pdf_pages(hidden_pdf, layout, tmp / "png", p03_png)
            count = build_pptx(layout, background_pngs, out)
            verification = verify_pptx(out, layout) if verify else None
        return count, kept_p03, verification
    finally:
        cleanup_paths(temp_files)


def cleanup_paths(paths: list[Path]) -> None:
    for path in paths:
        try:
            path.unlink()
        except FileNotFoundError:
            pass


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Export TickDeck deck.html to editable PPTX")
    parser.add_argument("deck_html", type=Path)
    parser.add_argument("-o", "--output", type=Path, default=None)
    parser.add_argument("--verify", action="store_true", help="read back the PPTX and verify structure/geometry")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    deck_html = args.deck_html.expanduser().resolve()
    if not deck_html.exists():
        raise SystemExit(f"NO_DECK: {deck_html}")
    out = (args.output.expanduser().resolve() if args.output else deck_html.with_name("deck.pptx"))

    ensure_runtime()
    count, kept_p03, verification = export_deck(deck_html, out, args.verify)
    size_mb = out.stat().st_size / 1024 / 1024
    print(f"PPTX_OK: {count} slides -> {out} ({size_mb:.1f}MB)")
    if kept_p03:
        print(f"HIDDEN_P03_PNG: {kept_p03}")
    else:
        print("HIDDEN_P03_PNG_SKIP: p03 not found")
    if verification:
        phrase = "yes" if verification["phrase_checked"] else "skip"
        delta = verification["max_delta_px"]
        delta_text = "n/a" if delta is None else f"{delta:.3f}px"
        print(
            "VERIFY_OK: "
            f"slides={verification['slides']} "
            f"text_boxes={verification['text_boxes']} "
            f"p03_phrase={phrase} "
            f"max_delta={delta_text}"
        )


if __name__ == "__main__":
    main()
