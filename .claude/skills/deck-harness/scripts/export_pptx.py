#!/usr/bin/env python3
"""덱 PDF → 완성형 PPTX (G4 결재 2026-07-03: "완성형 리포트를 PPTX로").

1차 방식 = 페이지 이미지 슬라이드: 렌더된 PDF 각 페이지를 16:9 슬라이드에
전면 이미지로 박는다. 폰트 portability 문제 원천 차단(수신 기기에 서체 불필요)·
시각 정합 100%(HTML 렌더 그대로). 편집 불가 — 부족하면 2차(네이티브 텍스트박스,
deck_spec 직접 빌드)로 확장(UPGRADE_PLAN.md 3단계).

사용: export_pptx.py <deck.pdf> [-o out.pptx] [--dpi 200]
의존: pdftoppm(poppler), python-pptx(scripts/.venv — uv venv로 생성)
"""

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ModuleNotFoundError:
    # 시스템 python으로 불렸으면 동봉 venv로 재실행 (uvicorn식 셔뱅 의존 X)
    venv_py = Path(__file__).parent / ".venv" / "bin" / "python"
    if venv_py.exists() and Path(sys.executable).resolve() != venv_py.resolve():
        raise SystemExit(subprocess.run([str(venv_py), *sys.argv]).returncode)
    raise SystemExit("python-pptx 없음 — scripts/.venv 생성 필요: uv venv && uv pip install python-pptx")

SLIDE_W = Emu(12192000)  # 13.333in — 16:9 표준 와이드
SLIDE_H = Emu(6858000)   # 7.5in


def export(pdf: Path, out: Path, dpi: int) -> int:
    with tempfile.TemporaryDirectory() as td:
        # JPEG: 슬라이드 전면 사진형 페이지라 PNG 대비 1/5 용량·화질 차이 비가시.
        r = subprocess.run(
            ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf), f"{td}/pg"],
            capture_output=True, text=True,
        )
        if r.returncode != 0:
            raise SystemExit(f"pdftoppm 실패: {r.stderr.strip()}")
        pages = sorted(Path(td).glob("pg-*.jpg"))
        if not pages:
            raise SystemExit(f"PDF에서 페이지를 못 뽑음: {pdf}")

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        blank = prs.slide_layouts[6]
        for img in pages:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(img), 0, 0, width=SLIDE_W, height=SLIDE_H)
        prs.save(str(out))
        return len(pages)


def main() -> None:
    ap = argparse.ArgumentParser(description="덱 PDF를 페이지 이미지 PPTX로 변환")
    ap.add_argument("pdf", type=Path)
    ap.add_argument("-o", "--output", type=Path, default=None)
    ap.add_argument("--dpi", type=int, default=200, help="페이지 이미지 해상도 (기본 200 ≈ 2666px 폭)")
    args = ap.parse_args()
    out = args.output or args.pdf.with_suffix(".pptx")
    n = export(args.pdf, out, args.dpi)
    size_mb = out.stat().st_size / 1024 / 1024
    print(f"PPTX_OK: {n} slides -> {out} ({size_mb:.1f}MB)")


if __name__ == "__main__":
    main()
