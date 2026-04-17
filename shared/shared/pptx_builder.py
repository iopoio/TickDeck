"""pptx_builder.py — python-pptx 서버사이드 PPTX 생성
타입별 전용 레이아웃: cover / problem / solution / how_it_works /
                     key_metrics / proof / why_us / cta / (fallback)
"""
import colorsys
import datetime
import io
import logging

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ── 슬라이드 사이즈 (16:9 와이드) ──────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# ── 폰트 ────────────────────────────────────────────────────────────
FONT_HEAD = "Pretendard"
FONT_BODY = "Pretendard"
FONT_EA = "Noto Sans KR"


# ══════════════════════════════════════════════════════════════════════
# 공통 유틸
# ══════════════════════════════════════════════════════════════════════

def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    hex_str = hex_str.lstrip('#')
    if len(hex_str) == 3:
        hex_str = ''.join([c * 2 for c in hex_str])
    return (int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor(*_hex_to_rgb(hex_str))


def _tinted_dark(primary_hex: str) -> RGBColor:
    """primary hue를 섞은 거의-검정 (다크 배경용)"""
    r, g, b = _hex_to_rgb(primary_hex)
    h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
    nr, ng, nb = colorsys.hls_to_rgb(h, 0.08, 0.25)
    return RGBColor(int(nr * 255), int(ng * 255), int(nb * 255))


def _tinted_light(primary_hex: str) -> RGBColor:
    """primary hue를 섞은 아주 밝은 색 (섹션 배경용)"""
    r, g, b = _hex_to_rgb(primary_hex)
    h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
    nr, ng, nb = colorsys.hls_to_rgb(h, 0.96, 0.30)
    return RGBColor(int(nr * 255), int(ng * 255), int(nb * 255))


def _tinted_gray(primary_hex: str) -> RGBColor:
    """primary hue를 섞은 중간 회색 (서브텍스트용)"""
    r, g, b = _hex_to_rgb(primary_hex)
    h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
    nr, ng, nb = colorsys.hls_to_rgb(h, 0.55, 0.08)
    return RGBColor(int(nr * 255), int(ng * 255), int(nb * 255))


def _headline_font_size(text: str, max_pt: int = 80) -> int:
    length = len(text)
    if length > 28:
        return min(max_pt, 48)
    elif length > 18:
        return min(max_pt, 60)
    else:
        return min(max_pt, 80)


def _apply_font(paragraph, latin: str = FONT_HEAD, ea: str = FONT_EA):
    """paragraph의 모든 run에 Latin + East Asian 폰트 지정"""
    for run in paragraph.runs:
        rPr = run._r.get_or_add_rPr()
        for el in rPr.findall(qn('a:latin')):
            rPr.remove(el)
        elem_latin = etree.SubElement(rPr, qn('a:latin'))
        elem_latin.set('typeface', latin)
        for el in rPr.findall(qn('a:ea')):
            rPr.remove(el)
        elem_ea = etree.SubElement(rPr, qn('a:ea'))
        elem_ea.set('typeface', ea)


def _bg(slide, color: RGBColor):
    """슬라이드 전체 배경 사각형"""
    rect = slide.shapes.add_shape(1, 0, 0, W, H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def _textbox(slide, x, y, w, h, text, size_pt, color: RGBColor,
             bold=False, align=PP_ALIGN.LEFT, wrap=True):
    """텍스트박스 1개 생성 + 스타일 지정 → paragraph 반환"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    _apply_font(p)
    return p


def _add_footer(slide, primary: RGBColor, gray: RGBColor,
                company_name: str, page_num: int, total: int):
    """하단 컬러 바 + 저작권 + 페이지 번호"""
    bar_h = int(H * 0.028)
    bar = slide.shapes.add_shape(1, 0, int(H - bar_h), W, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    copy_tb = slide.shapes.add_textbox(
        int(W * 0.04), int(H * 0.925), int(W * 0.50), Inches(0.22)
    )
    p = copy_tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = (
        f"© {datetime.datetime.now().year} {company_name}. All Rights Reserved."
    )
    run.font.size = Pt(7)
    run.font.color.rgb = gray
    _apply_font(p)

    pn_tb = slide.shapes.add_textbox(
        int(W * 0.86), int(H * 0.925), int(W * 0.11), Inches(0.22)
    )
    p = pn_tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num} / {total}"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(200, 200, 200)
    _apply_font(p)


def _left_panel(slide, primary: RGBColor, width_ratio=0.07):
    """좌측 primary 컬러 세로 패널"""
    panel = slide.shapes.add_shape(1, 0, 0, int(W * width_ratio), H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = primary
    panel.line.fill.background()


def _pill_badge(slide, text: str, primary: RGBColor, x, y):
    """회사명 pill badge (eyebrow)"""
    badge_w = Inches(max(len(text) * 0.11 + 0.5, 1.4))
    badge_h = Inches(0.30)
    badge = slide.shapes.add_shape(5, x, y, badge_w, badge_h)
    badge.fill.solid()
    badge.fill.fore_color.rgb = primary
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    _apply_font(p)


def _accent_bar(slide, primary: RGBColor, x, y, width_ratio=0.06, height=4):
    """가로 accent bar (섹션 구분용)"""
    bar = slide.shapes.add_shape(
        1, x, y, int(W * width_ratio), Pt(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()


# ══════════════════════════════════════════════════════════════════════
# 타입별 슬라이드 빌더
# ══════════════════════════════════════════════════════════════════════

def _build_cover(slide, info: dict, primary: RGBColor, gray: RGBColor,
                 bg_dark: RGBColor, company_name: str, page_num: int, total: int):
    """커버: tinted-dark 배경 + 좌측 컬러 패널 + pill badge + 오버사이즈 헤드라인"""
    _bg(slide, bg_dark)
    _left_panel(slide, primary)

    badge_x = int(W * 0.11)
    badge_y = int(H * 0.14)
    _pill_badge(slide, company_name, primary, badge_x, badge_y)

    headline = info.get("headline", company_name)
    hl_size = _headline_font_size(headline, max_pt=80)
    hl_x, hl_y = int(W * 0.11), int(H * 0.26)
    hl_w, hl_h = int(W * 0.72), int(H * 0.46)
    tb = slide.shapes.add_textbox(hl_x, hl_y, hl_w, hl_h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = headline
    run.font.size = Pt(hl_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    _apply_font(p)

    sub = info.get("subheadline", "")
    if sub:
        _textbox(slide, int(W * 0.11), int(H * 0.74), int(W * 0.72), int(H * 0.12),
                 sub, 18, gray)

    _add_footer(slide, primary, gray, company_name, page_num, total)


def _build_problem(slide, info: dict, primary: RGBColor, gray: RGBColor,
                   bg_dark: RGBColor, company_name: str, page_num: int, total: int):
    """문제: 흰 배경 + 상단 경고 빨간 accent + 헤드라인 + 체크리스트형 body"""
    white = RGBColor(255, 255, 255)
    dark_text = RGBColor(30, 30, 30)
    red_accent = RGBColor(239, 68, 68)  # Tailwind red-500

    _bg(slide, white)

    # 상단 좌측 빨간 사각 태그
    tag = slide.shapes.add_shape(1, 0, 0, int(W * 0.007), H)
    tag.fill.solid()
    tag.fill.fore_color.rgb = red_accent
    tag.line.fill.background()

    eyebrow = info.get("eyebrow", "")
    if eyebrow:
        _textbox(slide, int(W * 0.06), int(H * 0.10), int(W * 0.85), Inches(0.35),
                 eyebrow, 13, red_accent, bold=True)

    headline = info.get("headline", "")
    _textbox(slide, int(W * 0.06), int(H * 0.19), int(W * 0.85), int(H * 0.22),
             headline, 36, dark_text, bold=True)

    _accent_bar(slide, red_accent, int(W * 0.06), int(H * 0.43), width_ratio=0.05)

    body = info.get("body", [])
    if body:
        tb = slide.shapes.add_textbox(
            int(W * 0.06), int(H * 0.50), int(W * 0.85), int(H * 0.38)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"✗  {bullet}"
            run.font.size = Pt(17)
            run.font.color.rgb = RGBColor(80, 80, 80)
            _apply_font(p)

    _add_footer(slide, primary, gray, company_name, page_num, total)


def _build_solution(slide, info: dict, primary: RGBColor, gray: RGBColor,
                    bg_dark: RGBColor, company_name: str, page_num: int, total: int,
                    bg_light: RGBColor = None):
    """솔루션: 밝은 tinted 배경 + 좌측 컬러 패널 + 카드형 body"""
    if bg_light is None:
        bg_light = RGBColor(240, 245, 255)
    _bg(slide, bg_light)
    _left_panel(slide, primary, width_ratio=0.006)

    eyebrow = info.get("eyebrow", "")
    if eyebrow:
        _textbox(slide, int(W * 0.06), int(H * 0.10), int(W * 0.85), Inches(0.35),
                 eyebrow, 13, primary, bold=True)

    headline = info.get("headline", "")
    dark_text = RGBColor(20, 20, 20)
    _textbox(slide, int(W * 0.06), int(H * 0.19), int(W * 0.85), int(H * 0.22),
             headline, 36, dark_text, bold=True)

    _accent_bar(slide, primary, int(W * 0.06), int(H * 0.43), width_ratio=0.05)

    body = info.get("body", [])
    if body:
        n = len(body)
        card_w = int((W * 0.88) / n)
        card_h = int(H * 0.33)
        start_x = int(W * 0.06)
        card_y = int(H * 0.51)
        gap = int(W * 0.015)

        for i, item in enumerate(body):
            x = start_x + i * (card_w + gap)
            card = slide.shapes.add_shape(5, x, card_y, card_w - gap, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = primary
            card.line.width = Pt(0.75)

            tb = slide.shapes.add_textbox(
                x + Pt(12), card_y + Pt(14), card_w - gap - Pt(24), card_h - Pt(28)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(40, 40, 40)
            _apply_font(p)

    _add_footer(slide, primary, gray, company_name, page_num, total)


def _build_how_it_works(slide, info: dict, primary: RGBColor, gray: RGBColor,
                        bg_dark: RGBColor, company_name: str, page_num: int, total: int):
    """동작 방식: 흰 배경 + 번호 step 카드"""
    white = RGBColor(255, 255, 255)
    dark_text = RGBColor(20, 20, 20)
    _bg(slide, white)

    eyebrow = info.get("eyebrow", "")
    if eyebrow:
        _textbox(slide, int(W * 0.06), int(H * 0.08), int(W * 0.85), Inches(0.35),
                 eyebrow, 13, primary, bold=True)

    headline = info.get("headline", "")
    _textbox(slide, int(W * 0.06), int(H * 0.17), int(W * 0.85), int(H * 0.20),
             headline, 32, dark_text, bold=True)

    body = info.get("body", [])
    if body:
        n = min(len(body), 4)
        card_w = int((W * 0.88) / n)
        card_h = int(H * 0.42)
        start_x = int(W * 0.06)
        card_y = int(H * 0.44)
        gap = int(W * 0.012)

        for i in range(n):
            x = start_x + i * (card_w + gap)

            # 번호 원
            circle_r = int(H * 0.055)
            circle = slide.shapes.add_shape(
                9,  # OVAL = 9
                x + int((card_w - gap) / 2) - circle_r,
                card_y,
                circle_r * 2,
                circle_r * 2,
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = primary
            circle.line.fill.background()
            tf_c = circle.text_frame
            tf_c.word_wrap = False
            p_c = tf_c.paragraphs[0]
            p_c.alignment = PP_ALIGN.CENTER
            run_c = p_c.add_run()
            run_c.text = f"{i + 1:02d}"
            run_c.font.size = Pt(14)
            run_c.font.bold = True
            run_c.font.color.rgb = RGBColor(255, 255, 255)
            _apply_font(p_c)

            # 텍스트 카드
            text_y = card_y + circle_r * 2 + int(H * 0.02)
            tb = slide.shapes.add_textbox(
                x, text_y, card_w - gap, card_h - circle_r * 2 - int(H * 0.02)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = body[i]
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(50, 50, 50)
            _apply_font(p)

    _add_footer(slide, primary, gray, company_name, page_num, total)


def _build_key_metrics(slide, info: dict, primary: RGBColor, gray: RGBColor,
                       bg_dark: RGBColor, company_name: str, page_num: int, total: int):
    """핵심 수치: 숫자 대형 + 라벨 소형 (세로 나열)"""
    white = RGBColor(255, 255, 255)
    dark_text = RGBColor(20, 20, 20)
    _bg(slide, white)

    eyebrow = info.get("eyebrow", "")
    if eyebrow:
        _textbox(slide, int(W * 0.06), int(H * 0.08), int(W * 0.85), Inches(0.35),
                 eyebrow, 13, primary, bold=True)

    headline = info.get("headline", "")
    _textbox(slide, int(W * 0.06), int(H * 0.17), int(W * 0.85), int(H * 0.18),
             headline, 32, dark_text, bold=True)

    body = info.get("body", [])
    if body:
        n = len(body)
        avail_w = int(W * 0.88)
        col_w = avail_w // n
        start_x = int(W * 0.06)
        row_y = int(H * 0.42)

        for i, item in enumerate(body):
            x = start_x + i * col_w
            metric_h = int(H * 0.42)

            if ":" in item:
                label, val = item.split(":", 1)
                val = val.strip()
                label = label.strip()
            else:
                val = item
                label = ""

            # 수치 (대형)
            tb_val = slide.shapes.add_textbox(x, row_y, col_w, int(H * 0.26))
            tf_v = tb_val.text_frame
            p_v = tf_v.paragraphs[0]
            p_v.alignment = PP_ALIGN.CENTER
            run_v = p_v.add_run()
            run_v.text = val
            run_v.font.size = Pt(52)
            run_v.font.bold = True
            run_v.font.color.rgb = primary
            _apply_font(p_v)

            # 구분선
            line_y = row_y + int(H * 0.27)
            line_shape = slide.shapes.add_shape(
                1, x + int(col_w * 0.2), line_y, int(col_w * 0.6), Pt(1.5)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
            line_shape.line.fill.background()

            # 라벨
            if label:
                tb_lab = slide.shapes.add_textbox(
                    x, line_y + Pt(8), col_w, int(H * 0.14)
                )
                p_l = tb_lab.text_frame.paragraphs[0]
                p_l.alignment = PP_ALIGN.CENTER
                run_l = p_l.add_run()
                run_l.text = label
                run_l.font.size = Pt(14)
                run_l.font.color.rgb = RGBColor(100, 100, 100)
                _apply_font(p_l)

    _add_footer(slide, primary, gray, company_name, page_num, total)


def _build_proof(slide, info: dict, primary: RGBColor, gray: RGBColor,
                 bg_dark: RGBColor, company_name: str, page_num: int, total: int):
    """증거/성과: 다크 배경 + 큰 따옴표 + 인용/실적 리스트"""
    _bg(slide, bg_dark)
    white = RGBColor(255, 255, 255)

    # 큰 따옴표 장식
    tb_q = slide.shapes.add_textbox(int(W * 0.05), int(H * 0.04), Inches(1.2), Inches(1.0))
    p_q = tb_q.text_frame.paragraphs[0]
    run_q = p_q.add_run()
    run_q.text = "\u201c"  # left double quotation mark
    run_q.font.size = Pt(110)
    run_q.font.color.rgb = primary
    run_q.font.bold = True

    eyebrow = info.get("eyebrow", "")
    if eyebrow:
        _textbox(slide, int(W * 0.06), int(H * 0.11), int(W * 0.85), Inches(0.35),
                 eyebrow, 13, gray, bold=True)

    headline = info.get("headline", "")
    _textbox(slide, int(W * 0.06), int(H * 0.22), int(W * 0.85), int(H * 0.22),
             headline, 34, white, bold=True)

    body = info.get("body", [])
    if body:
        tb = slide.shapes.add_textbox(
            int(W * 0.06), int(H * 0.50), int(W * 0.85), int(H * 0.36)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"✓  {item}"
            run.font.size = Pt(16)
            run.font.color.rgb = gray
            _apply_font(p)

    _add_footer(slide, primary, gray, company_name, page_num, total)


def _build_why_us(slide, info: dict, primary: RGBColor, gray: RGBColor,
                  bg_dark: RGBColor, company_name: str, page_num: int, total: int):
    """차별점: 흰 배경 + 2열 비교 카드 또는 체크마크 리스트"""
    white = RGBColor(255, 255, 255)
    dark_text = RGBColor(20, 20, 20)
    _bg(slide, white)
    _left_panel(slide, primary, width_ratio=0.006)

    eyebrow = info.get("eyebrow", "")
    if eyebrow:
        _textbox(slide, int(W * 0.06), int(H * 0.08), int(W * 0.85), Inches(0.35),
                 eyebrow, 13, primary, bold=True)

    headline = info.get("headline", "")
    _textbox(slide, int(W * 0.06), int(H * 0.17), int(W * 0.85), int(H * 0.20),
             headline, 32, dark_text, bold=True)

    _accent_bar(slide, primary, int(W * 0.06), int(H * 0.40), width_ratio=0.05)

    body = info.get("body", [])
    if body:
        # 2열 레이아웃
        half = (len(body) + 1) // 2
        col_w = int(W * 0.42)
        col_gap = int(W * 0.06)
        left_x = int(W * 0.06)
        right_x = left_x + col_w + col_gap
        row_h = int(H * 0.11)
        start_y = int(H * 0.48)

        for i, item in enumerate(body):
            col = 0 if i < half else 1
            row = i if i < half else i - half
            x = left_x if col == 0 else right_x
            y = start_y + row * row_h

            tb = slide.shapes.add_textbox(x, y, col_w, row_h - Pt(4))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"✓  {item}"
            run.font.size = Pt(15)
            run.font.color.rgb = RGBColor(40, 40, 40)
            _apply_font(p)

    _add_footer(slide, primary, gray, company_name, page_num, total)


def _build_cta(slide, info: dict, primary: RGBColor, gray: RGBColor,
               bg_dark: RGBColor, company_name: str, page_num: int, total: int):
    """CTA: 다크 배경 + 중앙 헤드라인 + steps 카드"""
    _bg(slide, bg_dark)
    white = RGBColor(255, 255, 255)

    headline = info.get("headline", "")
    hl_size = min(_headline_font_size(headline, max_pt=64), 64)
    hl_w = int(W * 0.80)
    hl_h = int(H * 0.22)
    hl_x = int((W - hl_w) / 2)
    hl_y = int(H * 0.14)

    tb = slide.shapes.add_textbox(hl_x, hl_y, hl_w, hl_h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = headline
    run.font.size = Pt(hl_size)
    run.font.bold = True
    run.font.color.rgb = white
    _apply_font(p)

    sub = info.get("subheadline", "")
    if sub:
        sub_y = hl_y + hl_h + int(H * 0.01)
        _textbox(slide, int(W * 0.10), sub_y, int(W * 0.80), Inches(0.40),
                 sub, 17, gray, align=PP_ALIGN.CENTER)

    body = info.get("body", [])
    if body:
        n = min(len(body), 3)
        card_w = int(W * 0.26)
        card_h = int(H * 0.26)
        gap = int(W * 0.03)
        total_w = card_w * n + gap * (n - 1)
        start_x = int((W - total_w) / 2)
        card_y = int(H * 0.54)

        for i in range(n):
            x = start_x + i * (card_w + gap)
            card = slide.shapes.add_shape(5, x, card_y, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = bg_dark
            card.line.color.rgb = gray
            card.line.width = Pt(0.5)

            num_tb = slide.shapes.add_textbox(x + Pt(12), card_y + Pt(12), card_w - Pt(24), Pt(22))
            p_n = num_tb.text_frame.paragraphs[0]
            run_n = p_n.add_run()
            run_n.text = f"{i + 1:02d}"
            run_n.font.size = Pt(12)
            run_n.font.bold = True
            run_n.font.color.rgb = primary
            _apply_font(p_n)

            txt_tb = slide.shapes.add_textbox(
                x + Pt(12), card_y + Pt(40), card_w - Pt(24), card_h - Pt(55)
            )
            tf_t = txt_tb.text_frame
            tf_t.word_wrap = True
            p_t = tf_t.paragraphs[0]
            run_t = p_t.add_run()
            run_t.text = body[i]
            run_t.font.size = Pt(12)
            run_t.font.color.rgb = white
            _apply_font(p_t)

    _add_footer(slide, primary, gray, company_name, page_num, total)


def _build_default(slide, info: dict, primary: RGBColor, gray: RGBColor,
                   bg_dark: RGBColor, company_name: str, page_num: int, total: int):
    """폴백: 흰 배경 + eyebrow + 헤드라인 + 불릿 body"""
    white = RGBColor(255, 255, 255)
    dark_text = RGBColor(20, 20, 20)
    _bg(slide, white)

    eyebrow = info.get("eyebrow", "")
    if eyebrow:
        _textbox(slide, int(W * 0.06), int(H * 0.08), int(W * 0.85), Inches(0.35),
                 eyebrow, 13, primary, bold=True)

    headline = info.get("headline", "")
    _textbox(slide, int(W * 0.06), int(H * 0.19), int(W * 0.85), int(H * 0.22),
             headline, 34, dark_text, bold=True)

    _accent_bar(slide, primary, int(W * 0.06), int(H * 0.43), width_ratio=0.05)

    body = info.get("body", [])
    if body:
        tb = slide.shapes.add_textbox(
            int(W * 0.06), int(H * 0.50), int(W * 0.85), int(H * 0.36)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = f"•  {bullet}"
            run.font.size = Pt(17)
            run.font.color.rgb = RGBColor(60, 60, 60)
            _apply_font(p)

    _add_footer(slide, primary, gray, company_name, page_num, total)


# ══════════════════════════════════════════════════════════════════════
# 공개 진입점
# ══════════════════════════════════════════════════════════════════════

_SLIDE_BUILDERS = {
    "cover": _build_cover,
    "problem": _build_problem,
    "solution": _build_solution,
    "how_it_works": _build_how_it_works,
    "key_metrics": _build_key_metrics,
    "proof": _build_proof,
    "why_us": _build_why_us,
    "cta": _build_cta,
}


def build_pptx(slide_data: dict) -> bytes:
    """SlideData dict → PPTX bytes"""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    brand = slide_data.get("brand", {})
    primary_hex = brand.get("primaryColor", "#2563EB")
    company_name = brand.get("companyName", "TickDeck")

    primary = _rgb(primary_hex)
    bg_dark = _tinted_dark(primary_hex)
    bg_light = _tinted_light(primary_hex)
    gray = _tinted_gray(primary_hex)

    slides = slide_data.get("slides", [])
    total = len(slides)

    for page_num, slide_info in enumerate(slides, start=1):
        stype = slide_info.get("type", "content")
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        builder = _SLIDE_BUILDERS.get(stype, _build_default)
        try:
            if stype == "solution":
                builder(slide, slide_info, primary, gray, bg_dark, company_name, page_num, total, bg_light=bg_light)
            else:
                builder(slide, slide_info, primary, gray, bg_dark, company_name, page_num, total)
        except Exception as e:
            logger.warning(f"슬라이드 {page_num} ({stype}) 빌드 실패, 폴백: {e}")
            _build_default(slide, slide_info, primary, gray, bg_dark, company_name, page_num, total)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()
