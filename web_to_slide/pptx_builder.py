"""
pptx_builder.py — python-pptx 서버사이드 PPTX 생성
Phase 1: 커버 슬라이드 (코드 빌드 방식 — 템플릿 없이)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import io
import base64
import colorsys

# ── 컬러 유틸 ──────────────────────────────────────────

def _hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c * 2 for c in hex_color])
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def _tinted_dark(primary_hex):
    """primary hue를 섞은 거의-검정 (커버 배경)"""
    r, g, b = _hex_to_rgb(primary_hex)
    h, l, s = colorsys.rgb_to_hls(r/255, g/255, b/255)
    nr, ng, nb = colorsys.hls_to_rgb(h, 0.08, 0.25)
    return RGBColor(int(nr*255), int(ng*255), int(nb*255))

def _tinted_gray(primary_hex):
    """primary hue를 섞은 중간 회색 (서브텍스트)"""
    r, g, b = _hex_to_rgb(primary_hex)
    h, l, s = colorsys.rgb_to_hls(r/255, g/255, b/255)
    nr, ng, nb = colorsys.hls_to_rgb(h, 0.55, 0.08)
    return RGBColor(int(nr*255), int(ng*255), int(nb*255))

def _headline_font_size(text):
    length = len(text)
    if length > 28: return Pt(52)
    elif length > 18: return Pt(64)
    else: return Pt(80)

def _set_shape_transparency(shape, alpha_pct):
    """도형 fill에 투명도 적용 (0=불투명, 100=완전투명)"""
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.find(f'{{{ns}}}spPr') or shape._element.find('.//{{{ns}}}spPr'.format(ns=ns))
    if spPr is None:
        return
    solidFill = spPr.find(f'{{{ns}}}solidFill')
    if solidFill is None:
        return
    srgb = solidFill.find(f'{{{ns}}}srgbClr')
    if srgb is not None:
        alpha = etree.SubElement(srgb, f'{{{ns}}}alpha')
        alpha.set('val', str(int((100 - alpha_pct) * 1000)))

# ── 슬라이드 사이즈 ──────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# ── 커버 빌더 (코드 빌드) ──────────────────────────────────

def build_cover(brand, headline, sub, logo_b64=None):
    """
    커버 슬라이드 1장을 처음부터 코드로 빌드.
    Returns: bytes (PPTX)
    """
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    primary_hex = brand.get('primaryColor', '#1C3D5A')
    p_rgb = _hex_to_rgb(primary_hex)
    primary = RGBColor(*p_rgb)
    bg_color = _tinted_dark(primary_hex)
    gray = _tinted_gray(primary_hex)
    company_name = brand.get('name', 'Company')

    # ── 1. 배경 (tintedDark 단색 flat — Rule 7: 순수 검정 피함, 스타일가이드: 그라데이션 최소화) ──
    bg = slide.shapes.add_shape(
        1, 0, 0, W, H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # ── 2. 좌측 primary 컬러 패널 (패턴1: Color Panel Split — 슬라이드 좌 8%) ──
    panel = slide.shapes.add_shape(
        1, 0, 0, int(W * 0.08), H
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = primary
    panel.line.fill.background()

    # ── 3. Pill badge (회사명 — eyebrow 역할) ──
    badge_text = company_name.upper()
    badge_w = Inches(max(len(badge_text) * 0.12 + 0.5, 1.5))
    badge_h = Inches(0.35)
    badge_x = int(W * 0.12)  # 좌측 패널 오른쪽
    badge_y = int(H * 0.14)
    badge = slide.shapes.add_shape(
        5, badge_x, badge_y, badge_w, badge_h  # ROUNDED_RECTANGLE = 5
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = primary
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = badge_text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    # ── 4. 헤드라인 (오버사이즈 타이포 — 패턴2) ──
    hl_x = int(W * 0.12)
    hl_y = int(H * 0.26)
    hl_w = int(W * 0.70)
    hl_h = int(H * 0.42)
    hl_box = slide.shapes.add_textbox(hl_x, hl_y, hl_w, hl_h)
    tf = hl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = headline or company_name
    p.font.size = _headline_font_size(headline or company_name)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # ── 5. Accent line 제거 (PPTX 스킬: "AI 슬라이드의 특징" → 여백으로 대체) ──

    # ── 6. 서브헤드라인 (tintedGray) ──
    sub_x = int(W * 0.12)
    sub_y = int(H * 0.72)
    sub_w = int(W * 0.70)
    sub_h = int(H * 0.12)
    sub_box = slide.shapes.add_textbox(sub_x, sub_y, sub_w, sub_h)
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = sub or ''
    p.font.size = Pt(18)
    p.font.color.rgb = gray
    p.alignment = PP_ALIGN.LEFT

    # ── 7. 하단 bar (primary — 얇게) ──
    bar_h = int(H * 0.03)
    bar = slide.shapes.add_shape(
        1, 0, int(H - bar_h), W, bar_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    # ── 8. 저작권 ──
    copy_box = slide.shapes.add_textbox(
        int(W * 0.04), int(H * 0.92), int(W * 0.40), Inches(0.25)
    )
    p = copy_box.text_frame.paragraphs[0]
    import datetime
    p.text = f"© {datetime.datetime.now().year} {company_name}. All Rights Reserved."
    p.font.size = Pt(7)
    p.font.color.rgb = gray

    # ── 9. 페이지 번호 ──
    pn_box = slide.shapes.add_textbox(
        int(W * 0.85), int(H * 0.92), int(W * 0.12), Inches(0.25)
    )
    p = pn_box.text_frame.paragraphs[0]
    p.text = "1 / 9"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.RIGHT

    # ── 10. 로고 (있으면) ──
    if logo_b64:
        try:
            if ',' in logo_b64:
                logo_b64 = logo_b64.split(',')[1]
            image_data = base64.b64decode(logo_b64)
            image_stream = io.BytesIO(image_data)
            from PIL import Image
            with Image.open(image_stream) as img:
                w, h = img.size
                aspect = w / float(h or 1)
                max_w = Inches(1.8)
                max_h = Inches(0.5)
                if aspect >= max_w / max_h:
                    l_w = max_w
                    l_h = int(max_w / aspect)
                else:
                    l_h = max_h
                    l_w = int(max_h * aspect)
                l_x = int(W - l_w - Inches(0.5))
                l_y = int(H * 0.85)
                image_stream.seek(0)
                slide.shapes.add_picture(image_stream, l_x, l_y, width=l_w, height=l_h)
        except Exception:
            pass

    # 저장
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def merge_cover(pptx_bytes, brand, headline, sub, logo_b64=None):
    """
    Phase 1.5: PptxGenJS가 만든 PPTX에서 slide[0]을 python-pptx 커버로 교체.
    Returns: bytes (병합된 PPTX)
    """
    from copy import deepcopy

    # 1. PptxGenJS PPTX 열기
    src_stream = io.BytesIO(pptx_bytes)
    prs = Presentation(src_stream)

    # 2. python-pptx로 커버 PPTX 생성
    cover_bytes = build_cover(brand, headline, sub, logo_b64)
    cover_stream = io.BytesIO(cover_bytes)
    cover_prs = Presentation(cover_stream)
    cover_slide = cover_prs.slides[0]

    # 3. 기존 slide[0] 삭제
    if len(prs.slides) > 0:
        first_slide = prs.slides[0]
        rId = None
        for rel in prs.part.rels.values():
            if rel.target_part == first_slide.part:
                rId = rel.rId
                break
        if rId:
            sldIdLst = prs.presentation.sldIdLst
            for sldId in sldIdLst:
                if sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') == rId:
                    sldIdLst.remove(sldId)
                    break
            del prs.part.rels[rId]

    # 4. 새 커버를 추가하고 맨 앞으로 이동
    slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(slide_layout)
    new_slide_elem = new_slide._element
    cover_elem = cover_slide._element
    sp_tree_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
    old_cSld = new_slide_elem.find(sp_tree_tag)
    new_cSld = deepcopy(cover_elem.find(sp_tree_tag))
    if old_cSld is not None and new_cSld is not None:
        new_slide_elem.replace(old_cSld, new_cSld)

    sldIdLst = prs.presentation.sldIdLst
    sldId_elements = list(sldIdLst)
    if len(sldId_elements) >= 2:
        last = sldId_elements[-1]
        sldIdLst.remove(last)
        sldIdLst.insert(0, last)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
